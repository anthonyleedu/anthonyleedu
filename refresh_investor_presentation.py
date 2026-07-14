#!/usr/bin/env python3

from pathlib import Path
import tempfile
import zipfile

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_DATA_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


DECK_PATH = Path("/workspace/Samuel Space Organ and Tissue Engineering - Investor Presentation (2).pptx")

BG = RGBColor(247, 249, 252)
WHITE = RGBColor(255, 255, 255)
NAVY = RGBColor(17, 31, 61)
TEAL = RGBColor(19, 184, 169)
TEAL_DARK = RGBColor(0, 132, 122)
TEXT = RGBColor(36, 49, 74)
MUTED = RGBColor(96, 109, 130)
LINE = RGBColor(220, 228, 238)
SOFT = RGBColor(236, 248, 246)
SOFT_BLUE = RGBColor(238, 244, 252)
SOFT_NAVY = RGBColor(233, 239, 248)
SOFT_RED = RGBColor(250, 243, 243)
SERIES_BLUE = RGBColor(94, 134, 255)
SERIES_LIGHT = RGBColor(157, 191, 255)

MARGIN_X = 0.6
CONTENT_W = 12.133


def emu(value):
    return Inches(value)


def set_background(slide, color=BG):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def remove_defaults(slide):
    for shape in list(slide.shapes):
        shape._element.getparent().remove(shape._element)


def style_paragraph(paragraph, size, color=TEXT, bold=False, font_name="Aptos", align=PP_ALIGN.LEFT):
    paragraph.alignment = align
    paragraph.space_after = Pt(0)
    paragraph.space_before = Pt(0)
    for run in paragraph.runs:
        run.font.name = font_name
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = color


def add_box(slide, x, y, w, h, fill=WHITE, line=LINE, radius=True, line_width=1.0):
    kind = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, x, y, w, h=0.03, color=LINE):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(x), emu(y), emu(w), emu(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_circle(slide, x, y, d, fill=TEAL, line=TEAL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, emu(x), emu(y), emu(d), emu(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = line
    shape.line.width = Pt(1)
    return shape


def add_textbox(
    slide,
    x,
    y,
    w,
    h,
    text,
    size=12,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
    margin=0.04,
    valign=MSO_VERTICAL_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = emu(margin)
    frame.margin_right = emu(margin)
    frame.margin_top = emu(margin)
    frame.margin_bottom = emu(margin)
    frame.vertical_anchor = valign
    frame.auto_size = MSO_AUTO_SIZE.NONE
    frame.text = text
    for paragraph in frame.paragraphs:
        style_paragraph(paragraph, size=size, color=color, bold=bold, font_name=font_name, align=align)
    return box


def add_multiline_textbox(
    slide,
    x,
    y,
    w,
    h,
    lines,
    size=11,
    color=TEXT,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Aptos",
    margin=0.05,
    bullet=False,
    bullet_color=TEAL,
    line_gap=2,
):
    box = slide.shapes.add_textbox(emu(x), emu(y), emu(w), emu(h))
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = emu(margin)
    frame.margin_right = emu(margin)
    frame.margin_top = emu(margin)
    frame.margin_bottom = emu(margin)
    frame.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    frame.auto_size = MSO_AUTO_SIZE.NONE
    for idx, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = f"- {line}" if bullet else line
        paragraph.alignment = align
        paragraph.space_after = Pt(line_gap)
        for run in paragraph.runs:
            run.font.name = font_name
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = bullet_color if bullet and run.text.startswith("- ") else color
        if bullet:
            for run in paragraph.runs[1:]:
                run.font.color.rgb = color
    return box


def add_header(slide, section, title, subtitle, number):
    set_background(slide)
    add_line(slide, 0, 0, 13.333, 0.08, color=TEAL)
    add_textbox(slide, MARGIN_X, 0.22, 2.5, 0.22, section.upper(), size=10, color=TEAL_DARK, bold=True)
    add_textbox(slide, MARGIN_X, 0.48, 9.4, 0.52, title, size=27, color=NAVY, bold=True, font_name="Aptos Display")
    if subtitle:
        add_textbox(slide, MARGIN_X, 0.96, 11.7, 0.34, subtitle, size=11.2, color=MUTED)
    add_textbox(
        slide,
        12.35,
        7.0,
        0.4,
        0.18,
        str(number),
        size=8.5,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
        margin=0,
    )


def add_kpi_card(slide, x, y, w, h, value, label):
    add_box(slide, x, y, w, h, fill=WHITE, line=LINE)
    add_textbox(slide, x + 0.18, y + 0.14, w - 0.36, 0.48, value, size=24, color=NAVY, bold=True)
    add_textbox(slide, x + 0.18, y + 0.68, w - 0.36, h - 0.82, label, size=10.2, color=MUTED)


def add_info_card(slide, x, y, w, h, heading, body, tag=None, fill=WHITE):
    add_box(slide, x, y, w, h, fill=fill, line=LINE)
    add_textbox(slide, x + 0.18, y + 0.16, w - 0.36, 0.25, heading, size=13, color=NAVY, bold=True)
    add_textbox(slide, x + 0.18, y + 0.48, w - 0.36, h - 0.68, body, size=10, color=TEXT)
    if tag:
        add_textbox(slide, x + 0.18, y + h - 0.28, w - 0.36, 0.16, tag, size=8.8, color=TEAL_DARK, bold=True)


def add_pill(slide, x, y, w, h, text, fill=SOFT, color=TEAL_DARK, align=PP_ALIGN.CENTER, bold=True):
    pill = add_box(slide, x, y, w, h, fill=fill, line=fill)
    pill.line.color.rgb = fill
    add_textbox(slide, x + 0.02, y + 0.02, w - 0.04, h - 0.04, text, size=9, color=color, bold=bold, align=align, margin=0)


def style_chart(chart, *, has_legend=False, legend_bottom=False):
    chart.chart_style = None
    chart.has_title = False
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM if legend_bottom else XL_LEGEND_POSITION.TOP
        chart.legend.font.size = Pt(9)
        chart.legend.font.color.rgb = TEXT
        chart.legend.font.name = "Aptos"


def add_bar_chart(slide, x, y, w, h, categories, values):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Share (%)", values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, emu(x), emu(y), emu(w), emu(h), data)
    chart = graphic.chart
    style_chart(chart)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = TEAL
    series.format.line.color.rgb = TEAL
    plot = chart.plots[0]
    plot.gap_width = 45
    plot.has_data_labels = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.font.size = Pt(9)
    plot.data_labels.font.color.rgb = TEXT
    plot.data_labels.number_format = "0.0"
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.color.rgb = TEXT
    chart.value_axis.maximum_scale = 100
    chart.value_axis.minimum_scale = 0
    chart.value_axis.major_unit = 20
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = MUTED
    return chart


def add_column_chart(slide, x, y, w, h, categories, values):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Market ($B)", values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_CLUSTERED, emu(x), emu(y), emu(w), emu(h), data)
    chart = graphic.chart
    style_chart(chart)
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = NAVY
    series.format.line.color.rgb = NAVY
    plot = chart.plots[0]
    plot.gap_width = 45
    plot.has_data_labels = True
    plot.data_labels.position = XL_DATA_LABEL_POSITION.OUTSIDE_END
    plot.data_labels.number_format = '0.00"'
    plot.data_labels.font.size = Pt(8.5)
    plot.data_labels.font.color.rgb = TEXT
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.color.rgb = TEXT
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = MUTED
    chart.value_axis.major_unit = 0.5
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 4.0
    return chart


def add_stacked_chart(slide, x, y, w, h, categories, series_data):
    data = CategoryChartData()
    data.categories = categories
    for name, values in series_data:
        data.add_series(name, values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.COLUMN_STACKED, emu(x), emu(y), emu(w), emu(h), data)
    chart = graphic.chart
    style_chart(chart, has_legend=True)
    colors = [TEAL, NAVY, SERIES_BLUE]
    for idx, series in enumerate(chart.series):
        series.format.fill.solid()
        series.format.fill.fore_color.rgb = colors[idx]
        series.format.line.color.rgb = colors[idx]
    plot = chart.plots[0]
    plot.gap_width = 55
    chart.category_axis.tick_labels.font.size = Pt(9)
    chart.category_axis.tick_labels.font.color.rgb = TEXT
    chart.value_axis.has_major_gridlines = True
    chart.value_axis.major_gridlines.format.line.color.rgb = LINE
    chart.value_axis.tick_labels.font.size = Pt(9)
    chart.value_axis.tick_labels.font.color.rgb = MUTED
    chart.value_axis.minimum_scale = 0
    chart.value_axis.maximum_scale = 26
    chart.value_axis.major_unit = 5
    return chart


def add_doughnut_chart(slide, x, y, w, h, categories, values):
    data = CategoryChartData()
    data.categories = categories
    data.add_series("Use of Funds", values)
    graphic = slide.shapes.add_chart(XL_CHART_TYPE.DOUGHNUT, emu(x), emu(y), emu(w), emu(h), data)
    chart = graphic.chart
    style_chart(chart)
    colors = [TEAL, NAVY, SERIES_BLUE, SERIES_LIGHT, SOFT_NAVY]
    for idx, point in enumerate(chart.series[0].points):
        point.format.fill.solid()
        point.format.fill.fore_color.rgb = colors[idx]
        point.format.line.color.rgb = WHITE
    chart.plots[0].has_data_labels = True
    chart.plots[0].data_labels.number_format = '0"%"'
    chart.plots[0].data_labels.position = XL_DATA_LABEL_POSITION.CENTER
    chart.plots[0].data_labels.font.size = Pt(9)
    chart.plots[0].data_labels.font.color.rgb = WHITE
    chart.plots[0].doughnut_hole_size = 60
    return chart


def add_manual_table(slide, x, y, widths, rows, row_heights, header_fill=NAVY):
    current_y = y
    for row_idx, row in enumerate(rows):
        current_x = x
        for col_idx, cell_text in enumerate(row):
            fill = header_fill if row_idx == 0 else WHITE
            color = WHITE if row_idx == 0 else TEXT
            box = add_box(slide, current_x, current_y, widths[col_idx], row_heights[row_idx], fill=fill, line=LINE, radius=False, line_width=0.8)
            box.line.color.rgb = LINE if row_idx > 0 else header_fill
            add_textbox(
                slide,
                current_x + 0.08,
                current_y + 0.06,
                widths[col_idx] - 0.16,
                row_heights[row_idx] - 0.12,
                cell_text,
                size=9.2 if row_idx == 0 else 8.9,
                color=color,
                bold=row_idx == 0 or col_idx == 0,
            )
            current_x += widths[col_idx]
        current_y += row_heights[row_idx]


def extract_assets():
    asset_map = {
        "cover": "ppt/media/image-1-1.png",
        "close": "ppt/media/image-20-1.png",
        "earth": "ppt/media/image-5-1.png",
        "microgravity": "ppt/media/image-5-2.png",
        "printer": "ppt/media/image-6-6.png",
        "platform": "ppt/media/image-11-4.png",
        "cardiac": "ppt/media/image-7-1.png",
        "ortho": "ppt/media/image-7-2.png",
        "retina": "ppt/media/image-7-3.png",
        "wound": "ppt/media/image-7-4.png",
    }
    temp_dir = Path(tempfile.mkdtemp(prefix="samuel-space-assets-"))
    with zipfile.ZipFile(DECK_PATH) as source:
        for key, member in asset_map.items():
            (temp_dir / f"{key}.png").write_bytes(source.read(member))
    return {key: temp_dir / f"{key}.png" for key in asset_map}


def add_cover_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    set_background(slide, WHITE)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(8.6), 0, emu(4.733), emu(7.5)).fill.solid()
    hero_panel = slide.shapes[-1]
    hero_panel.fill.fore_color.rgb = SOFT
    hero_panel.line.fill.background()
    slide.shapes.add_picture(str(assets["cover"]), emu(8.78), emu(0), height=emu(7.5))
    add_line(slide, 0, 0, 13.333, 0.08, color=TEAL)
    add_textbox(slide, 0.78, 0.42, 2.8, 0.24, "CONFIDENTIAL INVESTOR PRESENTATION", size=10, color=TEAL_DARK, bold=True)
    add_textbox(slide, 0.78, 1.18, 6.9, 1.9, "Samuel Space\nOrgan and Tissue\nEngineering", size=28, color=NAVY, bold=True, font_name="Aptos Display")
    add_textbox(slide, 0.78, 3.34, 6.5, 0.7, "Manufacturing the next generation of human tissue - in orbit.", size=15.5, color=TEXT)
    add_box(slide, 0.78, 5.72, 3.18, 0.62, fill=SOFT_BLUE, line=SOFT_BLUE)
    add_textbox(slide, 0.96, 5.9, 2.85, 0.2, "Series A Financing  -  June 2026", size=11.2, color=NAVY, bold=True)
    add_textbox(slide, 0.78, 6.72, 5.5, 0.16, "Science-led orbital biomanufacturing platform for tissue engineering", size=9.5, color=MUTED)


def add_notice_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "NOTICE", "Confidential - Forward-Looking Statement", "", 2)
    add_box(slide, 0.78, 1.35, 11.8, 5.55, fill=WHITE, line=LINE)
    paragraphs = [
        "This presentation has been prepared by Samuel Space Organ and Tissue Engineering (the Company) solely to assist prospective investors in evaluating a potential investment. It is confidential and may not be reproduced or distributed, in whole or in part, without the Company's prior written consent.",
        "This presentation contains forward-looking statements, including projections of financial and operating performance, market size, regulatory timelines, and technology milestones. These statements reflect current expectations and assumptions and are subject to significant risks and uncertainties; actual results may differ materially. Nothing herein constitutes an offer to sell, or a solicitation of an offer to buy, any securities, and any such offer will be made only pursuant to definitive offering documents.",
        "Certain industry, market, and scientific data in this presentation has been compiled from third-party sources, including government agencies, peer-reviewed publications, and public company disclosures, as cited in the Appendix. The Company has not independently verified all such data and makes no representation as to its completeness or accuracy.",
    ]
    add_multiline_textbox(slide, 1.05, 1.68, 11.25, 4.6, paragraphs, size=10.8, color=TEXT, line_gap=7)
    add_pill(slide, 1.05, 6.18, 2.7, 0.36, "Investor discussion material")


def add_exec_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "EXECUTIVE SUMMARY", "Translating Proven Orbital Science Into a Commercial Tissue Platform", "", 3)
    kpis = [
        ("103,223", "People on the U.S. transplant waitlist today (OPTN, 2025)"),
        ("13 / day", "Americans die waiting for a donor organ"),
        ("$1.7B to $3.5B+", "Global bioprinting market from 2025 to 2030 (~16% CAGR)"),
        ("8 years", "Continuous in-orbit bioprinting validation"),
    ]
    card_w = 2.88
    for idx, (value, label) in enumerate(kpis):
        add_kpi_card(slide, 0.78 + idx * 3.02, 1.38, card_w, 1.35, value, label)
    info = [
        ("The Problem", "Organ and tissue shortages persist globally. Earth-based bioprinting hits a physical ceiling: gravity causes tissues to collapse without a temporary scaffold, limiting achievable geometry."),
        ("The Unlock", "Microgravity eliminates scaffold dependency. NASA, Redwire, ESA, LambdaVision and ROSCOSMOS have validated the science aboard the ISS continuously since 2018."),
        ("The Market", "A $1.7B market growing at ~16% CAGR, the first FDA-cleared tissue-engineered vascular graft (Symvess, Dec 2024), and four commercial stations entering service from 2026."),
        ("The Ask", "An $18M Series A to build Samuel's first flight-qualified payload platform, commercialise wound-care and orthopedic verticals, and file NASA InSPA Phase 1 grants."),
    ]
    positions = [(0.78, 3.02), (6.18, 3.02), (0.78, 5.03), (6.18, 5.03)]
    for (heading, body), (x, y) in zip(info, positions):
        add_info_card(slide, x, y, 5.37, 1.65, heading, body, fill=WHITE)


def add_problem_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "THE PROBLEM", "A Global Organ & Tissue Crisis", "", 4)
    add_box(slide, 0.78, 1.36, 6.0, 4.65, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 1.56, 4.6, 0.2, "U.S. Transplant Waiting List by Organ (2025)", size=11.5, color=NAVY, bold=True)
    add_bar_chart(slide, 1.0, 1.9, 5.55, 3.55, ["Kidney", "Liver", "Heart", "Lung", "Pancreas", "Other"], [87.0, 9.0, 2.0, 1.0, 0.5, 0.5])
    add_textbox(slide, 1.0, 5.56, 5.4, 0.18, "Source: OPTN / organdonor.gov, 2025", size=8.5, color=MUTED)
    stats = [
        ("941,652", "CVD deaths in the U.S. in 2022 - still the #1 cause of death, killing more people than all cancers combined"),
        ("~1 per 34 sec", "A cardiovascular death occurs in the U.S. on average every 34 seconds - roughly 2,500 per day"),
        ("10% met", "Of the estimated global transplant need. WHO reports ~130,000 transplants/yr vs. demand exceeding 1.3M"),
    ]
    for idx, (value, body) in enumerate(stats):
        add_info_card(slide, 7.0, 1.45 + idx * 1.38, 5.55, 1.15, value, body, fill=WHITE)
    add_textbox(slide, 7.08, 5.6, 5.3, 0.16, "Source: AHA 2025 Heart Disease & Stroke Statistics; WHO Global Observatory on Donation and Transplantation.", size=8.1, color=MUTED)
    add_box(slide, 0.78, 6.18, 11.77, 0.72, fill=SOFT_BLUE, line=SOFT_BLUE)
    add_textbox(slide, 1.0, 6.35, 11.3, 0.28, "Transplant recipients also face a lifetime of immunosuppressive therapy and rejection risk. Tissue printed from a patient's own cells eliminates cross-reactive rejection entirely - removing both the clinical burden and the ongoing drug cost.", size=9.6, color=TEXT)


def add_constraint_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "THE CORE CONSTRAINT", "Why Earth-Based Bioprinting Has a Hard Ceiling", "", 5)
    add_box(slide, 0.78, 1.32, 11.77, 0.72, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 1.5, 11.3, 0.34, "On Earth, gravity forces tissues to collapse under their own weight during printing. A scaffold must be built into every construct, adding complexity, introducing foreign material, and limiting achievable geometry. Microgravity removes this constraint entirely.", size=10.2, color=TEXT)
    add_box(slide, 0.78, 2.23, 5.7, 4.45, fill=WHITE, line=LINE)
    add_box(slide, 6.85, 2.23, 5.7, 4.45, fill=WHITE, line=LINE)
    add_textbox(slide, 1.02, 2.44, 2.6, 0.22, "EARTH GRAVITY", size=12, color=NAVY, bold=True)
    add_textbox(slide, 7.08, 2.44, 3.5, 0.22, "MICROGRAVITY  -  ISS ORBIT", size=12, color=NAVY, bold=True)
    add_box(slide, 1.0, 2.82, 5.26, 1.9, fill=SOFT_RED, line=SOFT_RED)
    slide.shapes.add_picture(str(assets["earth"]), emu(2.04), emu(2.93), height=emu(1.65))
    add_box(slide, 7.08, 2.82, 5.26, 1.9, fill=SOFT, line=SOFT)
    slide.shapes.add_picture(str(assets["microgravity"]), emu(8.12), emu(2.9), height=emu(1.66))
    add_multiline_textbox(
        slide,
        1.06,
        4.98,
        5.16,
        1.34,
        [
            "Cells slump and fuse unevenly under gravitational load",
            "Scaffold material printed first, then must be dissolved",
            "Complex hollow structures (e.g. heart chambers) are extremely difficult",
        ],
        size=10,
        color=TEXT,
        bullet=True,
    )
    add_multiline_textbox(
        slide,
        7.14,
        4.98,
        5.12,
        1.34,
        [
            "Tissue self-supports in true 3D - no scaffold required",
            "More uniform cell distribution and structural integrity",
            "Complex vascular and hollow geometries become achievable",
        ],
        size=10,
        color=TEXT,
        bullet=True,
    )


def add_technology_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "THE TECHNOLOGY", "How Space-Based Bioprinting Works: Five-Step Process", "", 6)
    steps = [
        ("1", "Cell Harvest", "Biopsy yields patient-specific cells, eliminating immunological rejection risk."),
        ("2", "Stem Cell Reprogramming", "Nobel-winning iPSC method converts adult cells to pluripotent stem cells."),
        ("3", "Bioink Formulation", "Cells are suspended in hydrogel optimised for viscosity and viability."),
        ("4", "In-Orbit Printing", "Layer-by-layer deposition in microgravity - no scaffold required."),
        ("5", "Tissue Maturation", "Bioreactor conditioning allows cells to differentiate into target tissue type."),
    ]
    card_w = 2.28
    for idx, (step, title, body) in enumerate(steps):
        x = 0.78 + idx * 2.41
        add_box(slide, x, 1.34, card_w, 1.72, fill=WHITE, line=LINE)
        add_circle(slide, x + 0.16, 1.48, 0.36, fill=TEAL, line=TEAL)
        add_textbox(slide, x + 0.16, 1.56, 0.36, 0.12, step, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_textbox(slide, x + 0.62, 1.47, 1.48, 0.34, title, size=11.2, color=NAVY, bold=True)
        add_textbox(slide, x + 0.16, 1.86, 1.95, 0.8, body, size=9.3, color=TEXT)
    add_box(slide, 0.78, 3.35, 6.0, 3.25, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 3.55, 5.5, 0.18, "Schematic: layer-by-layer orbital bioprinting hardware", size=9.6, color=TEAL_DARK, bold=True)
    slide.shapes.add_picture(str(assets["printer"]), emu(1.0), emu(3.8), width=emu(5.55))
    add_box(slide, 7.0, 3.35, 5.55, 2.0, fill=WHITE, line=LINE)
    add_textbox(slide, 7.22, 3.55, 2.2, 0.18, "Why orbit?", size=12, color=NAVY, bold=True)
    add_textbox(slide, 7.22, 3.82, 5.1, 1.25, "In the near-weightlessness of ISS orbit (~400 km altitude), tissues grow in true three dimensions with no gravitational compression. This principle has been validated continuously since 2018 across NASA, Redwire Corporation, ROSCOSMOS, LambdaVision Inc., and ESA/DLR - eight years of confirmation by organisations operating independently of one another.", size=10.1, color=TEXT)
    add_box(slide, 7.0, 5.58, 5.55, 1.02, fill=SOFT_BLUE, line=SOFT_BLUE)
    add_textbox(slide, 7.22, 5.77, 5.08, 0.62, "Note: Adult cardiac muscle cannot self-regenerate - damaged tissue is replaced by scar tissue that blocks electrical signals, a key driver of CVD mortality and why cardiac is Samuel's highest-value target vertical.", size=9.6, color=TEXT)


def add_product_strategy_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "PRODUCT STRATEGY", "Four High-Value Clinical Verticals", "", 7)
    cards = [
        ("Cardiac Tissue", "CVD kills ~2,500 Americans daily. Adult heart tissue cannot regenerate. Redwire returned the first live human cardiac tissue from orbit in May 2024 under BFF-Cardiac.", "Proven by: Redwire BFF-Cardiac, ISS, May 2024", "cardiac"),
        ("Orthopedic Grafts", "Meniscal tears are among the most common U.S. military injuries. Redwire printed the first human knee meniscus in orbit in September 2023, establishing the foundational protocol.", "Proven by: Redwire BFF-Meniscus-2, ISS, Sept 2023", "ortho"),
        ("Artificial Retinas", "30 million people worldwide suffer from degenerative retinal disease. LambdaVision has produced consistent 200-layer protein retina films in microgravity across 9 ISS missions.", "Proven by: LambdaVision / Space Tango, 9 ISS missions", "retina"),
        ("Wound Care Patches", "ESA and DLR demonstrated a handheld bioprinter creating customised skin patches aboard the ISS during the Cosmic Kiss mission. Lowest regulatory burden - Samuel's Phase 1 entry point.", "Proven by: ESA/DLR Bioprint FirstAid, ISS, Dec 2021", "wound"),
    ]
    positions = [(0.78, 1.45), (6.18, 1.45), (0.78, 4.0), (6.18, 4.0)]
    for (title, body, source, icon_name), (x, y) in zip(cards, positions):
        add_box(slide, x, y, 5.37, 2.25, fill=WHITE, line=LINE)
        add_box(slide, x + 0.18, y + 0.2, 0.78, 0.78, fill=SOFT, line=SOFT)
        slide.shapes.add_picture(str(assets[icon_name]), emu(x + 0.37), emu(y + 0.39), height=emu(0.4))
        add_textbox(slide, x + 1.08, y + 0.24, 4.0, 0.24, title, size=13, color=NAVY, bold=True)
        add_textbox(slide, x + 1.08, y + 0.6, 4.0, 0.96, body, size=10, color=TEXT)
        add_textbox(slide, x + 1.08, y + 1.72, 4.0, 0.18, source, size=8.8, color=TEAL_DARK, bold=True)


def add_validation_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "THIRD-PARTY VALIDATION", "Eight Years of Proven In-Orbit Bioprinting", "Every milestone below was publicly announced by an independent organisation - not by Samuel Space. They validate the market premise and demonstrate commercial feasibility.", 8)
    add_line(slide, 0.95, 3.77, 11.35, 0.04, color=LINE)
    milestones = [
        ("2018", "ROSCOSMOS Organ.Aut", "First magnetic 3D bioprinter on ISS; tissue constructs created 2018-2020 prove the in-orbit approach."),
        ("2019", "Redwire BFF Launched", "BioFabrication Facility arrives on ISS; later wins Popular Science Best of What's New 2023."),
        ("Dec '21", "ESA Bioprint FirstAid", "Portable handheld bioprinter demonstrated by astronaut Matthias Maurer on the Cosmic Kiss mission."),
        ("Sep '23", "First Knee Meniscus", "Redwire BFF prints the first human knee meniscus in orbit - a musculoskeletal milestone."),
        ("May '24", "First Cardiac Tissue", "Redwire returns the first live 3D-bioprinted human cardiac tissue from ISS under BFF-Cardiac."),
        ("Sep '25", "Grants & Contracts", "LambdaVision wins NASA InSPA Phase 2; Redwire is awarded a $25M NASA IDIQ for continued ISS work."),
        ("2026", "Commercial Pivot", "LambdaVision closes a $7M seed; commercial station payload slots are reserved ahead of ISS retirement."),
    ]
    top_positions = [0.9, 3.95, 7.0, 10.05]
    bottom_positions = [2.43, 5.48, 8.53]
    for idx, x in enumerate(top_positions):
        year, title, body = milestones[idx]
        add_circle(slide, x + 1.32, 3.66, 0.15, fill=TEAL, line=TEAL)
        add_box(slide, x, 1.6, 2.6, 1.7, fill=WHITE, line=LINE)
        add_pill(slide, x + 0.15, 1.78, 0.72, 0.28, year)
        add_textbox(slide, x + 0.15, 2.16, 2.28, 0.26, title, size=11.1, color=NAVY, bold=True)
        add_textbox(slide, x + 0.15, 2.48, 2.28, 0.62, body, size=9.2, color=TEXT)
    for idx, x in enumerate(bottom_positions, start=4):
        year, title, body = milestones[idx]
        add_circle(slide, x + 1.32, 3.66, 0.15, fill=TEAL, line=TEAL)
        add_box(slide, x, 4.15, 2.6, 1.7, fill=WHITE, line=LINE)
        add_pill(slide, x + 0.15, 4.33, 0.78, 0.28, year)
        add_textbox(slide, x + 0.15, 4.71, 2.28, 0.26, title, size=11.1, color=NAVY, bold=True)
        add_textbox(slide, x + 0.15, 5.03, 2.28, 0.62, body, size=9.2, color=TEXT)


def add_market_opportunity_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "MARKET OPPORTUNITY", "A Rapidly Expanding Global Addressable Market", "", 9)
    add_box(slide, 0.78, 1.36, 6.15, 4.95, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 1.56, 4.9, 0.2, "Global 3D Bioprinting Market 2023-2030 (USD $B)", size=11.5, color=NAVY, bold=True)
    add_column_chart(slide, 1.0, 1.9, 5.75, 3.88, ["2023", "2024", "2025E", "2026E", "2027E", "2028E", "2029E", "2030E"], [1.28, 1.46, 1.7, 1.97, 2.28, 2.64, 3.06, 3.49])
    add_textbox(slide, 1.0, 5.82, 5.7, 0.18, "Source: Mordor Intelligence, 3D Bioprinting Market Report, Jan 2026  |  ~15.9% CAGR", size=8.5, color=MUTED)
    cards = [
        ("$3.5B+", "Projected global 3D bioprinting market by 2030 at ~15.9% CAGR (Mordor Intelligence, Jan 2026)"),
        ("Dec 2024", "FDA cleared Symvess - the first tissue-engineered vascular graft - setting a direct regulatory precedent"),
        ("$420M+", "Committed by pharma companies to commercial-station research programs through 2030 (industry research)"),
        ("$200-400M", "Samuel's estimated initial SAM by 2030 (cardiac patches, ortho grafts, retinal constructs)"),
    ]
    for idx, (value, body) in enumerate(cards):
        add_info_card(slide, 7.15, 1.45 + idx * 1.2, 5.4, 0.95, value, body, fill=WHITE)


def add_market_timing_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "MARKET TIMING", "Four Converging Tailwinds Create a Unique Entry Window", "", 10)
    cards = [
        ("~89%", "Launch Costs Have Collapsed", "Cost-to-orbit has dropped from ~$54,000/kg (Space Shuttle era) to ~$6,000/kg (Falcon 9). Starship is expected to reduce this further. Research payloads that were cost-prohibitive a decade ago are now commercially viable."),
        ("4", "Commercial Stations Are Coming", "ISS retires around 2030. NASA's CLD program has funded four successors: Starlab, Axiom Station, Orbital Reef, and Vast Haven-1. Samuel is designed to remain station-agnostic from day one."),
        ("Dec 2024", "Regulatory Precedent Is Set", "The FDA cleared Symvess - the first tissue-engineered synthetic vascular graft - in December 2024. This is the most important regulatory precedent for clinical bioprinted tissue and validates the FDA's review pathway."),
        ("$329M", "Capital Is Already Flowing", "Varda Space Industries has raised $329M total, including a $187M Series C in July 2025. LambdaVision closed a $7M seed atop $15M+ in non-dilutive grants."),
    ]
    positions = [(0.78, 1.45), (6.18, 1.45), (0.78, 4.05), (6.18, 4.05)]
    for (value, heading, body), (x, y) in zip(cards, positions):
        add_box(slide, x, y, 5.37, 2.18, fill=WHITE, line=LINE)
        add_pill(slide, x + 0.18, y + 0.18, 0.9, 0.3, value)
        add_textbox(slide, x + 0.18, y + 0.62, 4.8, 0.25, heading, size=13, color=NAVY, bold=True)
        add_textbox(slide, x + 0.18, y + 0.98, 4.95, 0.95, body, size=9.8, color=TEXT)


def add_platform_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "OUR PLATFORM", "Station-Agnostic Payload Platform", "", 11)
    add_box(slide, 0.78, 1.32, 11.77, 0.72, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 1.5, 11.3, 0.34, "Samuel Space is built as a payload integrator and tissue manufacturing platform - hardware-compatible across all emerging commercial LEO stations. The same flight-qualified platform, bioink libraries, and protocols deploy whether the payload flies on the ISS, Starlab, Axiom, Orbital Reef, or Vast.", size=10.1, color=TEXT)
    pillars = [
        ("Modular Payload Design", "Hardware footprint matches standard double-locker form factor compatible with ISS, Starlab, and Axiom interfaces - one build, many stations."),
        ("Proprietary Cell-Line IP", "Bioink formulations and printing protocols optimised for microgravity are the core defensible IP asset, accumulated across every mission flown."),
        ("Service-First Revenue", "Fee-for-service bioprinting generates near-term revenue from pharma and academic customers while building the IP base that creates long-term enterprise value."),
    ]
    for idx, (title, body) in enumerate(pillars):
        add_info_card(slide, 0.78, 2.25 + idx * 1.36, 5.2, 1.15, title, body, fill=WHITE)
    add_box(slide, 6.2, 2.25, 6.35, 4.43, fill=WHITE, line=LINE)
    add_textbox(slide, 6.45, 2.44, 3.2, 0.22, "STATION ECOSYSTEM COMPATIBILITY", size=11.5, color=NAVY, bold=True)
    add_textbox(slide, 6.45, 2.72, 3.4, 0.18, "Engineered for compatibility with next-generation orbital infrastructure.", size=9.5, color=MUTED)
    slide.shapes.add_picture(str(assets["platform"]), emu(6.45), emu(3.02), width=emu(5.85))
    ecosystem = [
        ("ISS", "Redwire BFF - active now through ~2030"),
        ("Starlab", "Voyager/Airbus JV - targeting 2029 launch"),
        ("Axiom Station", "First commercial modules ~2026-2028"),
        ("Orbital Reef", "Blue Origin/Sierra Space - CLD-funded"),
        ("Vast Haven-1", "First dedicated commercial station, 2026"),
    ]
    base_y = 5.12
    for idx, (name, detail) in enumerate(ecosystem):
        row_y = base_y + idx * 0.27
        add_textbox(slide, 6.45, row_y, 1.55, 0.16, name, size=8.8, color=NAVY, bold=True)
        add_textbox(slide, 7.95, row_y, 4.18, 0.16, detail, size=8.6, color=MUTED)
    add_textbox(slide, 6.45, 6.52, 5.7, 0.14, "Station names shown as market context only - not as confirmed Samuel partnerships.", size=7.8, color=MUTED)


def add_execution_plan_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "EXECUTION PLAN", "Four-Phase Technology & Commercialisation Roadmap", "", 12)
    phases = [
        ("PHASE 1", "Years 1-2", "Platform & Wound Care", ["ISS payload hardware certification", "Commercial wound-care patch trials", "NASA InSPA Phase 1 grant applications", "Research tissue models for pharma CROs"], "TRL 4 -> 6"),
        ("PHASE 2", "Years 3-5", "Orthopedic & Cardiac", ["510(k) filing for wound-care products", "Orthopedic graft BFF-compatible program", "Cardiac patch preclinical studies", "Series B fundraise"], "TRL 6 -> 8"),
        ("PHASE 3", "Years 6-9", "Retinal & Vascular", ["Clinical trials for retinal constructs", "FDA PMA pathway engagement", "Expand to Orbital Reef / Axiom", "Series C / pre-IPO raise"], "TRL 7 -> 9"),
        ("PHASE 4", "Year 10+", "Whole-Organ Platform", ["Vascularised full-organ constructs", "Commercial station scale-up", "Platform licensing to pharma majors", "IPO or strategic exit"], "TRL 9+"),
    ]
    card_w = 2.85
    fills = [WHITE, WHITE, WHITE, WHITE]
    for idx, (phase, years, focus, bullets, trl) in enumerate(phases):
        x = 0.78 + idx * 3.0
        add_box(slide, x, 1.45, card_w, 5.2, fill=fills[idx], line=LINE)
        add_pill(slide, x + 0.16, 1.66, 0.86, 0.28, phase)
        add_textbox(slide, x + 0.16, 2.08, 1.55, 0.18, years, size=9.2, color=MUTED, bold=True)
        add_textbox(slide, x + 0.16, 2.38, 2.45, 0.28, focus, size=12.2, color=NAVY, bold=True)
        add_multiline_textbox(slide, x + 0.16, 2.82, 2.45, 2.8, bullets, size=9.6, color=TEXT, bullet=True)
        add_box(slide, x + 0.16, 6.0, 1.12, 0.28, fill=SOFT_BLUE, line=SOFT_BLUE)
        add_textbox(slide, x + 0.24, 6.07, 1.0, 0.12, trl, size=8.8, color=NAVY, bold=True)


def add_competitive_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "COMPETITIVE LANDSCAPE", "Samuel Space Targets the Space-Based Clinical Quadrant", "", 13)
    add_box(slide, 1.05, 1.45, 11.0, 4.85, fill=WHITE, line=LINE)
    add_line(slide, 1.4, 3.88, 10.3, 0.03, color=LINE)
    add_line(slide, 6.55, 1.8, 0.03, 4.15, color=LINE)
    add_textbox(slide, 1.18, 1.66, 2.2, 0.16, "EARTH-BASED", size=9.5, color=MUTED, bold=True)
    add_textbox(slide, 6.76, 1.66, 2.2, 0.16, "SPACE-BASED", size=9.5, color=MUTED, bold=True)
    add_textbox(slide, 1.1, 1.98, 1.9, 0.16, "CLINICAL / COMMERCIAL", size=9.5, color=MUTED, bold=True)
    add_textbox(slide, 1.1, 5.78, 1.55, 0.16, "EARLY RESEARCH", size=9.5, color=MUTED, bold=True)
    points = [
        (2.2, 2.8, "3D Systems (Med)", NAVY),
        (1.7, 3.8, "Cellink / BICO", NAVY),
        (2.0, 4.7, "Organovo", NAVY),
        (8.7, 3.95, "Redwire BFF", NAVY),
        (10.0, 3.05, "LambdaVision", NAVY),
        (7.6, 5.0, "ROSCOSMOS / 3D BPS", NAVY),
        (9.95, 2.2, "Samuel Space (target)", TEAL_DARK),
    ]
    for x, y, label, color in points:
        size = 0.18 if "Samuel" not in label else 0.24
        add_circle(slide, x, y, size, fill=color, line=color)
        add_textbox(slide, x + 0.24, y - 0.02, 1.95, 0.22, label, size=9.4, color=TEXT, bold="Samuel" in label)
    add_textbox(slide, 1.05, 6.52, 7.2, 0.14, "Positions are illustrative, based on public technology and commercialisation-stage disclosures.", size=8, color=MUTED)


def add_revenue_model_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "REVENUE MODEL", "Three Complementary Revenue Streams", "", 14)
    streams = [
        ("REVENUE STREAM 1", "Contract Bioprinting Services", "Fee-for-service orbital bioprinting missions for pharmaceutical companies, academic medical centres, and government research agencies, modelled on how Redwire's BFF and Space Tango currently operate commercially. Revenue begins in Year 1.", "Near-Term Revenue"),
        ("REVENUE STREAM 2", "Non-Dilutive Grant Funding", "NASA InSPA, ESA BIC, DoD SBIR/STTR, and National Eye Institute mechanisms provide substantial non-dilutive capital. LambdaVision has secured $15M+ cumulatively; Redwire holds a $25M NASA IDIQ. Samuel targets $4-8M in grants across Years 1-4.", "Reduces Equity Burn"),
        ("REVENUE STREAM 3", "Tissue IP Licensing & Royalties", "As proprietary bioink formulations and protocols demonstrate clinical-grade output, Samuel licenses IP to pharmaceutical and device manufacturers. This becomes the primary long-term value driver, targeted from Year 3 onward.", "Long-Term Value Driver"),
    ]
    for idx, (label, title, body, footer) in enumerate(streams):
        x = 0.78 + idx * 4.02
        add_box(slide, x, 1.62, 3.75, 4.6, fill=WHITE, line=LINE)
        add_pill(slide, x + 0.18, 1.83, 1.22, 0.28, label)
        add_textbox(slide, x + 0.18, 2.28, 3.1, 0.28, title, size=13, color=NAVY, bold=True)
        add_textbox(slide, x + 0.18, 2.72, 3.18, 2.45, body, size=10, color=TEXT)
        add_box(slide, x + 0.18, 5.48, 1.76, 0.32, fill=SOFT_BLUE, line=SOFT_BLUE)
        add_textbox(slide, x + 0.26, 5.55, 1.6, 0.14, footer, size=8.8, color=NAVY, bold=True)


def add_go_to_market_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "GO-TO-MARKET", "Sequenced Market Entry Across Four Partner Categories", "", 15)
    cards = [
        ("1", "Commercial Station Operators", "Redwire (active BFF operator), Starlab (Voyager/Airbus JV), Axiom Space, Blue Origin Orbital Reef, and Vast Haven-1. Samuel negotiates research-service agreements and payload manifesting rights with each.", "Redwire, Starlab, Axiom, Vast"),
        ("2", "Payload Integrators & CROs", "Partners such as Space Tango, LambdaVision's enabling integrator, provide environmental control, crew interface, and downlink services. Samuel outsources non-core integration to move faster and preserve capital.", "Space Tango, Nanoracks/Voyager"),
        ("3", "Pharma & Academic Medical Centres", "Initial contract-service revenue comes from pharmaceutical R&D divisions seeking orbital tissue models for drug screening, and from NIH-funded academic biotech groups with institutional mandates.", "Pharma R&D, NIH-funded labs"),
        ("4", "Government & Grant Agencies", "NASA (InSPA, SBIR/STTR), ESA, NIH National Eye Institute, and DoD MTEC provide non-dilutive grant capital that de-risks early development and validates the technology for commercial buyers.", "NASA, ESA, NIH NEI, DoD MTEC"),
    ]
    positions = [(0.78, 1.58), (6.18, 1.58), (0.78, 4.1), (6.18, 4.1)]
    for (num, title, body, footer), (x, y) in zip(cards, positions):
        add_box(slide, x, y, 5.37, 2.18, fill=WHITE, line=LINE)
        add_circle(slide, x + 0.2, y + 0.18, 0.36, fill=TEAL, line=TEAL)
        add_textbox(slide, x + 0.2, y + 0.26, 0.36, 0.12, num, size=10, color=WHITE, bold=True, align=PP_ALIGN.CENTER, margin=0)
        add_textbox(slide, x + 0.68, y + 0.18, 4.25, 0.25, title, size=12.8, color=NAVY, bold=True)
        add_textbox(slide, x + 0.2, y + 0.62, 4.95, 0.95, body, size=9.9, color=TEXT)
        add_pill(slide, x + 0.2, y + 1.72, 2.2, 0.3, footer, fill=SOFT_BLUE, color=NAVY, align=PP_ALIGN.LEFT)


def add_financials_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "FINANCIALS - ILLUSTRATIVE", "Five-Year Revenue Model (Illustrative Projections)", "Figures are illustrative projections for investor discussion purposes only. Actual results will vary materially.", 16)
    add_box(slide, 0.78, 1.42, 6.4, 5.05, fill=WHITE, line=LINE)
    add_stacked_chart(
        slide,
        1.0,
        1.76,
        5.95,
        4.25,
        ["Yr 1", "Yr 2", "Yr 3", "Yr 4", "Yr 5"],
        [
            ("Contract Services ($M)", [0.5, 1.2, 3.0, 7.0, 14.0]),
            ("Grant Revenue ($M)", [0.3, 1.0, 1.5, 2.0, 2.5]),
            ("IP Licensing ($M)", [0.0, 0.2, 1.6, 4.5, 7.7]),
        ],
    )
    add_box(slide, 7.35, 1.42, 5.2, 5.05, fill=WHITE, line=LINE)
    rows = [
        ["Year", "Total Rev.", "Key Driver"],
        ["Year 1", "$0.8M", "First service contracts + NASA grant"],
        ["Year 2", "$2.4M", "Expanded services + IP licensing launch"],
        ["Year 3", "$6.1M", "Cardiac/ortho milestones + Starlab start"],
        ["Year 4", "$13.5M", "IP licensing deals + multi-station ops"],
        ["Year 5", "$24.2M", "Full vertical scale-up across four verticals"],
    ]
    add_manual_table(slide, 7.55, 1.74, [1.0, 1.0, 2.95], rows, [0.42, 0.58, 0.58, 0.58, 0.58, 0.58])


def add_ask_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "THE ASK", "An $18M Series A to Reach First Commercial Flight", "", 17)
    add_box(slide, 0.78, 1.45, 3.15, 5.0, fill=WHITE, line=LINE)
    add_textbox(slide, 1.04, 1.72, 1.4, 0.18, "ROUND SIZE", size=10.5, color=TEAL_DARK, bold=True)
    add_textbox(slide, 1.04, 2.0, 2.1, 0.6, "$18M", size=30, color=NAVY, bold=True)
    add_textbox(slide, 1.04, 2.62, 2.1, 0.2, "Series A Preferred Equity", size=10.4, color=TEXT, bold=True)
    metrics = [
        ("~30 months", "Runway to next value-inflection milestone"),
        ("2", "Product verticals reaching clinical pilot stage"),
        ("3-5", "Targeted commercial flight missions secured"),
    ]
    base_y = 3.2
    for idx, (value, body) in enumerate(metrics):
        add_box(slide, 1.04, base_y + idx * 0.92, 2.6, 0.72, fill=SOFT_BLUE, line=SOFT_BLUE)
        add_textbox(slide, 1.18, base_y + 0.11 + idx * 0.92, 0.95, 0.18, value, size=12.5, color=NAVY, bold=True)
        add_textbox(slide, 2.1, base_y + 0.11 + idx * 0.92, 1.4, 0.3, body, size=8.8, color=TEXT)
    add_box(slide, 4.18, 1.45, 3.7, 5.0, fill=WHITE, line=LINE)
    add_textbox(slide, 4.44, 1.72, 2.0, 0.18, "Use of Funds", size=10.5, color=TEAL_DARK, bold=True)
    add_doughnut_chart(
        slide,
        4.6,
        2.02,
        2.85,
        2.55,
        ["R&D / Payload Dev", "Flight & Integration", "Regulatory & Clinical", "Team Build-Out", "G&A"],
        [40.0, 25.0, 15.0, 15.0, 5.0],
    )
    labels = [
        ("R&D / Payload Development", TEAL),
        ("Flight & Integration Costs", NAVY),
        ("Regulatory & Clinical Prep", SERIES_BLUE),
        ("Team Build-Out", SERIES_LIGHT),
        ("G&A", SOFT_NAVY),
    ]
    legend_y = 4.85
    for idx, (label, color) in enumerate(labels):
        row_y = legend_y + idx * 0.26
        add_circle(slide, 4.5, row_y, 0.12, fill=color, line=color)
        add_textbox(slide, 4.68, row_y - 0.01, 2.65, 0.14, label, size=8.6, color=TEXT)
    add_box(slide, 8.12, 1.45, 4.43, 5.0, fill=WHITE, line=LINE)
    add_textbox(slide, 8.38, 1.72, 2.6, 0.18, "Key Use-of-Funds Milestones", size=10.5, color=TEAL_DARK, bold=True)
    add_multiline_textbox(
        slide,
        8.34,
        2.08,
        3.85,
        3.65,
        [
            "Flight-qualify modular payload hardware",
            "Complete 3-5 commissioned ISS missions",
            "File first NASA InSPA Phase 1 grant",
            "Secure Starlab payload reservation",
            "Build core scientific and engineering team",
        ],
        size=10,
        color=TEXT,
        bullet=True,
    )
    add_textbox(slide, 0.78, 6.72, 7.0, 0.14, "Allocation figures are illustrative and subject to change based on diligence and negotiated terms.", size=8.1, color=MUTED)


def add_risk_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "RISK FACTORS", "Key Risks and Mitigating Strategies", "", 18)
    rows = [
        ["Risk Category", "Risk", "Mitigant"],
        ["Technical", "Scaling from validated tissue patches to full vascularised organs remains scientifically unsolved industry-wide.", "Phased roadmap monetises near-term, lower-complexity verticals (wound care, orthopedic) before attempting whole-organ printing."],
        ["Regulatory", "FDA approval pathways for novel tissue-engineered products are long, costly, and precedent-limited.", "Symvess clearance (Dec 2024) establishes the first direct precedent. Samuel engages regulatory counsel from Year 1 and targets lower-risk 510(k) products first."],
        ["Platform / Access", "ISS retirement (~2030) creates dependency risk if commercial successor stations are delayed.", "Station-agnostic hardware design supports ISS, Starlab, Axiom, Orbital Reef, and Vast simultaneously, diversifying single-platform exposure."],
        ["Capital Intensity", "Flight missions, payload hardware, and clinical trials require sustained, significant capital investment.", "Blended funding strategy combines equity with non-dilutive NASA/ESA/DoD grants, following the precedent set by LambdaVision and Redwire."],
        ["Competitive", "Incumbent operators (Redwire, LambdaVision) hold first-mover relationships with NASA and station operators.", "Samuel differentiates on multi-vertical platform breadth and station-agnostic flexibility rather than competing head-on for any single contract."],
    ]
    add_manual_table(slide, 0.78, 1.48, [2.0, 4.25, 5.52], rows, [0.44, 0.92, 0.92, 0.92, 0.92, 0.92])


def add_organization_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "ORGANISATION", "Team & Scientific Advisory Approach", "", 19)
    add_box(slide, 0.78, 1.32, 11.77, 0.72, fill=WHITE, line=LINE)
    add_textbox(slide, 1.0, 1.5, 11.3, 0.34, "Samuel Space is building a founding team and Scientific Advisory Board around the specific expertise this platform requires. Roles below reflect the Company's hiring plan and target advisory profile for the Series A period.", size=10.2, color=TEXT)
    cards = [
        ("Founder / Chief Executive Officer", "Commercial space or biotech venture leadership; track record building regulated hardware or therapeutics businesses from seed to scale."),
        ("Chief Scientific Officer", "PhD-level expertise in tissue engineering, stem cell biology, or regenerative medicine; prior peer-reviewed bioprinting research."),
        ("VP, Payload Engineering", "Aerospace hardware engineering background; experience flight-qualifying instrumentation for ISS or commercial LEO payloads."),
        ("Head of Regulatory Affairs", "FDA medical device / biologics regulatory pathway expertise; experience navigating 510(k) and PMA submissions."),
    ]
    positions = [(0.78, 2.3), (6.18, 2.3), (0.78, 4.2), (6.18, 4.2)]
    for (title, body), (x, y) in zip(cards, positions):
        add_info_card(slide, x, y, 5.37, 1.62, title, body, fill=WHITE)
    add_box(slide, 0.78, 6.0, 11.77, 0.7, fill=SOFT_BLUE, line=SOFT_BLUE)
    add_textbox(slide, 1.0, 6.15, 2.8, 0.18, "Scientific Advisory Board - In Formation", size=10.8, color=NAVY, bold=True)
    add_textbox(slide, 1.0, 6.35, 11.25, 0.18, "Samuel is actively recruiting an SAB spanning stem-cell and cardiac tissue biology, microgravity payload design, and protein-based biomanufacturing. Target SAB appointments will be announced as the round closes.", size=9.2, color=TEXT)


def add_closing_slide(prs, assets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    set_background(slide, WHITE)
    slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, emu(8.4), 0, emu(4.933), emu(7.5)).fill.solid()
    hero_panel = slide.shapes[-1]
    hero_panel.fill.fore_color.rgb = SOFT_BLUE
    hero_panel.line.fill.background()
    slide.shapes.add_picture(str(assets["close"]), emu(8.52), emu(0), height=emu(7.5))
    add_line(slide, 0, 0, 13.333, 0.08, color=TEAL)
    add_textbox(slide, 0.78, 0.58, 2.5, 0.2, "THE OPPORTUNITY", size=10, color=TEAL_DARK, bold=True)
    add_textbox(slide, 0.78, 1.3, 6.85, 1.5, "Every organ printed in orbit\nis a life that didn't have to wait.", size=26.5, color=NAVY, bold=True, font_name="Aptos Display")
    add_textbox(slide, 0.78, 3.28, 6.9, 1.15, "Samuel Space stands at the intersection of two converging trends: a mature, validated microgravity bioprinting science base, and the arrival of commercial space stations built for exactly this kind of work. The science already works. What's needed now is a company built to commercialise it.", size=11.4, color=TEXT)
    add_box(slide, 0.78, 5.55, 4.4, 0.78, fill=SOFT, line=SOFT)
    add_textbox(slide, 1.0, 5.74, 3.9, 0.16, "Samuel Space Organ and Tissue Engineering", size=10.4, color=NAVY, bold=True)
    add_textbox(slide, 1.0, 5.98, 3.9, 0.16, "investors@samuelspace.com  (placeholder)", size=9.6, color=MUTED)
    add_textbox(slide, 12.35, 7.0, 0.4, 0.18, "20", size=8.5, color=MUTED, align=PP_ALIGN.RIGHT, margin=0)


def add_appendix_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    remove_defaults(slide)
    add_header(slide, "APPENDIX", "Sources & References", "", 21)
    add_box(slide, 0.78, 1.42, 5.65, 4.55, fill=WHITE, line=LINE)
    add_box(slide, 6.9, 1.42, 5.65, 4.55, fill=WHITE, line=LINE)
    add_textbox(slide, 1.02, 1.64, 3.0, 0.18, "Bioprinting Science & Technology", size=11.2, color=NAVY, bold=True)
    add_multiline_textbox(
        slide,
        1.0,
        1.94,
        5.15,
        3.7,
        [
            'NASA - "3D Bioprinting" research overview, iss-research program page',
            "NASA - BioFabrication Facility mission page, Redwire Corporation",
            "Redwire Corporation - BFF-Meniscus first human knee meniscus printed in orbit, Sept 2023",
            "Redwire Corporation - BFF-Cardiac investigation announcements and ISS return coverage, 2024",
            "Redwire Corporation - NASA IDIQ contract award coverage, Sept 2025",
            "ESA / DLR - Bioprint FirstAid investigation overview, Cosmic Kiss mission materials",
            "LambdaVision Inc. - Press releases on Space Tango partnership, NASA InSPA Phase 2 award, and seed financing",
        ],
        size=8.9,
        color=TEXT,
        bullet=True,
    )
    add_textbox(slide, 7.14, 1.64, 3.2, 0.18, "Market, Regulatory & Industry Data", size=11.2, color=NAVY, bold=True)
    add_multiline_textbox(
        slide,
        7.12,
        1.94,
        5.15,
        3.7,
        [
            "ROSCOSMOS - 3D MBP / Organ.Aut magnetic bioprinter mission summaries, 2018-2020",
            "American Heart Association - 2025 Heart Disease & Stroke Statistics Update",
            "OPTN / organdonor.gov (HRSA) - U.S. transplant waiting list data, 2025",
            "World Health Organization - Global Observatory on Donation and Transplantation",
            "Mordor Intelligence - 3D Bioprinting Market Report, updated January 2026",
            "U.S. FDA - Symvess clearance announcement, Dec 2024",
            "NASA Commercial LEO Destinations program pages covering Starlab, Axiom, Orbital Reef, and Vast",
        ],
        size=8.9,
        color=TEXT,
        bullet=True,
    )
    add_box(slide, 0.78, 6.08, 11.77, 0.7, fill=SOFT_BLUE, line=SOFT_BLUE)
    add_textbox(slide, 1.0, 6.23, 1.9, 0.18, "A Note on Methodology", size=10.8, color=NAVY, bold=True)
    add_textbox(slide, 2.52, 6.2, 9.7, 0.22, "All third-party scientific milestones, dates, and figures in this presentation are drawn from the public sources listed above and were current as of the citation date shown. References to NASA, Redwire, LambdaVision, ROSCOSMOS, and ESA/DLR describe industry validation of the underlying science, not partnerships, endorsements, or affiliations with Samuel Space.", size=8.8, color=TEXT)


def main():
    assets = extract_assets()
    prs = Presentation()
    prs.slide_width = emu(13.333333333)
    prs.slide_height = emu(7.5)

    add_cover_slide(prs, assets)
    add_notice_slide(prs)
    add_exec_summary_slide(prs)
    add_problem_slide(prs)
    add_constraint_slide(prs, assets)
    add_technology_slide(prs, assets)
    add_product_strategy_slide(prs, assets)
    add_validation_slide(prs)
    add_market_opportunity_slide(prs)
    add_market_timing_slide(prs)
    add_platform_slide(prs, assets)
    add_execution_plan_slide(prs)
    add_competitive_slide(prs)
    add_revenue_model_slide(prs)
    add_go_to_market_slide(prs)
    add_financials_slide(prs)
    add_ask_slide(prs)
    add_risk_slide(prs)
    add_organization_slide(prs)
    add_closing_slide(prs, assets)
    add_appendix_slide(prs)

    temp_path = DECK_PATH.with_suffix(".tmp.pptx")
    prs.save(temp_path)
    temp_path.replace(DECK_PATH)


if __name__ == "__main__":
    main()
