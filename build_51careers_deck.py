#!/usr/bin/env python3
"""Rebuild 51Careers.AI investor presentation — institutional visual system."""

from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

OUT = Path("/workspace/51Careers_AI_Investor_Presentation.pptx")
SW, SH = Inches(13.333), Inches(7.5)

# Typography — different from Century Schoolbook / Calibri
FONT_DISPLAY = "Georgia"
FONT_BODY = "Segoe UI"

# Palette — institutional navy + deep teal (no purple / cream clichés)
NAVY = RGBColor(0x0B, 0x1F, 0x33)
INK = RGBColor(0x14, 0x28, 0x3C)
BODY = RGBColor(0x2C, 0x3A, 0x4B)
MUTED = RGBColor(0x5E, 0x6E, 0x80)
LINE = RGBColor(0xD4, 0xDE, 0xE8)
SOFT = RGBColor(0xEE, 0xF3, 0xF8)
PANEL = RGBColor(0xF5, 0xF8, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
TEAL = RGBColor(0x0A, 0x6B, 0x63)
TEAL_DK = RGBColor(0x07, 0x52, 0x4C)
TEAL_SOFT = RGBColor(0xD7, 0xEB, 0xE8)
BLUE = RGBColor(0x2F, 0x5D, 0x8A)
BRONZE = RGBColor(0x8A, 0x6D, 0x3B)
SLATE = RGBColor(0x6B, 0x7C, 0x8F)
CHART_COLORS = [TEAL, BLUE, BRONZE, SLATE, RGBColor(0xA8, 0xB8, 0xC6)]


def rgb_hex(c: RGBColor) -> str:
    return f"{c[0]:02X}{c[1]:02X}{c[2]:02X}"


def set_run(run, text, font=FONT_BODY, size=11, bold=False, color=BODY, italic=False):
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color


def clear_tf(tf):
    tf.clear()
    p = tf.paragraphs[0]
    p.clear()
    return p


def add_text(shape, paragraphs, valign=MSO_ANCHOR.TOP):
    """paragraphs: list of dicts with text/font/size/bold/color/align/space_after/italic"""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        shape.text_frame._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(valign, "t"))
    except Exception:
        pass
    first = True
    for spec in paragraphs:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = spec.get("align", PP_ALIGN.LEFT)
        if "space_after" in spec:
            p.space_after = Pt(spec["space_after"])
        if "space_before" in spec:
            p.space_before = Pt(spec["space_before"])
        run = p.add_run()
        set_run(
            run,
            spec["text"],
            font=spec.get("font", FONT_BODY),
            size=spec.get("size", 11),
            bold=spec.get("bold", False),
            color=spec.get("color", BODY),
            italic=spec.get("italic", False),
        )


def fill_solid(shape, color):
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    try:
        shape.line.fill.background()
    except (AttributeError, TypeError, ValueError):
        pass


def stroke(shape, color=LINE, width_pt=1.0):
    shape.line.color.rgb = color
    shape.line.width = Pt(width_pt)


def rect(slide, l, t, w, h, fill=None, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, l, t, w, h)
    if fill is None:
        s.fill.background()
    else:
        fill_solid(s, fill)
    if line is None:
        s.line.fill.background()
    else:
        stroke(s, line, line_w)
    return s


def round_rect(slide, l, t, w, h, fill=WHITE, line=None, line_w=1.0):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    fill_solid(s, fill)
    # tighten radius
    try:
        s.adjustments[0] = 0.08
    except Exception:
        pass
    if line is None:
        s.line.fill.background()
    else:
        stroke(s, line, line_w)
    return s


def textbox(slide, l, t, w, h, paragraphs, valign=MSO_ANCHOR.TOP):
    s = slide.shapes.add_textbox(l, t, w, h)
    add_text(s, paragraphs, valign=valign)
    return s


def accent_bar(slide, l, t, w=Inches(0.08), h=Inches(0.35), color=TEAL):
    return rect(slide, l, t, w, h, fill=color)


def footer(slide, page: int, total: int = 15):
    rect(slide, Inches(0), SH - Inches(0.42), SW, Inches(0.42), fill=PANEL)
    rect(slide, Inches(0), SH - Inches(0.42), SW, Pt(1), fill=LINE)
    textbox(
        slide,
        Inches(0.55),
        SH - Inches(0.36),
        Inches(6),
        Inches(0.28),
        [{"text": "51CAREERS.AI   ·   CONFIDENTIAL", "size": 8, "color": MUTED, "bold": True}],
    )
    textbox(
        slide,
        SW - Inches(1.1),
        SH - Inches(0.36),
        Inches(0.7),
        Inches(0.28),
        [{"text": str(page), "size": 9, "color": MUTED, "align": PP_ALIGN.RIGHT}],
    )


def section_header(slide, eyebrow: str, title: str, subtitle: str | None = None):
    rect(slide, Inches(0), Inches(0), SW, Inches(0.08), fill=TEAL)
    textbox(
        slide,
        Inches(0.55),
        Inches(0.28),
        Inches(12),
        Inches(0.28),
        [{"text": eyebrow.upper(), "size": 10, "bold": True, "color": TEAL, "font": FONT_BODY}],
    )
    textbox(
        slide,
        Inches(0.55),
        Inches(0.52),
        Inches(12.2),
        Inches(0.55),
        [{"text": title, "size": 28, "bold": True, "color": NAVY, "font": FONT_DISPLAY}],
    )
    if subtitle:
        textbox(
            slide,
            Inches(0.55),
            Inches(1.05),
            Inches(12.2),
            Inches(0.35),
            [{"text": subtitle, "size": 12, "color": MUTED}],
        )


def style_chart(chart, has_legend=True, point_colors=None):
    chart.has_legend = has_legend
    if has_legend:
        chart.legend.include_in_layout = False
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.font.size = Pt(9)
        chart.legend.font.name = FONT_BODY
        chart.legend.font.color.rgb = MUTED
    try:
        plot = chart.plots[0]
        plot.has_data_labels = False
    except Exception:
        pass
    # color series
    try:
        for i, series in enumerate(chart.series):
            color = CHART_COLORS[i % len(CHART_COLORS)]
            solid = series.format.fill
            solid.solid()
            solid.fore_color.rgb = color
            try:
                series.format.line.color.rgb = color
            except Exception:
                pass
            if point_colors:
                for j, pt in enumerate(series.points):
                    c = point_colors[j % len(point_colors)]
                    pt.format.fill.solid()
                    pt.format.fill.fore_color.rgb = c
    except Exception:
        pass
    # axes
    try:
        cat = chart.category_axis
        cat.has_major_gridlines = False
        cat.tick_labels.font.size = Pt(9)
        cat.tick_labels.font.name = FONT_BODY
        cat.tick_labels.font.color.rgb = MUTED
        cat.format.line.color.rgb = LINE
    except Exception:
        pass
    try:
        val = chart.value_axis
        val.has_major_gridlines = True
        val.major_gridlines.format.line.color.rgb = LINE
        val.tick_labels.font.size = Pt(9)
        val.tick_labels.font.name = FONT_BODY
        val.tick_labels.font.color.rgb = MUTED
        val.format.line.color.rgb = LINE
    except Exception:
        pass


def add_chart(slide, chart_type, left, top, width, height, categories, series_data, legend=True, point_colors=None):
    """series_data: list of (name, values)"""
    data = CategoryChartData()
    data.categories = categories
    for name, values in series_data:
        data.add_series(name, values)
    chart_shape = slide.shapes.add_chart(chart_type, left, top, width, height, data)
    style_chart(chart_shape.chart, has_legend=legend, point_colors=point_colors)
    return chart_shape


def blank_slide(prs):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    fill_solid(slide.background, PANEL)
    return slide


# ───────────────────────────── SLIDES ─────────────────────────────


def slide_cover(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_solid(slide.background, NAVY)
    # atmospheric bands
    rect(slide, Inches(0), Inches(0), SW, Inches(0.12), fill=TEAL)
    rect(slide, Inches(0), SH - Inches(1.35), SW, Inches(1.35), fill=RGBColor(0x08, 0x18, 0x28))
    # left accent panel
    rect(slide, Inches(0), Inches(0), Inches(0.18), SH, fill=TEAL)

    textbox(
        slide,
        Inches(0.7),
        Inches(0.55),
        Inches(8),
        Inches(0.3),
        [{"text": "CONFIDENTIAL INVESTOR PRESENTATION", "size": 11, "bold": True, "color": TEAL_SOFT}],
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(1.35),
        Inches(10),
        Inches(0.85),
        [{"text": "51CAREERS.AI", "size": 48, "bold": True, "color": WHITE, "font": FONT_DISPLAY}],
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(2.2),
        Inches(10),
        Inches(0.4),
        [{"text": "The AI-Powered Global Career Platform", "size": 18, "color": RGBColor(0xC5, 0xD5, 0xE0)}],
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(3.0),
        Inches(8),
        Inches(0.35),
        [{"text": "Growth Financing", "size": 14, "bold": True, "color": TEAL}],
    )
    textbox(
        slide,
        Inches(0.7),
        Inches(3.35),
        Inches(9),
        Inches(0.35),
        [{"text": "US$5.0 Million Private Placement   ·   July 2026", "size": 14, "color": WHITE}],
    )

    # KPI strip
    kpis = [
        ("2016", "FOUNDED · NEW YORK CITY"),
        ("8+", "COUNTRIES OF OPERATION"),
        ("US$10M+", "ANNUAL REVENUE (EST.)†"),
    ]
    x = Inches(0.7)
    for val, label in kpis:
        card = round_rect(slide, x, Inches(4.2), Inches(3.5), Inches(1.35), fill=RGBColor(0x12, 0x2A, 0x40))
        textbox(
            slide,
            x + Inches(0.25),
            Inches(4.35),
            Inches(3.0),
            Inches(0.6),
            [{"text": val, "size": 32, "bold": True, "color": WHITE, "font": FONT_DISPLAY}],
        )
        textbox(
            slide,
            x + Inches(0.25),
            Inches(5.05),
            Inches(3.0),
            Inches(0.35),
            [{"text": label, "size": 9, "bold": True, "color": TEAL_SOFT}],
        )
        x += Inches(3.75)

    textbox(
        slide,
        Inches(0.7),
        Inches(5.75),
        Inches(10),
        Inches(0.25),
        [{"text": "† Third-party estimates; see p.9", "size": 8, "color": SLATE}],
    )
    textbox(
        slide,
        Inches(0.7),
        SH - Inches(0.85),
        Inches(12),
        Inches(0.35),
        [{
            "text": "NEW YORK  ·  SHANGHAI  ·  LONDON  ·  SINGAPORE  ·  SEOUL  ·  TORONTO  ·  DELHI",
            "size": 10,
            "bold": True,
            "color": RGBColor(0x9A, 0xB0, 0xC0),
            "align": PP_ALIGN.LEFT,
        }],
    )


def slide_disclaimer(prs):
    slide = blank_slide(prs)
    section_header(slide, "Important Notice", "Disclaimer")
    body = (
        "This presentation has been prepared by 51 Careers Inc. and its affiliates (together, “51Careers.AI” or the “Company”) "
        "solely for informational purposes in connection with a proposed private financing. It is confidential and intended "
        "exclusively for the named recipient; it may not be reproduced, distributed, or disclosed, in whole or in part, without "
        "the Company’s prior written consent.\n\n"
        "This document does not constitute an offer to sell, or a solicitation of an offer to buy, any securities in any "
        "jurisdiction. Any such offer will be made only pursuant to definitive documentation and only to qualified investors in "
        "compliance with applicable securities laws.\n\n"
        "Certain statements herein are forward-looking and reflect current expectations of management. Such statements involve "
        "known and unknown risks and uncertainties, and actual results may differ materially. Certain market, operating, and "
        "financial figures are estimates — including third-party estimates — presented for illustration only and remain subject "
        "to confirmation in due diligence. No representation or warranty, express or implied, is made as to the accuracy or "
        "completeness of the information contained herein, and nothing herein constitutes legal, tax, financial, or investment "
        "advice. Recipients should conduct their own independent investigation and analysis."
    )
    panel = round_rect(slide, Inches(0.55), Inches(1.35), Inches(12.2), Inches(4.8), fill=WHITE, line=LINE)
    textbox(
        slide,
        Inches(0.85),
        Inches(1.6),
        Inches(11.6),
        Inches(4.2),
        [{"text": body, "size": 12, "color": BODY}],
    )
    textbox(
        slide,
        Inches(0.55),
        Inches(6.35),
        Inches(8),
        Inches(0.3),
        [{"text": "Prepared by 51Careers.AI — July 2026", "size": 10, "color": MUTED}],
    )
    footer(slide, 2)


def slide_highlights(prs):
    slide = blank_slide(prs)
    section_header(slide, "Executive Summary", "Investment Highlights")
    items = [
        ("01", "Large & Expanding Market",
         "Global recruitment and career-development spend exceeds US$200B, with AI-enabled segments compounding at double-digit rates."),
        ("02", "A Decade of Operating Proof",
         "Founded in New York in 2016; an established, revenue-generating career-services franchise now scaling as an AI-first platform."),
        ("03", "AI-Native Product Suite",
         "Personalized mock interviews, proprietary question banks, online assessments, mentor coaching, and intelligent job matching."),
        ("04", "Global Footprint",
         "Teams and delivery capability across the U.S., China, U.K., Singapore, South Korea, Canada, and India."),
        ("05", "Strategic AI Partnership",
         "Deep collaboration with Alibaba Cloud AI — Tongyi LLM, NLP, semantic search — powering advisor, matching, and screening systems."),
        ("06", "Compelling Entry Point",
         "US$5.0M primary round at a US$50.0M valuation to fund the product roadmap and global go-to-market."),
    ]
    positions = [
        (0.55, 1.35), (4.55, 1.35), (8.55, 1.35),
        (0.55, 4.05), (4.55, 4.05), (8.55, 4.05),
    ]
    for (num, title, desc), (x, y) in zip(items, positions):
        round_rect(slide, Inches(x), Inches(y), Inches(3.75), Inches(2.4), fill=WHITE, line=LINE)
        accent_bar(slide, Inches(x), Inches(y), w=Inches(0.08), h=Inches(2.4))
        textbox(slide, Inches(x + 0.25), Inches(y + 0.2), Inches(3.2), Inches(0.4),
                [{"text": num, "size": 18, "bold": True, "color": TEAL, "font": FONT_DISPLAY}])
        textbox(slide, Inches(x + 0.25), Inches(y + 0.7), Inches(3.2), Inches(0.45),
                [{"text": title, "size": 14, "bold": True, "color": NAVY}])
        textbox(slide, Inches(x + 0.25), Inches(y + 1.2), Inches(3.2), Inches(1.0),
                [{"text": desc, "size": 11, "color": BODY}])
    footer(slide, 3)


def slide_problem(prs):
    slide = blank_slide(prs)
    section_header(slide, "Market Context", "Hiring Is Broken on Both Sides of the Table")

    # Left narrative cards
    narratives = [
        ("For candidates",
         "Hiring processes are opaque and intensely competitive. Interview preparation remains fragmented across scattered resources, and quality human coaching is expensive and does not scale."),
        ("For international talent",
         "Language, cultural, and visa complexity compound an already difficult search — a large, motivated population is structurally underserved by generic tools."),
        ("For employers",
         "Application volumes overwhelm recruiting teams while offering limited signal. Screening is slow, costly, and inconsistent, and strong candidates are lost in the noise."),
    ]
    y = 1.35
    for title, body in narratives:
        round_rect(slide, Inches(0.55), Inches(y), Inches(6.3), Inches(1.5), fill=WHITE, line=LINE)
        accent_bar(slide, Inches(0.55), Inches(y), h=Inches(1.5))
        textbox(slide, Inches(0.85), Inches(y + 0.18), Inches(5.7), Inches(0.35),
                [{"text": title, "size": 14, "bold": True, "color": NAVY}])
        textbox(slide, Inches(0.85), Inches(y + 0.55), Inches(5.7), Inches(0.85),
                [{"text": body, "size": 11, "color": BODY}])
        y += 1.65

    # Right KPI metrics (separate scales — not mixed on one axis)
    round_rect(slide, Inches(7.1), Inches(1.35), Inches(5.65), Inches(5.0), fill=WHITE, line=LINE)
    textbox(slide, Inches(7.35), Inches(1.5), Inches(5.2), Inches(0.35),
            [{"text": "Structural friction indicators", "size": 13, "bold": True, "color": NAVY}])
    metrics = [
        ("250+", "Applications per corporate job opening", 0.72, TEAL),
        ("1.1M+", "International students in the U.S. alone", 0.95, BLUE),
        ("75%", "Of résumés screened out before human review", 0.75, BRONZE),
    ]
    y = 2.0
    for value, label, frac, color in metrics:
        textbox(slide, Inches(7.45), Inches(y), Inches(5.1), Inches(0.45),
                [{"text": value, "size": 28, "bold": True, "color": NAVY, "font": FONT_DISPLAY}])
        textbox(slide, Inches(7.45), Inches(y + 0.48), Inches(5.1), Inches(0.3),
                [{"text": label.upper(), "size": 9, "bold": True, "color": MUTED}])
        # proportional intensity bar
        rect(slide, Inches(7.45), Inches(y + 0.85), Inches(5.0), Inches(0.12), fill=SOFT)
        rect(slide, Inches(7.45), Inches(y + 0.85), Inches(5.0 * frac), Inches(0.12), fill=color)
        y += 1.25
    textbox(slide, Inches(7.35), Inches(5.85), Inches(5.2), Inches(0.35),
            [{"text": "Sources: Open Doors 2024; Glassdoor; industry estimates.", "size": 8, "color": MUTED}])
    footer(slide, 4)


def slide_market(prs):
    slide = blank_slide(prs)
    section_header(slide, "Market Opportunity", "A US$200B+ Global Talent Market in Structural Shift")

    # TAM / SAM / SOM as nested visual bars + values
    markets = [
        ("TOTAL ADDRESSABLE MARKET", "US$200B+", "Global recruitment, talent-acquisition, and career-development spend.", TEAL, 1.0),
        ("SERVICEABLE MARKET", "US$30B", "AI-enabled hiring technology and digital career preparation.", BLUE, 0.55),
        ("INITIAL TARGET SEGMENT", "US$1.5B", "Cross-border early-career talent across current Company markets.", BRONZE, 0.28),
    ]
    y = 1.35
    for label, value, desc, color, width_frac in markets:
        round_rect(slide, Inches(0.55), Inches(y), Inches(6.4), Inches(1.35), fill=WHITE, line=LINE)
        bar_w = Inches(5.9 * width_frac)
        rect(slide, Inches(0.75), Inches(y + 1.0), bar_w, Inches(0.16), fill=color)
        textbox(slide, Inches(0.75), Inches(y + 0.12), Inches(5.9), Inches(0.25),
                [{"text": label, "size": 9, "bold": True, "color": MUTED}])
        textbox(slide, Inches(0.75), Inches(y + 0.35), Inches(5.9), Inches(0.4),
                [{"text": value, "size": 26, "bold": True, "color": NAVY, "font": FONT_DISPLAY}])
        textbox(slide, Inches(0.75), Inches(y + 0.72), Inches(5.9), Inches(0.28),
                [{"text": desc, "size": 10, "color": BODY}])
        y += 1.5

    # Growth chart
    round_rect(slide, Inches(7.15), Inches(1.35), Inches(5.6), Inches(4.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(7.35), Inches(1.5), Inches(5.2), Inches(0.35),
            [{"text": "AI hiring & career-tech spend (illustrative, US$B)", "size": 12, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(7.25), Inches(1.9), Inches(5.35), Inches(3.7),
        ["2024", "2025E", "2026E", "2027E", "2028E"],
        [("AI hiring tech", (18, 22, 27, 33, 40))],
        legend=False,
    )
    textbox(slide, Inches(0.55), Inches(6.35), Inches(12), Inches(0.3),
            [{"text": "Management estimates informed by public industry research; presented for illustration.", "size": 8, "color": MUTED}])
    footer(slide, 5)


def slide_platform(prs):
    slide = blank_slide(prs)
    section_header(slide, "The Platform", "One Platform Spanning Preparation, Placement & Hiring")

    # Two columns
    round_rect(slide, Inches(0.55), Inches(1.35), Inches(6.0), Inches(4.55), fill=WHITE, line=LINE)
    rect(slide, Inches(0.55), Inches(1.35), Inches(6.0), Inches(0.5), fill=TEAL)
    textbox(slide, Inches(0.75), Inches(1.45), Inches(5.6), Inches(0.35),
            [{"text": "FOR CANDIDATES", "size": 12, "bold": True, "color": WHITE}])
    cand = [
        ("AI Mock Interviews", "Role- and company-specific practice with personalized, real-time feedback."),
        ("Question Banks & Assessments", "Proprietary libraries of real interview questions and timed online assessments."),
        ("Mentor Coaching", "One-on-one guidance from practitioners at leading global employers."),
        ("Referrals & Placement", "Internship and full-time referral programs, including outcome-based offerings."),
    ]
    y = 2.05
    for t, d in cand:
        textbox(slide, Inches(0.85), Inches(y), Inches(5.4), Inches(0.3),
                [{"text": t, "size": 13, "bold": True, "color": NAVY}])
        textbox(slide, Inches(0.85), Inches(y + 0.28), Inches(5.4), Inches(0.4),
                [{"text": d, "size": 11, "color": BODY}])
        y += 0.85

    round_rect(slide, Inches(6.8), Inches(1.35), Inches(6.0), Inches(4.55), fill=WHITE, line=LINE)
    rect(slide, Inches(6.8), Inches(1.35), Inches(6.0), Inches(0.5), fill=BLUE)
    textbox(slide, Inches(7.0), Inches(1.45), Inches(5.6), Inches(0.35),
            [{"text": "FOR EMPLOYERS", "size": 12, "bold": True, "color": WHITE}])
    emp = [
        ("AI Candidate Screening", "Semantic evaluation of fit, capability, and readiness at scale."),
        ("Job Posting Optimization", "AI-assisted role definition and sourcing precision."),
        ("Talent Evaluation", "Structured assessment and interview analytics for confident decisions."),
    ]
    y = 2.15
    for t, d in emp:
        textbox(slide, Inches(7.1), Inches(y), Inches(5.4), Inches(0.3),
                [{"text": t, "size": 13, "bold": True, "color": NAVY}])
        textbox(slide, Inches(7.1), Inches(y + 0.28), Inches(5.4), Inches(0.45),
                [{"text": d, "size": 11, "color": BODY}])
        y += 1.0

    # flywheel note
    round_rect(slide, Inches(0.55), Inches(6.05), Inches(12.25), Inches(0.55), fill=TEAL_SOFT, line=None)
    textbox(slide, Inches(0.75), Inches(6.15), Inches(11.8), Inches(0.4),
            [{"text": "The flywheel: candidate outcomes generate proprietary data that sharpens matching and screening for employers — a compounding advantage on both sides of the market.",
              "size": 11, "bold": True, "color": TEAL_DK}])
    footer(slide, 6)


def slide_tech(prs):
    slide = blank_slide(prs)
    section_header(slide, "Technology", "Enterprise-Grade AI, Built with Alibaba Cloud")

    round_rect(slide, Inches(0.55), Inches(1.35), Inches(7.4), Inches(2.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.8), Inches(1.5), Inches(7.0), Inches(0.35),
            [{"text": "Alibaba Cloud AI partnership — July 2025", "size": 14, "bold": True, "color": TEAL}])
    textbox(slide, Inches(0.8), Inches(1.95), Inches(6.9), Inches(1.1),
            [{"text": "A deep collaboration integrating frontier AI capability across the platform — the Tongyi large language model, natural-language processing, LLM inference, intelligent recommendation, and semantic search — on enterprise-grade compute and security infrastructure.",
              "size": 12, "color": BODY}])
    bullets = [
        "AI Career Advisor delivering personalized career-path recommendations",
        "Intelligent matching and screening for precise candidate–role fit",
        "Enterprise data-security framework protecting user privacy",
    ]
    y = 3.15
    for b in bullets:
        textbox(slide, Inches(0.8), Inches(y), Inches(6.9), Inches(0.28),
                [{"text": "▸  " + b, "size": 11, "color": INK}])
        y += 0.28

    round_rect(slide, Inches(8.15), Inches(1.35), Inches(4.65), Inches(2.55), fill=NAVY)
    textbox(slide, Inches(8.4), Inches(1.55), Inches(4.2), Inches(0.35),
            [{"text": "Proprietary data moat", "size": 14, "bold": True, "color": WHITE, "font": FONT_DISPLAY}])
    textbox(slide, Inches(8.4), Inches(2.1), Inches(4.2), Inches(1.5),
            [{"text": "A decade of accumulated interview questions, outcome-labeled placement data, and mentor feedback loops — assets that generic models and job boards cannot readily replicate.",
              "size": 12, "color": RGBColor(0xC5, 0xD5, 0xE0)}])

    # Architecture stack visual
    layers = [
        ("APPLICATIONS", "Mock Interviews · Career Advisor · Matching · Screening", TEAL),
        ("INTELLIGENCE", "Tongyi LLM · NLP · Semantic Search · Recommendation", BLUE),
        ("DATA & INFRASTRUCTURE", "Proprietary question corpus · Outcome data · Alibaba Cloud", BRONZE),
    ]
    y = 4.15
    textbox(slide, Inches(0.55), Inches(4.05), Inches(6), Inches(0.3),
            [{"text": "Platform architecture (simplified)", "size": 12, "bold": True, "color": NAVY}])
    for label, detail, color in layers:
        round_rect(slide, Inches(0.55), Inches(y), Inches(12.25), Inches(0.7), fill=WHITE, line=LINE)
        rect(slide, Inches(0.55), Inches(y), Inches(0.12), Inches(0.7), fill=color)
        textbox(slide, Inches(0.9), Inches(y + 0.08), Inches(3.5), Inches(0.28),
                [{"text": label, "size": 11, "bold": True, "color": color}])
        textbox(slide, Inches(4.5), Inches(y + 0.18), Inches(8.0), Inches(0.35),
                [{"text": detail, "size": 12, "color": BODY}])
        y += 0.8
    footer(slide, 7)


def slide_model(prs):
    slide = blank_slide(prs)
    section_header(slide, "Business Model", "Diversified, Outcome-Aligned Revenue")

    pillars = [
        ("01 · CONSUMER PLATFORM", "Subscriptions & Prep",
         "Self-serve AI interview practice, question banks, and assessments sold as subscriptions and credit packs.",
         "Recurring · High margin · Scales globally", TEAL),
        ("02 · PREMIUM SERVICES", "Coaching & Placement",
         "Mentor-led programs and outcome-based placement offerings — the Company’s proven, high-ARPU franchise since 2016.",
         "High ARPU · Proven demand · Cash generative", BLUE),
        ("03 · ENTERPRISE & INSTITUTIONS", "Employer SaaS",
         "AI screening, talent evaluation, and posting optimization for employers; partnerships with universities and institutions.",
         "B2B · Land-and-expand · Data flywheel", BRONZE),
    ]
    x = 0.55
    for eyebrow, title, desc, tag, color in pillars:
        round_rect(slide, Inches(x), Inches(1.35), Inches(3.95), Inches(2.9), fill=WHITE, line=LINE)
        rect(slide, Inches(x), Inches(1.35), Inches(3.95), Inches(0.12), fill=color)
        textbox(slide, Inches(x + 0.2), Inches(1.6), Inches(3.55), Inches(0.3),
                [{"text": eyebrow, "size": 10, "bold": True, "color": color}])
        textbox(slide, Inches(x + 0.2), Inches(1.95), Inches(3.55), Inches(0.4),
                [{"text": title, "size": 16, "bold": True, "color": NAVY, "font": FONT_DISPLAY}])
        textbox(slide, Inches(x + 0.2), Inches(2.45), Inches(3.55), Inches(1.1),
                [{"text": desc, "size": 11, "color": BODY}])
        textbox(slide, Inches(x + 0.2), Inches(3.7), Inches(3.55), Inches(0.35),
                [{"text": tag, "size": 10, "bold": True, "color": MUTED}])
        x += 4.15

    # Revenue mix evolution chart
    round_rect(slide, Inches(0.55), Inches(4.45), Inches(12.25), Inches(2.15), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(4.55), Inches(5), Inches(0.3),
            [{"text": "Illustrative revenue mix evolution (%)", "size": 12, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_STACKED_100,
        Inches(0.7), Inches(4.8), Inches(8.2), Inches(1.7),
        ["Today", "2027E", "2029E"],
        [
            ("Premium services", (70, 45, 25)),
            ("Consumer platform", (25, 35, 40)),
            ("Employer SaaS", (5, 20, 35)),
        ],
        legend=True,
    )
    textbox(slide, Inches(9.1), Inches(5.15), Inches(3.4), Inches(1.2),
            [{"text": "Model evolution", "size": 12, "bold": True, "color": NAVY, "space_after": 6},
             {"text": "Services-led today → blended by 2027 → software-led by 2029, expanding gross margin as the platform scales.",
              "size": 11, "color": BODY}])
    footer(slide, 8)


def slide_traction(prs):
    slide = blank_slide(prs)
    section_header(slide, "Traction", "A Real Business, Not a Concept")

    kpis = [
        ("10 yrs", "OPERATING HISTORY SINCE 2016"),
        ("US$10M+", "ANNUAL REVENUE (2025, EST.)†"),
        ("≈50", "PROFESSIONALS WORLDWIDE†"),
        ("8+", "COUNTRIES OF OPERATION"),
    ]
    x = 0.55
    for val, label in kpis:
        round_rect(slide, Inches(x), Inches(1.3), Inches(2.95), Inches(1.25), fill=WHITE, line=LINE)
        textbox(slide, Inches(x + 0.15), Inches(1.4), Inches(2.65), Inches(0.55),
                [{"text": val, "size": 26, "bold": True, "color": NAVY, "font": FONT_DISPLAY}])
        textbox(slide, Inches(x + 0.15), Inches(2.05), Inches(2.65), Inches(0.35),
                [{"text": label, "size": 9, "bold": True, "color": MUTED}])
        x += 3.15

    # Timeline
    round_rect(slide, Inches(0.55), Inches(2.75), Inches(7.5), Inches(3.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(2.9), Inches(7), Inches(0.3),
            [{"text": "Selected milestones", "size": 13, "bold": True, "color": NAVY}])
    # timeline line
    rect(slide, Inches(1.15), Inches(3.45), Inches(0.04), Inches(2.55), fill=TEAL)
    milestones = [
        ("2016", "Founded in New York as a premium career-consulting firm for international students"),
        ("2016–24", "Scaled the coaching, referral, and placement franchise across the U.S. and Greater China"),
        ("Jul 2025", "Strategic AI partnership with Alibaba Cloud (Tongyi LLM, NLP, semantic search)"),
        ("2025", "Launched 51Careers.AI — the Company’s AI-powered global career platform"),
        ("2026", "Global platform expansion underway across eight-plus countries"),
    ]
    y = 3.35
    for date, desc in milestones:
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(1.05), Inches(y + 0.08), Inches(0.24), Inches(0.24))
        fill_solid(oval, TEAL)
        textbox(slide, Inches(1.5), Inches(y), Inches(1.2), Inches(0.3),
                [{"text": date, "size": 11, "bold": True, "color": TEAL}])
        textbox(slide, Inches(2.75), Inches(y), Inches(5.0), Inches(0.5),
                [{"text": desc, "size": 11, "color": BODY}])
        y += 0.55

    # Why it matters + mini chart of scale
    round_rect(slide, Inches(8.25), Inches(2.75), Inches(4.55), Inches(3.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(8.45), Inches(2.9), Inches(4.2), Inches(0.3),
            [{"text": "Why it matters", "size": 13, "bold": True, "color": NAVY}])
    reasons = [
        "Durable, decade-long demand with proven willingness to pay",
        "Cross-border delivery capability that is difficult to replicate",
        "AI now layers operating leverage onto an established service core",
        "Existing revenue de-risks the platform transition",
    ]
    y = 3.35
    for r in reasons:
        textbox(slide, Inches(8.45), Inches(y), Inches(4.15), Inches(0.55),
                [{"text": "▸  " + r, "size": 11, "color": BODY}])
        y += 0.55

    textbox(slide, Inches(0.55), Inches(6.4), Inches(12), Inches(0.3),
            [{"text": "† Third-party estimates (ZoomInfo; LeadIQ, 2025–26) ranging US$9.7–15M revenue and ≈30–50 staff; to be confirmed in due diligence.",
              "size": 8, "color": MUTED}])
    footer(slide, 9)


def slide_competition(prs):
    slide = blank_slide(prs)
    section_header(slide, "Competitive Positioning", "Differentiated at the Intersection of AI and Outcomes")

    # Capability score chart
    round_rect(slide, Inches(0.55), Inches(1.3), Inches(5.7), Inches(3.6), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(1.4), Inches(5.3), Inches(0.3),
            [{"text": "Capability coverage score (illustrative)", "size": 12, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.BAR_CLUSTERED,
        Inches(0.65), Inches(1.75), Inches(5.45), Inches(2.95),
        ["51Careers.AI", "Job Boards", "Prep Apps", "Coaching Firms", "Staffing Agencies"],
        [("Coverage (0–6)", (6, 1, 2, 3, 2))],
        legend=False,
    )

    # Matrix table
    headers = ["Capability", "51Careers", "Job Boards", "Prep Apps", "Coaching", "Staffing"]
    rows = [
        ["AI interview practice", "●", "—", "○", "○", "—"],
        ["Real-question banks", "●", "—", "○", "○", "—"],
        ["Human mentor network", "●", "—", "—", "●", "○"],
        ["Outcome placement", "●", "—", "—", "○", "●"],
        ["Employer AI screening", "●", "○", "—", "—", "○"],
        ["Cross-border focus", "●", "○", "—", "○", "○"],
    ]
    table_shape = slide.shapes.add_table(7, 6, Inches(6.45), Inches(1.3), Inches(6.35), Inches(3.6))
    table = table_shape.table
    for c, h in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = h
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
            for run in p.runs:
                run.font.name = FONT_BODY
                run.font.size = Pt(9)
                run.font.bold = True
                run.font.color.rgb = WHITE
        cell.fill.solid()
        cell.fill.fore_color.rgb = NAVY
    for r, row in enumerate(rows, 1):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.CENTER if c else PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT_BODY
                    run.font.size = Pt(10)
                    run.font.bold = (c == 1 and val == "●")
                    run.font.color.rgb = TEAL if val == "●" else (BODY if c == 0 else MUTED)
            cell.fill.solid()
            cell.fill.fore_color.rgb = WHITE if r % 2 else SOFT

    textbox(slide, Inches(0.55), Inches(5.1), Inches(12.2), Inches(0.35),
            [{"text": "●  Full capability      ○  Partial      —  Not offered", "size": 10, "color": MUTED}])
    round_rect(slide, Inches(0.55), Inches(5.5), Inches(12.25), Inches(0.85), fill=TEAL_SOFT)
    textbox(slide, Inches(0.75), Inches(5.65), Inches(11.85), Inches(0.55),
            [{"text": "51Careers.AI is the only platform combining AI-native preparation, human mentorship, and outcome-based placement across borders.",
              "size": 13, "bold": True, "color": TEAL_DK}])
    textbox(slide, Inches(0.55), Inches(6.4), Inches(12), Inches(0.25),
            [{"text": "Illustrative category comparison based on management assessment of publicly available offerings.", "size": 8, "color": MUTED}])
    footer(slide, 10)


def slide_growth(prs):
    slide = blank_slide(prs)
    section_header(slide, "Growth Strategy", "A Phased Path to Global Scale")

    phases = [
        ("2026", "Strengthen the core", TEAL, [
            "Scale the AI platform across the U.S. and Greater China",
            "Migrate the coaching client base onto the platform",
            "Expand question-bank and assessment coverage by role and employer",
        ]),
        ("2027", "Extend the footprint", BLUE, [
            "Full go-to-market in the U.K., Singapore, South Korea, Canada, and India",
            "University and institutional partnership programs",
            "Convert employer pilots into paid SaaS relationships",
        ]),
        ("2028+", "Compound the platform", BRONZE, [
            "Employer SaaS at scale across core geographies",
            "Additional languages, verticals, and seniority segments",
            "Selective M&A of regional career-services assets",
        ]),
    ]
    x = 0.55
    for year, title, color, bullets in phases:
        round_rect(slide, Inches(x), Inches(1.35), Inches(3.95), Inches(3.7), fill=WHITE, line=LINE)
        rect(slide, Inches(x), Inches(1.35), Inches(3.95), Inches(0.9), fill=color)
        textbox(slide, Inches(x + 0.25), Inches(1.45), Inches(3.45), Inches(0.4),
                [{"text": year, "size": 26, "bold": True, "color": WHITE, "font": FONT_DISPLAY}])
        textbox(slide, Inches(x + 0.25), Inches(1.9), Inches(3.45), Inches(0.3),
                [{"text": title, "size": 13, "bold": True, "color": WHITE}])
        y = 2.5
        for b in bullets:
            textbox(slide, Inches(x + 0.25), Inches(y), Inches(3.45), Inches(0.7),
                    [{"text": "▸  " + b, "size": 12, "color": BODY}])
            y += 0.75
        x += 4.15

    # Progress visual chart — market expansion intensity
    round_rect(slide, Inches(0.55), Inches(5.2), Inches(12.25), Inches(1.2), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(5.3), Inches(4), Inches(0.28),
            [{"text": "Geographic expansion intensity (illustrative)", "size": 11, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(5.0), Inches(5.15), Inches(7.6), Inches(1.25),
        ["2025", "2026", "2027", "2028", "2029"],
        [("Active markets", (3, 4, 8, 8, 10))],
        legend=False,
    )
    textbox(slide, Inches(0.55), Inches(6.45), Inches(12), Inches(0.25),
            [{"text": "Sequencing follows the existing footprint — new capital extends proven demand in markets the Company already serves, rather than funding market discovery.",
              "size": 11, "color": MUTED}])
    footer(slide, 11)


def slide_financials(prs):
    slide = blank_slide(prs)
    section_header(slide, "Financial Outlook", "Illustrative Five-Year Trajectory")

    round_rect(slide, Inches(0.55), Inches(1.3), Inches(7.9), Inches(3.15), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(1.4), Inches(5), Inches(0.28),
            [{"text": "Revenue, US$ millions", "size": 12, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.COLUMN_STACKED,
        Inches(0.65), Inches(1.65), Inches(7.65), Inches(2.65),
        ["2025E", "2026E", "2027E", "2028E", "2029E"],
        [
            ("Estimated base", (10, None, None, None, None)),
            ("Projected", (None, 14.5, 22.5, 34, 52)),
        ],
        legend=True,
    )

    # Gross margin trajectory
    round_rect(slide, Inches(0.55), Inches(4.6), Inches(7.9), Inches(1.7), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.75), Inches(4.68), Inches(5), Inches(0.25),
            [{"text": "Illustrative gross margin (%)", "size": 11, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(0.65), Inches(4.85), Inches(7.65), Inches(1.4),
        ["2025E", "2026E", "2027E", "2028E", "2029E"],
        [("Gross margin", (42, 48, 55, 60, 65))],
        legend=False,
    )

    # KPI + drivers
    round_rect(slide, Inches(8.65), Inches(1.3), Inches(4.15), Inches(1.55), fill=NAVY)
    textbox(slide, Inches(8.85), Inches(1.45), Inches(3.75), Inches(0.65),
            [{"text": "≈50%", "size": 36, "bold": True, "color": WHITE, "font": FONT_DISPLAY}])
    textbox(slide, Inches(8.85), Inches(2.2), Inches(3.75), Inches(0.4),
            [{"text": "REVENUE CAGR, 2025E–2029E", "size": 10, "bold": True, "color": TEAL_SOFT}])

    round_rect(slide, Inches(8.65), Inches(3.05), Inches(4.15), Inches(3.25), fill=WHITE, line=LINE)
    textbox(slide, Inches(8.85), Inches(3.2), Inches(3.75), Inches(0.3),
            [{"text": "Key drivers", "size": 13, "bold": True, "color": NAVY}])
    drivers = [
        "Self-serve platform subscriptions compounding across markets",
        "Mix shift from services to software lifting gross margin toward 65%+",
        "Employer SaaS ramping from 2027",
        "AI-assisted delivery driving operating leverage",
    ]
    y = 3.6
    for d in drivers:
        textbox(slide, Inches(8.85), Inches(y), Inches(3.75), Inches(0.6),
                [{"text": "▸  " + d, "size": 11, "color": BODY}])
        y += 0.6

    textbox(slide, Inches(0.55), Inches(6.4), Inches(12.2), Inches(0.3),
            [{"text": "Illustrative management projections, not a forecast; 2025 base reflects third-party revenue estimates. Actual results will differ.",
              "size": 8, "color": MUTED}])
    footer(slide, 12)


def slide_offering(prs):
    slide = blank_slide(prs)
    section_header(slide, "The Offering", "US$5.0 Million to Accelerate Global Leadership")

    # Terms panel
    round_rect(slide, Inches(0.55), Inches(1.3), Inches(5.9), Inches(4.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(0.8), Inches(1.5), Inches(5.4), Inches(0.3),
            [{"text": "INDICATIVE TERMS", "size": 11, "bold": True, "color": TEAL}])
    terms = [
        ("Round size", "US$5.0 million (primary)"),
        ("Valuation", "US$50.0 million (pre-money)"),
        ("Structure", "Preferred equity"),
        ("Use of funds", "Product & AI, expansion, team"),
        ("Target close", "2H 2026"),
    ]
    y = 2.0
    for label, value in terms:
        rect(slide, Inches(0.8), Inches(y + 0.55), Inches(5.4), Pt(1), fill=LINE)
        textbox(slide, Inches(0.8), Inches(y), Inches(2.4), Inches(0.4),
                [{"text": label, "size": 12, "color": MUTED}])
        textbox(slide, Inches(3.2), Inches(y), Inches(3.0), Inches(0.4),
                [{"text": value, "size": 13, "bold": True, "color": NAVY}])
        y += 0.65
    textbox(slide, Inches(0.8), Inches(5.4), Inches(5.4), Inches(0.3),
            [{"text": "Indicative only; subject to definitive documentation.", "size": 9, "italic": True, "color": MUTED}])

    # Doughnut
    round_rect(slide, Inches(6.7), Inches(1.3), Inches(6.1), Inches(4.55), fill=WHITE, line=LINE)
    textbox(slide, Inches(6.95), Inches(1.5), Inches(5.6), Inches(0.3),
            [{"text": "Use of proceeds", "size": 13, "bold": True, "color": NAVY}])
    add_chart(
        slide,
        XL_CHART_TYPE.DOUGHNUT,
        Inches(6.9), Inches(1.85), Inches(5.7), Inches(3.5),
        [
            "Product & AI development (35%)",
            "International expansion (30%)",
            "Talent & team (20%)",
            "Brand & growth marketing (10%)",
            "Working capital (5%)",
        ],
        [("Use of Proceeds", (35, 30, 20, 10, 5))],
        legend=True,
        point_colors=CHART_COLORS,
    )

    round_rect(slide, Inches(0.55), Inches(6.0), Inches(12.25), Inches(0.55), fill=TEAL_SOFT)
    textbox(slide, Inches(0.75), Inches(6.1), Inches(11.85), Inches(0.4),
            [{"text": "What this capital achieves — next-generation platform release · go-to-market in three additional countries · employer SaaS general availability · scaled global mentor network.",
              "size": 12, "bold": True, "color": TEAL_DK}])
    footer(slide, 13)


def slide_leadership(prs):
    """Executive leadership — six condensed investor-ready profiles on one slide."""
    slide = blank_slide(prs)
    section_header(slide, "Leadership", "Executive Leadership Team")

    leaders = [
        (
            "Rocky Chen",
            "FOUNDER & CEO",
            [
                "Founder of 51 Careers and 51Careers.AI; investor & Managing Director at publicly listed Helio",
                "Decade of international business, capital markets, and global resource integration across AI, education, and related industries",
            ],
        ),
        (
            "Stephanie Li",
            "CO-FOUNDER & CFO",
            [
                "Co-Founder since 2016; U.S. CPA — M.S. Accounting (Pace); B.A. Accounting (Dongbei University of Finance and Economics)",
                "Oversees financial strategy, corporate governance, and sustainable growth; prior senior finance leadership in U.S. industry",
            ],
        ),
        (
            "Gavin Ding",
            "CO-FOUNDER, CTO & COO",
            [
                "Serial entrepreneur (10+ yrs) across SaaS, AI, and digital business; B.S. Computer Science, ECUST",
                "Prior CTO roles (Youpindao, Roubeibei) and multiple co-founded ventures; leads technology, operations, and AI strategy",
            ],
        ),
        (
            "Robin Zhu",
            "HEAD OF PRODUCT",
            [
                "Leads product strategy, architecture, and development of the AI-powered platform",
                "Former CTO & Director of Product, Ci Finance; senior product/tech roles at CPIC, Allinpay, Noah Holdings, and Tianrang Intelligence",
            ],
        ),
        (
            "Chris Lin",
            "NORTH AMERICAN PARTNER",
            [
                "Full-stack engineer; previously at Amazon Web Services (AWS) on mission-critical systems",
                "B.S. Mechanical Engineering (Northwestern), M.S. Robotics; extensive hiring-panel experience and career mentorship",
            ],
        ),
        (
            "Jon Serbin",
            "SENIOR ADVISOR",
            [
                "Harvard & MIT; former senior executive at Morgan Stanley; founder of Cedar (Wall Street investment banking)",
                "40+ years advising tech M&A, capital raising, and listings; advises 51Careers.AI on growth, capital markets, and expansion",
            ],
        ),
    ]

    # 2 × 3 grid
    card_w, card_h = Inches(4.0), Inches(2.35)
    gap_x, gap_y = Inches(0.2), Inches(0.18)
    origin_x, origin_y = Inches(0.55), Inches(1.28)
    accents = [TEAL, BLUE, BRONZE, TEAL, BLUE, BRONZE]

    for i, (name, title, bullets) in enumerate(leaders):
        col, row = i % 3, i // 3
        x = origin_x + col * (card_w + gap_x)
        y = origin_y + row * (card_h + gap_y)
        round_rect(slide, x, y, card_w, card_h, fill=WHITE, line=LINE)
        accent_bar(slide, x, y, w=Inches(0.08), h=card_h, color=accents[i])
        textbox(
            slide,
            x + Inches(0.22),
            y + Inches(0.14),
            card_w - Inches(0.35),
            Inches(0.35),
            [{"text": name, "size": 15, "bold": True, "color": NAVY, "font": FONT_DISPLAY}],
        )
        textbox(
            slide,
            x + Inches(0.22),
            y + Inches(0.48),
            card_w - Inches(0.35),
            Inches(0.28),
            [{"text": title, "size": 10, "bold": True, "color": TEAL}],
        )
        by = y + Inches(0.85)
        for bullet in bullets:
            textbox(
                slide,
                x + Inches(0.22),
                by,
                card_w - Inches(0.38),
                Inches(0.65),
                [{"text": "▸  " + bullet, "size": 10, "color": BODY}],
            )
            by += Inches(0.68)

    textbox(
        slide,
        Inches(0.55),
        Inches(6.3),
        Inches(12.2),
        Inches(0.3),
        [{
            "text": "Supported by ≈50 professionals across New York, Shanghai, London, Singapore, Seoul, Toronto, and Delhi. Full biographies available in the data room.",
            "size": 9,
            "color": MUTED,
        }],
    )
    footer(slide, 14)


def slide_close(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    fill_solid(slide.background, NAVY)
    rect(slide, Inches(0), Inches(0), SW, Inches(0.12), fill=TEAL)
    rect(slide, Inches(0), Inches(0), Inches(0.18), SH, fill=TEAL)

    textbox(
        slide,
        Inches(1.0),
        Inches(1.8),
        Inches(11.3),
        Inches(1.5),
        [{
            "text": "“Our mission is to make great jobs accessible to everyone.”",
            "size": 28,
            "bold": True,
            "color": WHITE,
            "font": FONT_DISPLAY,
            "align": PP_ALIGN.CENTER,
        }],
        valign=MSO_ANCHOR.MIDDLE,
    )
    textbox(
        slide,
        Inches(1.0),
        Inches(3.5),
        Inches(11.3),
        Inches(0.5),
        [{"text": "51CAREERS.AI", "size": 22, "bold": True, "color": TEAL, "font": FONT_DISPLAY, "align": PP_ALIGN.CENTER}],
    )
    textbox(
        slide,
        Inches(1.0),
        Inches(4.2),
        Inches(11.3),
        Inches(0.35),
        [{"text": "Rocky Chen  ·  Founder & Chief Executive Officer", "size": 13, "color": RGBColor(0xC5, 0xD5, 0xE0), "align": PP_ALIGN.CENTER}],
    )
    textbox(
        slide,
        Inches(1.0),
        Inches(5.0),
        Inches(11.3),
        Inches(0.35),
        [{"text": "info@51careers.com      ·      www.51careers.ai", "size": 12, "color": WHITE, "align": PP_ALIGN.CENTER}],
    )
    textbox(
        slide,
        Inches(1.0),
        Inches(5.4),
        Inches(11.3),
        Inches(0.35),
        [{"text": "48 Wall Street, 11th Floor, New York, NY 10005", "size": 12, "color": SLATE, "align": PP_ALIGN.CENTER}],
    )
    textbox(
        slide,
        Inches(1.0),
        SH - Inches(0.7),
        Inches(11.3),
        Inches(0.3),
        [{"text": "51CAREERS.AI   ·   CONFIDENTIAL   ·   15", "size": 9, "bold": True, "color": SLATE, "align": PP_ALIGN.CENTER}],
    )


def build():
    prs = Presentation()
    prs.slide_width = SW
    prs.slide_height = SH

    slide_cover(prs)
    slide_disclaimer(prs)
    slide_highlights(prs)
    slide_problem(prs)
    slide_market(prs)
    slide_platform(prs)
    slide_tech(prs)
    slide_model(prs)
    slide_traction(prs)
    slide_competition(prs)
    slide_growth(prs)
    slide_financials(prs)
    slide_offering(prs)
    slide_leadership(prs)
    slide_close(prs)

    prs.save(OUT)
    print(f"Wrote {OUT} ({OUT.stat().st_size:,} bytes) with {len(prs.slides)} slides")
    # count charts
    charts = 0
    for s in prs.slides:
        for sh in s.shapes:
            if sh.has_chart:
                charts += 1
    print(f"Native charts: {charts}")


if __name__ == "__main__":
    build()
