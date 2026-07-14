#!/usr/bin/env python3
"""Rebuild Samuel Space investor deck: light, professional, dense, no emoji icons."""

from __future__ import annotations

import io
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.dml.color import RGBColor
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import nsmap, qn
from pptx.util import Emu, Inches, Pt
from lxml import etree

# ---------------------------------------------------------------------------
# Design tokens — light institutional palette
# ---------------------------------------------------------------------------
W, H = 13.333, 7.5
MARGIN = 0.28
CONTENT_TOP = 1.18
CONTENT_BOTTOM = 7.12
CONTENT_H = CONTENT_BOTTOM - CONTENT_TOP
CONTENT_W = W - 2 * MARGIN

BG = RGBColor(0xF7, 0xF9, 0xFC)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
PANEL = RGBColor(0xEE, 0xF3, 0xF8)
PANEL_ALT = RGBColor(0xE4, 0xEC, 0xF3)
BORDER = RGBColor(0xD0, 0xDA, 0xE6)
NAVY = RGBColor(0x0B, 0x1F, 0x33)
BODY = RGBColor(0x3A, 0x4A, 0x5C)
MUTED = RGBColor(0x6B, 0x7C, 0x8F)
ACCENT = RGBColor(0x0A, 0x6B, 0x63)
ACCENT_DARK = RGBColor(0x07, 0x52, 0x4C)
ACCENT_SOFT = RGBColor(0xD7, 0xEB, 0xE8)
ACCENT_MID = RGBColor(0x1A, 0x8A, 0x80)
LINE = RGBColor(0xC5, 0xD1, 0xDE)
CHART_1 = RGBColor(0x0A, 0x6B, 0x63)
CHART_2 = RGBColor(0x2F, 0x5D, 0x8A)
CHART_3 = RGBColor(0x8A, 0x6D, 0x3B)
CHART_4 = RGBColor(0x5A, 0x6B, 0x7C)
CHART_5 = RGBColor(0xA8, 0xB5, 0xC2)

FONT_DISPLAY = "Georgia"
FONT_BODY = "Calibri"

OUT = Path("/workspace/Samuel Space Organ and Tissue Engineering - Investor Presentation (2).pptx")
ASSET_DIR = Path("/workspace/assets")


# ---------------------------------------------------------------------------
# Low-level helpers
# ---------------------------------------------------------------------------
def set_run(run, text, size, color, bold=False, font=FONT_BODY):
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font


def clear_paragraphs(tf):
    p = tf.paragraphs[0]
    p.clear()
    for extra in list(tf.paragraphs[1:]):
        p_elem = extra._p
        p_elem.getparent().remove(p_elem)
    return tf.paragraphs[0]


def add_text_box(slide, left, top, width, height, text, size, color,
                 bold=False, font=FONT_BODY, align=PP_ALIGN.LEFT,
                 anchor=MSO_ANCHOR.TOP, lines=None):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        tf._txBody.bodyPr.set("anchor", {
            MSO_ANCHOR.TOP: "t",
            MSO_ANCHOR.MIDDLE: "ctr",
            MSO_ANCHOR.BOTTOM: "b",
        }.get(anchor, "t"))
    except Exception:
        pass
    if lines is None:
        lines = [text]
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_before = Pt(0)
        p.space_after = Pt(0)
        run = p.add_run()
        if isinstance(line, tuple):
            t, sz, col, b, f = (line + (bold, font))[:5]
            set_run(run, t, sz, col, b, f)
        else:
            set_run(run, line, size, color, bold, font)
    return shape


def _strip_shadow(shape):
    """Remove any default theme shadow/effect list from a shape."""
    sp_pr = shape._element.find(qn("p:spPr"))
    if sp_pr is None:
        return
    for effect in list(sp_pr.findall(qn("a:effectLst"))):
        sp_pr.remove(effect)
    # Explicit empty effect list prevents theme shadow inheritance
    etree.SubElement(sp_pr, qn("a:effectLst"))


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    # Prefer sharp rectangles for a more institutional look (no soft card shadows)
    shape_type = MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    _strip_shadow(shape)
    return shape


def add_line(slide, left, top, width, color=LINE, thickness=1.25):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(left), Inches(top), Inches(width), Pt(thickness))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def set_slide_bg(slide, color=BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def page_number(slide, n, total=21):
    add_text_box(slide, W - MARGIN - 0.7, 7.15, 0.7, 0.25,
                 f"{n}", 9, MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def section_header(slide, eyebrow, title, n, subtitle=None):
    set_slide_bg(slide)
    # Top accent bar full width
    add_rect(slide, 0, 0, W, 0.055, ACCENT)
    add_text_box(slide, MARGIN, 0.14, CONTENT_W, 0.24, eyebrow.upper(), 10, ACCENT, bold=True, font=FONT_BODY)
    add_text_box(slide, MARGIN, 0.36, CONTENT_W - 0.2, 0.48, title, 24, NAVY, bold=True, font=FONT_DISPLAY)
    if subtitle:
        add_text_box(slide, MARGIN, 0.84, CONTENT_W, 0.26, subtitle, 11, MUTED)
    page_number(slide, n)
    # Footer rule
    add_line(slide, MARGIN, 7.18, CONTENT_W, LINE, 1)
    add_text_box(slide, MARGIN, 7.22, 5.0, 0.22,
                 "Samuel Space  |  Confidential", 8.5, MUTED)


def numbered_badge(slide, left, top, num, size=0.36, fill=ACCENT, font_size=12):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    _strip_shadow(shape)
    tf = shape.text_frame
    tf.word_wrap = False
    p = clear_paragraphs(tf)
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    set_run(run, str(num), font_size, WHITE, True)
    try:
        tf._txBody.bodyPr.set("anchor", "ctr")
    except Exception:
        pass
    return shape


def style_table(table, header=True):
    for r_idx, row in enumerate(table.rows):
        for cell in row.cells:
            # fill
            tc = cell._tc
            tcPr = tc.get_or_add_tcPr()
            # clear existing solid fills
            for child in list(tcPr):
                if "solidFill" in child.tag or " gradFill" in child.tag:
                    tcPr.remove(child)
            solid = etree.SubElement(tcPr, qn("a:solidFill"))
            srgb = etree.SubElement(solid, qn("a:srgbClr"))
            if header and r_idx == 0:
                srgb.set("val", "0A6B63")
            elif r_idx % 2 == 1:
                srgb.set("val", "EEF3F8")
            else:
                srgb.set("val", "FFFFFF")
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = FONT_BODY
                    if header and r_idx == 0:
                        run.font.bold = True
                        run.font.color.rgb = WHITE
                        run.font.size = Pt(10.5)
                    else:
                        run.font.color.rgb = BODY
                        run.font.size = Pt(10.5)
                        run.font.bold = False


def add_bullets(slide, left, top, width, height, items, size=10.5, color=BODY, spacing=6):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(spacing if i else 0)
        p.space_after = Pt(0)
        p.level = 0
        run = p.add_run()
        set_run(run, f"–  {item}", size, color, False)
    return shape


def card(slide, left, top, width, height, fill=WHITE, border=BORDER):
    return add_rect(slide, left, top, width, height, fill, border, radius=True)


def accent_bar(slide, left, top, height, width=0.07):
    return add_rect(slide, left, top, width, height, ACCENT)


# ---------------------------------------------------------------------------
# Generated visuals (light, geometric — not emoji icons)
# ---------------------------------------------------------------------------
def gen_earth_vs_orbit():
    """Side-by-side scientific comparison diagrams on light backgrounds."""
    def make(kind):
        img = Image.new("RGB", (720, 720), (247, 249, 252))
        d = ImageDraw.Draw(img, "RGBA")
        # soft panel
        d.rounded_rectangle([40, 40, 680, 680], radius=28, fill=(238, 243, 248), outline=(208, 218, 230), width=2)
        cx, cy = 360, 340
        if kind == "earth":
            # gravity arrows down
            for x in (220, 360, 500):
                d.line([(x, 120), (x, 210)], fill=(107, 124, 143), width=4)
                d.polygon([(x - 12, 210), (x + 12, 210), (x, 232)], fill=(107, 124, 143))
            # slumped cells
            cells = [(300, 380, 70), (360, 400, 80), (420, 385, 65), (330, 450, 55), (400, 455, 60), (360, 500, 50)]
            for x, y, r in cells:
                d.ellipse([x - r, y - r, x + r, y + r], fill=(10, 107, 99, 90), outline=(10, 107, 99, 180), width=2)
            # scaffold bars
            d.line([(200, 520), (520, 520)], fill=(90, 107, 124), width=6)
            d.line([(240, 560), (480, 560)], fill=(90, 107, 124), width=4)
        else:
            # free-float cells in 3D structure
            positions = [
                (300, 260, 48), (380, 250, 52), (340, 320, 58), (420, 330, 46),
                (280, 370, 44), (360, 390, 50), (440, 400, 42), (320, 450, 46),
                (400, 460, 48), (360, 520, 40),
            ]
            for x, y, r in positions:
                d.ellipse([x - r, y - r, x + r, y + r], fill=(10, 107, 99, 110), outline=(7, 82, 76, 200), width=2)
            # orbit ring
            d.ellipse([160, 160, 560, 560], outline=(26, 138, 128, 120), width=3)
            d.ellipse([200, 200, 520, 520], outline=(26, 138, 128, 70), width=2)
        path = ASSET_DIR / f"diagram_{kind}.png"
        img.save(path, "PNG")
        return path

    return make("earth"), make("orbit")


def gen_bioprint_schematic():
    img = Image.new("RGB", (900, 520), (247, 249, 252))
    d = ImageDraw.Draw(img, "RGBA")
    d.rounded_rectangle([20, 20, 880, 500], radius=20, fill=(238, 243, 248), outline=(208, 218, 230), width=2)
    # printer head
    d.rounded_rectangle([120, 80, 280, 160], radius=8, fill=(11, 31, 51))
    d.rectangle([180, 160, 220, 220], fill=(10, 107, 99))
    d.ellipse([190, 215, 210, 235], fill=(26, 138, 128))
    # layers building up
    for i, y in enumerate(range(360, 250, -28)):
        w = 180 + i * 18
        x0 = 450 - w // 2
        d.rounded_rectangle([x0, y, x0 + w, y + 22], radius=6, fill=(10, 107, 99, 160 + i * 10), outline=(7, 82, 76))
    # nozzle stream
    d.line([(200, 235), (200, 360)], fill=(26, 138, 128), width=3)
    # labels as simple bars
    d.rounded_rectangle([620, 100, 820, 140], radius=6, fill=(215, 235, 232))
    d.rounded_rectangle([620, 170, 820, 210], radius=6, fill=(215, 235, 232))
    d.rounded_rectangle([620, 240, 820, 280], radius=6, fill=(215, 235, 232))
    path = ASSET_DIR / "bioprint_schematic.png"
    img.save(path, "PNG")
    return path


def gen_title_panel():
    """Abstract light geometric panel for title slide — no cartoon icons."""
    img = Image.new("RGB", (900, 1350), (238, 243, 248))
    d = ImageDraw.Draw(img, "RGBA")
    # subtle grid
    for x in range(0, 900, 40):
        d.line([(x, 0), (x, 1350)], fill=(208, 218, 230, 90), width=1)
    for y in range(0, 1350, 40):
        d.line([(0, y), (900, y)], fill=(208, 218, 230, 90), width=1)
    # large arc (planet horizon suggestion)
    d.ellipse([-200, 700, 1100, 1800], fill=(215, 235, 232), outline=(10, 107, 99, 60), width=3)
    d.ellipse([-120, 820, 1020, 1700], outline=(10, 107, 99, 90), width=2)
    # concentric rings
    d.ellipse([200, 180, 700, 680], outline=(10, 107, 99, 100), width=3)
    d.ellipse([260, 240, 640, 620], outline=(10, 107, 99, 70), width=2)
    d.ellipse([320, 300, 580, 560], outline=(10, 107, 99, 50), width=2)
    # center node
    d.ellipse([420, 400, 480, 460], fill=(10, 107, 99))
    # orbital path
    d.arc([80, 120, 820, 860], start=200, end=340, fill=(10, 107, 99, 140), width=3)
    path = ASSET_DIR / "title_panel.png"
    img.save(path, "PNG")
    return path


# ---------------------------------------------------------------------------
# Slides
# ---------------------------------------------------------------------------
def slide_01(prs, title_img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, 0.08, ACCENT)
    # Left content — vertically distributed to use the full plane
    add_text_box(slide, MARGIN, 0.55, 7.6, 0.35,
                 "CONFIDENTIAL INVESTOR PRESENTATION", 12, ACCENT, bold=True)
    add_text_box(slide, MARGIN, 1.35, 7.8, 2.8, "", 44, NAVY, bold=True, font=FONT_DISPLAY,
                 lines=[
                     ("Samuel Space", 46, NAVY, True, FONT_DISPLAY),
                     ("Organ and Tissue", 46, NAVY, True, FONT_DISPLAY),
                     ("Engineering", 46, NAVY, True, FONT_DISPLAY),
                 ])
    add_line(slide, MARGIN, 4.45, 4.0, ACCENT, 3)
    add_text_box(slide, MARGIN, 4.75, 7.6, 0.9,
                 "Manufacturing the next generation of human tissue — in orbit.",
                 18, BODY)
    add_rect(slide, MARGIN, 5.9, 7.6, 0.95, PANEL)
    accent_bar(slide, MARGIN, 5.9, 0.95, 0.08)
    add_text_box(slide, MARGIN + 0.3, 6.05, 7.1, 0.3, "Series A Financing", 14, NAVY, True)
    add_text_box(slide, MARGIN + 0.3, 6.4, 7.1, 0.3, "June 2026", 13, MUTED)
    # Right visual panel — full height
    slide.shapes.add_picture(str(title_img), Inches(8.35), Inches(0.08), Inches(4.98), Inches(7.34))
    add_rect(slide, 8.35, 0.08, 0.08, 7.34, ACCENT)
    page_number(slide, 1)


def slide_02(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Notice", "Confidential — Forward-Looking Statement", 2)
    card(slide, MARGIN, CONTENT_TOP, CONTENT_W, 5.55)
    accent_bar(slide, MARGIN, CONTENT_TOP, 5.55, 0.08)
    paras = [
        'This presentation has been prepared by Samuel Space Organ and Tissue Engineering ("the Company") solely to assist prospective investors in evaluating a potential investment. It is confidential and may not be reproduced or distributed, in whole or in part, without the Company\'s prior written consent.',
        "This presentation contains forward-looking statements, including projections of financial and operating performance, market size, regulatory timelines, and technology milestones. These statements reflect current expectations and assumptions and are subject to significant risks and uncertainties; actual results may differ materially. Nothing herein constitutes an offer to sell, or a solicitation of an offer to buy, any securities, and any such offer will be made only pursuant to definitive offering documents.",
        "Certain industry, market, and scientific data in this presentation has been compiled from third-party sources — including government agencies, peer-reviewed publications, and public company disclosures — as cited in the Appendix. The Company has not independently verified all such data and makes no representation as to its completeness or accuracy.",
    ]
    shape = slide.shapes.add_textbox(Inches(MARGIN + 0.35), Inches(CONTENT_TOP + 0.35),
                                    Inches(CONTENT_W - 0.6), Inches(5.0))
    tf = shape.text_frame
    tf.word_wrap = True
    for i, para in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(14 if i else 0)
        p.space_after = Pt(0)
        run = p.add_run()
        set_run(run, para, 13.5, BODY)


def slide_03(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Executive Summary",
                   "Translating Proven Orbital Science Into a Commercial Tissue Platform", 3)
    metrics = [
        ("103,223", "People on the U.S. transplant\nwaitlist today (OPTN, 2025)"),
        ("13 / day", "Americans die waiting for\na donor organ"),
        ("$1.7B–$3.5B+", "Global bioprinting market\n2025 to 2030 (~16% CAGR)"),
        ("8 years", "Continuous in-orbit\nbioprinting validation"),
    ]
    gap = 0.14
    mw = (CONTENT_W - 3 * gap) / 4
    mh = 1.45
    for i, (val, lab) in enumerate(metrics):
        x = MARGIN + i * (mw + gap)
        card(slide, x, CONTENT_TOP, mw, mh)
        accent_bar(slide, x, CONTENT_TOP, mh, 0.07)
        add_text_box(slide, x + 0.18, CONTENT_TOP + 0.18, mw - 0.3, 0.45, val, 20, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.18, CONTENT_TOP + 0.7, mw - 0.3, 0.65,
                     "", 11.5, BODY, lines=[(ln, 11.5, BODY, False, FONT_BODY) for ln in lab.split("\n")])

    pillars = [
        ("The Problem", "Organ and tissue shortages persist globally. Earth-based bioprinting hits a physical ceiling: gravity causes tissues to collapse without a temporary scaffold, limiting achievable geometry."),
        ("The Unlock", "Microgravity eliminates scaffold dependency. NASA, Redwire, ESA, LambdaVision and ROSCOSMOS have validated the science aboard the ISS continuously since 2018."),
        ("The Market", "A $1.7B market growing at ~16% CAGR, the first FDA-cleared tissue-engineered vascular graft (Symvess, Dec 2024), and four commercial stations entering service from 2026."),
        ("The Ask", "An $18M Series A to build Samuel's first flight-qualified payload platform, commercialise wound-care and orthopedic verticals, and file NASA InSPA Phase 1 grants."),
    ]
    py = CONTENT_TOP + mh + 0.14
    ph = CONTENT_BOTTOM - py
    pw = (CONTENT_W - 3 * gap) / 4
    for i, (title, body) in enumerate(pillars):
        x = MARGIN + i * (pw + gap)
        card(slide, x, py, pw, ph)
        add_rect(slide, x, py, pw, 0.08, ACCENT)
        numbered_badge(slide, x + 0.2, py + 0.28, i + 1, 0.34, ACCENT, 11)
        add_text_box(slide, x + 0.2, py + 0.78, pw - 0.4, 0.4, title, 16, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.2, py + 1.35, pw - 0.4, ph - 1.55, body, 12.5, BODY)


def slide_04(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "The Problem", "A Global Organ & Tissue Crisis", 4)
    # Chart left
    add_text_box(slide, MARGIN, CONTENT_TOP, 6.5, 0.3,
                 "U.S. Transplant Waiting List by Organ (2025)", 12, NAVY, True)
    chart_data = CategoryChartData()
    chart_data.categories = ["Kidney", "Liver", "Heart", "Lung", "Pancreas", "Other"]
    chart_data.add_series("Share (%)", (87.0, 9.0, 2.0, 1.0, 0.5, 0.5))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.BAR_CLUSTERED, Inches(MARGIN), Inches(CONTENT_TOP + 0.35),
        Inches(6.5), Inches(3.7), chart_data
    ).chart
    chart.has_legend = False
    plot = chart.plots[0]
    plot.has_data_labels = False
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT
    add_text_box(slide, MARGIN, CONTENT_TOP + 4.1, 6.5, 0.25,
                 "Source: OPTN / organdonor.gov, 2025", 8.5, MUTED)

    # Right stats
    stats = [
        ("941,652", "CVD deaths in the U.S. in 2022 — still the #1 cause of death, killing more people than all cancers combined"),
        ("~1 per 34 sec", "A cardiovascular death occurs in the U.S. on average every 34 seconds — roughly 2,500 per day"),
        ("10% met", "Of the estimated global transplant need. WHO reports ~130,000 transplants/yr vs. demand exceeding 1.3M"),
    ]
    rx, rw = 7.2, 5.75
    sh = 1.25
    for i, (v, t) in enumerate(stats):
        y = CONTENT_TOP + i * (sh + 0.15)
        card(slide, rx, y, rw, sh)
        accent_bar(slide, rx, y, sh, 0.07)
        add_text_box(slide, rx + 0.25, y + 0.2, 2.1, 0.85, v, 18, ACCENT, True, FONT_DISPLAY, anchor=MSO_ANCHOR.MIDDLE)
        add_text_box(slide, rx + 2.4, y + 0.2, rw - 2.65, 0.85, t, 11, BODY, anchor=MSO_ANCHOR.MIDDLE)
    add_text_box(slide, rx, CONTENT_TOP + 4.2, rw, 0.35,
                 "Source: AHA 2025 Heart Disease & Stroke Statistics; WHO Global Observatory on Donation and Transplantation.",
                 8, MUTED)

    # Bottom callout full width
    add_rect(slide, MARGIN, 6.15, CONTENT_W, 0.82, ACCENT_DARK)
    add_text_box(slide, MARGIN + 0.3, 6.28, CONTENT_W - 0.6, 0.55,
                 "Transplant recipients also face a lifetime of immunosuppressive therapy and rejection risk. Tissue printed from a patient's own cells eliminates cross-reactive rejection entirely — removing both the clinical burden and the ongoing drug cost.",
                 12, WHITE)


def slide_05(prs, earth_img, orbit_img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "The Core Constraint",
                   "Why Earth-Based Bioprinting Has a Hard Ceiling", 5)
    add_text_box(slide, MARGIN, CONTENT_TOP, CONTENT_W, 0.55,
                 "On Earth, gravity forces tissues to collapse under their own weight during printing. A scaffold must be built into every construct, adding complexity, introducing foreign material, and limiting achievable geometry. Microgravity removes this constraint entirely.",
                 12.5, BODY)
    left_x, right_x = MARGIN, MARGIN + CONTENT_W / 2 + 0.1
    pw = CONTENT_W / 2 - 0.1
    py, ph = CONTENT_TOP + 0.6, CONTENT_H - 0.65
    # Earth card
    card(slide, left_x, py, pw, ph)
    add_rect(slide, left_x, py, pw, 0.48, PANEL_ALT)
    add_text_box(slide, left_x, py, pw, 0.48, "EARTH GRAVITY", 13, BODY, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    img_size = min(3.1, pw - 0.8)
    slide.shapes.add_picture(str(earth_img), Inches(left_x + (pw - img_size) / 2), Inches(py + 0.65),
                             Inches(img_size), Inches(img_size))
    add_bullets(slide, left_x + 0.3, py + 0.75 + img_size, pw - 0.5, 1.5, [
        "Cells slump and fuse unevenly under gravitational load",
        "Scaffold material printed first, then must be dissolved",
        "Complex hollow structures (e.g. heart chambers) extremely difficult",
    ], 12, BODY, 5)

    # Orbit card
    card(slide, right_x, py, pw, ph)
    add_rect(slide, right_x, py, pw, 0.48, ACCENT)
    add_text_box(slide, right_x, py, pw, 0.48, "MICROGRAVITY — ISS ORBIT", 13, WHITE, True,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    slide.shapes.add_picture(str(orbit_img), Inches(right_x + (pw - img_size) / 2), Inches(py + 0.65),
                             Inches(img_size), Inches(img_size))
    add_bullets(slide, right_x + 0.3, py + 0.75 + img_size, pw - 0.5, 1.5, [
        "Tissue self-supports in true 3D — no scaffold required",
        "More uniform cell distribution and structural integrity",
        "Complex vascular and hollow geometries become achievable",
    ], 12, ACCENT_DARK, 5)


def slide_06(prs, schematic):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "The Technology",
                   "How Space-Based Bioprinting Works: Five-Step Process", 6)
    steps = [
        ("Cell Harvest", "Biopsy yields patient-specific cells, eliminating immunological rejection risk"),
        ("Stem Cell Reprogramming", "Nobel-winning iPSC method converts adult cells to pluripotent stem cells"),
        ("Bioink Formulation", "Cells suspended in hydrogel optimised for printing viscosity and cell viability"),
        ("In-Orbit Printing", "Layer-by-layer deposition in microgravity — no scaffold required"),
        ("Tissue Maturation", "Bioreactor conditioning allows cells to differentiate into target tissue type"),
    ]
    gap = 0.14
    sw = (CONTENT_W - 4 * gap) / 5
    for i, (title, body) in enumerate(steps):
        x = MARGIN + i * (sw + gap)
        card(slide, x, CONTENT_TOP, sw, 2.35)
        numbered_badge(slide, x + 0.15, CONTENT_TOP + 0.18, i + 1)
        # connector
        if i < 4:
            add_rect(slide, x + sw - 0.02, CONTENT_TOP + 1.05, gap + 0.04, 0.04, ACCENT_MID)
        add_text_box(slide, x + 0.15, CONTENT_TOP + 0.7, sw - 0.3, 0.55, title, 13, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.15, CONTENT_TOP + 1.3, sw - 0.3, 0.9, body, 10, BODY)

    # Bottom: schematic + why orbit
    by = CONTENT_TOP + 2.55
    bh = 3.0
    card(slide, MARGIN, by, 5.0, bh)
    slide.shapes.add_picture(str(schematic), Inches(MARGIN + 0.25), Inches(by + 0.25),
                             Inches(4.5), Inches(2.2))
    add_text_box(slide, MARGIN + 0.25, by + 2.5, 4.5, 0.35,
                 "Schematic: layer-by-layer orbital bioprinting hardware", 9.5, MUTED)

    card(slide, MARGIN + 5.2, by, CONTENT_W - 5.2, bh)
    accent_bar(slide, MARGIN + 5.2, by, bh, 0.08)
    add_text_box(slide, MARGIN + 5.5, by + 0.25, CONTENT_W - 5.7, 0.4,
                 "Why orbit?", 16, ACCENT, True, FONT_DISPLAY)
    add_text_box(slide, MARGIN + 5.5, by + 0.75, CONTENT_W - 5.7, 1.35,
                 "In the near-weightlessness of ISS orbit (~400 km altitude), tissues grow in true three dimensions with no gravitational compression. This principle has been validated continuously since 2018 across NASA, Redwire Corporation, ROSCOSMOS, LambdaVision Inc., and ESA/DLR — eight years of confirmation by organisations operating entirely independently of one another.",
                 11.5, BODY)
    add_line(slide, MARGIN + 5.5, by + 2.2, CONTENT_W - 5.9, LINE, 1)
    add_text_box(slide, MARGIN + 5.5, by + 2.35, CONTENT_W - 5.7, 0.55,
                 "Note: Adult cardiac muscle cannot self-regenerate — damaged tissue is replaced by scar tissue that blocks electrical signals, a key driver of CVD mortality and why cardiac is Samuel's highest-value target vertical.",
                 10, MUTED)


def slide_07(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Product Strategy", "Four High-Value Clinical Verticals", 7)
    items = [
        ("01", "Cardiac Tissue",
         "CVD kills ~2,500 Americans daily. Adult heart tissue cannot regenerate. Redwire returned the first live human cardiac tissue from orbit in May 2024 under BFF-Cardiac.",
         "Proven by: Redwire BFF-Cardiac, ISS, May 2024"),
        ("02", "Orthopedic Grafts",
         "Meniscal tears are among the most common U.S. military injuries. Redwire printed the first human knee meniscus in orbit in September 2023, establishing the foundational protocol.",
         "Proven by: Redwire BFF-Meniscus-2, ISS, Sept 2023"),
        ("03", "Artificial Retinas",
         "30 million people worldwide suffer from degenerative retinal disease. LambdaVision has produced consistent 200-layer protein retina films in microgravity across 9 ISS missions.",
         "Proven by: LambdaVision / Space Tango, 9 ISS missions"),
        ("04", "Wound Care Patches",
         "ESA and DLR demonstrated a handheld bioprinter creating customised skin patches aboard the ISS during the Cosmic Kiss mission. Lowest regulatory burden — Samuel's Phase 1 entry point.",
         "Proven by: ESA/DLR Bioprint FirstAid, ISS, Dec 2021"),
    ]
    gap = 0.18
    cw = (CONTENT_W - gap) / 2
    ch = (CONTENT_H - gap) / 2 - 0.05
    for i, (num, title, body, proof) in enumerate(items):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + gap)
        y = CONTENT_TOP + row * (ch + gap)
        card(slide, x, y, cw, ch)
        accent_bar(slide, x, y, ch, 0.08)
        add_text_box(slide, x + 0.3, y + 0.25, 0.8, 0.35, num, 18, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, x + 1.1, y + 0.28, cw - 1.4, 0.4, title, 17, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.3, y + 0.85, cw - 0.55, 1.15, body, 12, BODY)
        add_rect(slide, x + 0.3, y + ch - 0.55, cw - 0.55, 0.35, ACCENT_SOFT)
        add_text_box(slide, x + 0.4, y + ch - 0.55, cw - 0.75, 0.35, proof, 10, ACCENT_DARK, True,
                     anchor=MSO_ANCHOR.MIDDLE)


def slide_08(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Third-Party Validation",
                   "Eight Years of Proven In-Orbit Bioprinting", 8,
                   "Every milestone below was publicly announced by an independent organisation — not by Samuel Space.")
    events = [
        ("2018", "ROSCOSMOS Organ.Aut", "First magnetic 3D bioprinter on ISS; tissue constructs created 2018–2020 prove in-orbit approach"),
        ("2019", "Redwire BFF Launched", "BioFabrication Facility (Techshot/nScrypt) arrives on ISS; later wins Popular Science Best of What's New 2023"),
        ("Dec '21", "ESA Bioprint FirstAid", "Portable handheld bioprinter demonstrated by astronaut Matthias Maurer on Cosmic Kiss mission"),
        ("Sep '23", "First Knee Meniscus", "Redwire BFF prints first human knee meniscus in orbit — musculoskeletal milestone"),
        ("May '24", "First Cardiac Tissue", "Redwire returns first live 3D-bioprinted human cardiac tissue from ISS under BFF-Cardiac"),
        ("Sep '25", "Grants & Contracts", "LambdaVision wins NASA InSPA Phase 2; Redwire awarded $25M NASA IDIQ for continued ISS work"),
        ("2026", "Commercial Pivot", "LambdaVision closes $7M seed; commercial station payload slots reserved ahead of ISS 2030 retirement"),
    ]
    # Full-width timeline using the content band
    ty = CONTENT_TOP + CONTENT_H / 2
    add_line(slide, MARGIN, ty, CONTENT_W, ACCENT, 3)
    n = len(events)
    slot = CONTENT_W / n
    for i, (date, title, body) in enumerate(events):
        x = MARGIN + i * slot + 0.04
        cw = slot - 0.08
        node = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x + cw / 2 - 0.11), Inches(ty - 0.11),
                                     Inches(0.22), Inches(0.22))
        node.fill.solid()
        node.fill.fore_color.rgb = ACCENT
        node.line.fill.background()
        _strip_shadow(node)
        above = (i % 2 == 0)
        card_h = 2.25
        cy = ty - 0.28 - card_h if above else ty + 0.28
        # stem
        stem_y = cy + card_h if above else ty + 0.11
        stem_h = (ty - 0.11) - (cy + card_h) if above else cy - (ty + 0.11)
        if stem_h > 0.02:
            add_rect(slide, x + cw / 2 - 0.01, stem_y, 0.02, stem_h, ACCENT_MID)
        card(slide, x, cy, cw, card_h)
        add_text_box(slide, x + 0.1, cy + 0.12, cw - 0.2, 0.25, date, 10.5, ACCENT, True)
        add_text_box(slide, x + 0.1, cy + 0.42, cw - 0.2, 0.5, title, 12, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.1, cy + 1.0, cw - 0.2, 1.1, body, 10, BODY)


def slide_09(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Market Opportunity",
                   "A Rapidly Expanding Global Addressable Market", 9)
    add_text_box(slide, MARGIN, CONTENT_TOP, 6.8, 0.3,
                 "Global 3D Bioprinting Market 2023–2030 (USD $B)", 12, NAVY, True)
    chart_data = CategoryChartData()
    chart_data.categories = ["2023", "2024", "2025E", "2026E", "2027E", "2028E", "2029E", "2030E"]
    chart_data.add_series("Market ($B)", (1.28, 1.46, 1.7, 1.97, 2.28, 2.64, 3.06, 3.49))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED, Inches(MARGIN), Inches(CONTENT_TOP + 0.35),
        Inches(6.9), Inches(4.5), chart_data
    ).chart
    chart.has_legend = False
    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = ACCENT
    add_text_box(slide, MARGIN, 6.55, 6.9, 0.3,
                 "Source: Mordor Intelligence, 3D Bioprinting Market Report, Jan 2026  ·  ~15.9% CAGR",
                 8.5, MUTED)

    stats = [
        ("$3.5B+", "Projected global 3D bioprinting market by 2030 at ~15.9% CAGR (Mordor Intelligence, Jan 2026)"),
        ("Dec 2024", "FDA cleared Symvess — first tissue-engineered vascular graft — setting a direct regulatory precedent"),
        ("$420M+", "Committed by pharma companies to commercial-station research programs through 2030 (industry research)"),
        ("$200–400M", "Samuel's estimated initial SAM by 2030 (cardiac patches, ortho grafts, retinal constructs)"),
    ]
    rx, rw = 7.55, 5.4
    sh = 1.25
    for i, (v, t) in enumerate(stats):
        y = CONTENT_TOP + i * (sh + 0.12)
        card(slide, rx, y, rw, sh)
        accent_bar(slide, rx, y, sh, 0.07)
        add_text_box(slide, rx + 0.25, y + 0.18, rw - 0.45, 0.35, v, 18, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, rx + 0.25, y + 0.55, rw - 0.45, 0.6, t, 11, BODY)


def slide_10(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Market Timing",
                   "Four Converging Tailwinds Create a Unique Entry Window", 10)
    items = [
        ("~89%", "reduction in launch cost/kg since 2010", "Launch Costs Have Collapsed",
         "Cost-to-orbit has dropped from ~$54,000/kg (Space Shuttle era) to ~$6,000/kg (Falcon 9). Starship is expected to reduce this further. Research payloads that were cost-prohibitive a decade ago are now commercially viable."),
        ("4", "NASA-funded commercial station programs", "Commercial Stations Are Coming",
         "ISS retires ~2030. NASA's CLD program has funded four successors: Starlab (Voyager/Airbus, targeting 2029), Axiom Station, Orbital Reef (Blue Origin/Sierra Space), and Vast Haven-1. Samuel is designed station-agnostic from day one."),
        ("Dec 2024", "First FDA-cleared tissue-engineered vascular graft", "Regulatory Precedent Is Set",
         "The FDA cleared Symvess — the first tissue-engineered synthetic vascular graft — in December 2024. This is the single most important regulatory precedent for clinical bioprinted tissue and validates the FDA's review pathway."),
        ("$329M", "Raised by Varda Space (in-space manufacturing)", "Capital Is Already Flowing",
         "Varda Space Industries has raised $329M total, including a $187M Series C (July 2025) led by Natural Capital and Shrug Capital with Founders Fund and Khosla Ventures. LambdaVision closed a $7M seed atop $15M+ in non-dilutive grants."),
    ]
    gap = 0.16
    cw = (CONTENT_W - gap) / 2
    ch = (CONTENT_H - gap) / 2 - 0.02
    for i, (stat, sub, title, body) in enumerate(items):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + gap)
        y = CONTENT_TOP + row * (ch + gap)
        card(slide, x, y, cw, ch)
        accent_bar(slide, x, y, ch, 0.08)
        add_text_box(slide, x + 0.3, y + 0.22, cw - 0.5, 0.4, stat, 22, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.3, y + 0.65, cw - 0.5, 0.28, sub, 10, MUTED)
        add_text_box(slide, x + 0.3, y + 1.05, cw - 0.5, 0.35, title, 15, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.3, y + 1.5, cw - 0.5, 1.15, body, 11, BODY)


def slide_11(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Our Platform", "Station-Agnostic Payload Platform", 11)
    add_text_box(slide, MARGIN, CONTENT_TOP, 7.6, 0.7,
                 "Samuel Space is built as a payload integrator and tissue manufacturing platform — hardware-compatible across all emerging commercial LEO stations. The same flight-qualified platform, bioink libraries, and protocols deploy whether the payload flies on the ISS, Starlab, Axiom, Orbital Reef, or Vast.",
                 12, BODY)
    features = [
        ("01", "Modular Payload Design",
         "Hardware footprint matches standard double-locker form factor compatible with ISS, Starlab, and Axiom interfaces — one build, many stations."),
        ("02", "Proprietary Cell-Line IP",
         "Bioink formulations and printing protocols optimised for microgravity are the core defensible IP asset, accumulated across every mission flown."),
        ("03", "Service-First Revenue",
         "Fee-for-service bioprinting generates near-term revenue from pharma and academic customers while building the IP base that creates long-term enterprise value."),
    ]
    for i, (num, title, body) in enumerate(features):
        y = CONTENT_TOP + 0.85 + i * 1.15
        card(slide, MARGIN, y, 7.6, 1.05)
        accent_bar(slide, MARGIN, y, 1.05, 0.08)
        add_text_box(slide, MARGIN + 0.25, y + 0.15, 0.6, 0.3, num, 14, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, MARGIN + 0.95, y + 0.15, 6.3, 0.3, title, 14, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, MARGIN + 0.95, y + 0.5, 6.3, 0.45, body, 11, BODY)

    # Right rail — stations
    rx, rw = 8.3, 4.65
    add_text_box(slide, rx, CONTENT_TOP, rw, 0.3, "STATION ECOSYSTEM COMPATIBILITY", 10, ACCENT, True)
    add_rect(slide, rx, CONTENT_TOP + 0.4, rw, 0.55, ACCENT)
    add_text_box(slide, rx, CONTENT_TOP + 0.4, rw, 0.55, "Samuel Space Platform", 14, WHITE, True,
                 FONT_DISPLAY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    stations = [
        ("ISS", "Redwire BFF — active now through ~2030"),
        ("Starlab", "Voyager/Airbus JV — targeting 2029 launch"),
        ("Axiom Station", "First commercial modules ~2026–2028"),
        ("Orbital Reef", "Blue Origin/Sierra Space — CLD-funded"),
        ("Vast Haven-1", "First dedicated commercial station, 2026"),
    ]
    for i, (name, desc) in enumerate(stations):
        y = CONTENT_TOP + 1.1 + i * 0.78
        card(slide, rx, y, rw, 0.7)
        add_rect(slide, rx, y, 0.08, 0.7, ACCENT_MID)
        add_text_box(slide, rx + 0.25, y + 0.08, rw - 0.4, 0.28, name, 12.5, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, rx + 0.25, y + 0.36, rw - 0.4, 0.28, desc, 10, MUTED)
    add_text_box(slide, MARGIN, 6.75, 7.6, 0.25,
                 "Station names shown as market context only — not as confirmed Samuel partnerships.",
                 9, MUTED)


def slide_12(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Execution Plan",
                   "Four-Phase Technology & Commercialisation Roadmap", 12)
    phases = [
        ("PHASE 1", "Years 1–2", "Platform & Wound Care",
         ["ISS payload hardware certification", "Commercial wound-care patch trials",
          "NASA InSPA Phase 1 grant applications", "Research tissue models for pharma CROs"],
         "TRL 4 → 6", True),
        ("PHASE 2", "Years 3–5", "Orthopedic & Cardiac",
         ["510(k) filing for wound-care products", "Orthopedic graft BFF-compatible program",
          "Cardiac patch preclinical studies", "Series B fundraise"],
         "TRL 6 → 8", False),
        ("PHASE 3", "Years 6–9", "Retinal & Vascular",
         ["Clinical trials for retinal constructs", "FDA PMA pathway engagement",
          "Expand to Orbital Reef / Axiom", "Series C / pre-IPO raise"],
         "TRL 7 → 9", False),
        ("PHASE 4", "Year 10+", "Whole-Organ Platform",
         ["Vascularised full-organ constructs", "Commercial station scale-up",
          "Platform licensing to pharma majors", "IPO or strategic exit"],
         "TRL 9+", False),
    ]
    gap = 0.16
    pw = (CONTENT_W - 3 * gap) / 4
    ph = CONTENT_H - 0.05
    for i, (phase, years, title, items, trl, highlight) in enumerate(phases):
        x = MARGIN + i * (pw + gap)
        fill = ACCENT_SOFT if highlight else WHITE
        card(slide, x, CONTENT_TOP, pw, ph, fill=fill)
        if highlight:
            add_rect(slide, x, CONTENT_TOP, pw, 0.08, ACCENT)
        else:
            add_rect(slide, x, CONTENT_TOP, pw, 0.08, PANEL_ALT)
        add_text_box(slide, x + 0.18, CONTENT_TOP + 0.25, pw - 0.35, 0.25, phase, 11, ACCENT if highlight else NAVY, True)
        add_text_box(slide, x + 0.18, CONTENT_TOP + 0.52, pw - 0.35, 0.25, years, 11, MUTED)
        add_line(slide, x + 0.18, CONTENT_TOP + 0.85, pw - 0.36, ACCENT if highlight else LINE, 1.5)
        add_text_box(slide, x + 0.18, CONTENT_TOP + 1.05, pw - 0.35, 0.55, title, 15, NAVY, True, FONT_DISPLAY)
        add_bullets(slide, x + 0.18, CONTENT_TOP + 1.7, pw - 0.35, 2.8, items, 11.5,
                    ACCENT_DARK if highlight else BODY, 8)
        add_rect(slide, x + 0.18, CONTENT_TOP + ph - 0.55, pw - 0.36, 0.38, ACCENT if highlight else NAVY)
        add_text_box(slide, x + 0.18, CONTENT_TOP + ph - 0.55, pw - 0.36, 0.38, trl, 11, WHITE, True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_13(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Competitive Landscape",
                   "Samuel Space Targets the Space-Based Clinical Quadrant", 13)
    mx, my, mw, mh = MARGIN, CONTENT_TOP, CONTENT_W, CONTENT_H - 0.3
    card(slide, mx, my, mw, mh, fill=WHITE)
    mid_x = mx + mw / 2
    mid_y = my + 0.45 + (mh - 0.7) / 2

    # Target quadrant fill first (under axes/points)
    add_rect(slide, mid_x + 0.02, my + 0.45, mw / 2 - 0.22, (mh - 0.7) / 2, ACCENT_SOFT)

    add_line(slide, mx + 0.15, mid_y, mw - 0.3, ACCENT, 1.75)
    vline = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(mid_x), Inches(my + 0.45),
                                   Pt(1.75), Inches(mh - 0.7))
    vline.fill.solid()
    vline.fill.fore_color.rgb = ACCENT
    vline.line.fill.background()
    _strip_shadow(vline)

    add_text_box(slide, mx + 0.25, my + 0.12, mw / 2 - 0.4, 0.28, "EARTH-BASED", 12, MUTED, True)
    add_text_box(slide, mid_x + 0.2, my + 0.12, mw / 2 - 0.4, 0.28, "SPACE-BASED", 12, ACCENT, True)
    add_text_box(slide, mx + 0.2, my + 0.55, 2.4, 0.25, "CLINICAL / COMMERCIAL", 10, BODY, True)
    add_text_box(slide, mx + 0.2, my + mh - 0.35, 2.2, 0.25, "EARLY RESEARCH", 10, MUTED, True)

    points = [
        (mx + mw * 0.16, my + mh * 0.40, "Cellink / BICO", False),
        (mx + mw * 0.22, my + mh * 0.58, "Organovo", False),
        (mx + mw * 0.28, my + mh * 0.28, "3D Systems (Med)", False),
        (mx + mw * 0.66, my + mh * 0.48, "Redwire BFF", False),
        (mx + mw * 0.80, my + mh * 0.34, "LambdaVision", False),
        (mx + mw * 0.58, my + mh * 0.74, "ROSCOSMOS / 3D BPS", False),
        (mx + mw * 0.64, my + mh * 0.64, "Entry", False),
        (mx + mw * 0.74, my + mh * 0.20, "Samuel Space (target)", True),
    ]
    for x, y, label, highlight in points:
        size = 0.32 if highlight else 0.20
        oval = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(size), Inches(size))
        oval.fill.solid()
        oval.fill.fore_color.rgb = ACCENT if highlight else CHART_4
        oval.line.color.rgb = WHITE if highlight else CHART_4
        oval.line.width = Pt(1.5)
        _strip_shadow(oval)
        add_text_box(slide, x - 0.2, y + size + 0.04, 1.7, 0.4, label,
                     11 if highlight else 10, ACCENT if highlight else BODY, highlight)

    add_text_box(slide, MARGIN, 6.85, CONTENT_W, 0.25,
                 "Positions are illustrative, based on public technology and commercialisation-stage disclosures.",
                 9, MUTED)


def slide_14(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Revenue Model", "Three Complementary Revenue Streams", 14)
    streams = [
        ("01", "REVENUE STREAM 1", "Contract Bioprinting Services",
         "Fee-for-service orbital bioprinting missions for pharmaceutical companies, academic medical centres, and government research agencies, modelled on how Redwire's BFF and Space Tango currently operate commercially. Revenue begins in Year 1.",
         "Near-Term Revenue"),
        ("02", "REVENUE STREAM 2", "Non-Dilutive Grant Funding",
         "NASA InSPA, ESA BIC, DoD SBIR/STTR, and National Eye Institute mechanisms provide substantial non-dilutive capital. LambdaVision has secured $15M+ cumulatively; Redwire holds a $25M NASA IDIQ. Samuel targets $4–8M in grants across Years 1–4.",
         "Reduces Equity Burn"),
        ("03", "REVENUE STREAM 3", "Tissue IP Licensing & Royalties",
         "As proprietary bioink formulations and protocols demonstrate clinical-grade output, Samuel licenses IP to pharmaceutical and device manufacturers. This becomes the primary long-term value driver, targeted from Year 3 onward.",
         "Long-Term Value Driver"),
    ]
    sh = (CONTENT_H - 0.24) / 3
    for i, (num, label, title, body, tag) in enumerate(streams):
        y = CONTENT_TOP + i * (sh + 0.12)
        card(slide, MARGIN, y, CONTENT_W, sh)
        accent_bar(slide, MARGIN, y, sh, 0.08)
        add_text_box(slide, MARGIN + 0.3, y + 0.18, 1.0, 0.3, num, 16, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, MARGIN + 1.3, y + 0.18, 3.5, 0.25, label, 10, ACCENT, True)
        add_text_box(slide, MARGIN + 1.3, y + 0.45, 8.0, 0.35, title, 17, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, MARGIN + 1.3, y + 0.9, 8.2, 0.7, body, 12, BODY)
        add_rect(slide, MARGIN + CONTENT_W - 2.6, y + 0.35, 2.25, 0.35, ACCENT_SOFT)
        add_text_box(slide, MARGIN + CONTENT_W - 2.6, y + 0.35, 2.25, 0.35, tag, 10, ACCENT_DARK, True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_15(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Go-to-Market",
                   "Sequenced Market Entry Across Four Partner Categories", 15)
    items = [
        ("1", "Commercial Station Operators",
         "Redwire (active BFF operator), Starlab (Voyager/Airbus JV), Axiom Space, Blue Origin Orbital Reef, and Vast Haven-1. Samuel negotiates research-service agreements and payload manifesting rights with each.",
         "Redwire, Starlab, Axiom, Vast"),
        ("2", "Payload Integrators & CROs",
         "Partners such as Space Tango, LambdaVision's enabling integrator, provide environmental control, crew interface, and downlink services. Samuel outsources non-core integration to move faster and preserve capital.",
         "Space Tango, Nanoracks/Voyager"),
        ("3", "Pharma & Academic Medical Centres",
         "Initial contract-service revenue comes from pharmaceutical R&D divisions seeking orbital tissue models for drug screening, and from NIH-funded academic biotech groups with institutional mandates.",
         "Pharma R&D, NIH-funded labs"),
        ("4", "Government & Grant Agencies",
         "NASA (InSPA, SBIR/STTR), ESA, NIH National Eye Institute, and DoD MTEC provide non-dilutive grant capital that de-risks early development and validates the technology for commercial buyers.",
         "NASA, ESA, NIH NEI, DoD MTEC"),
    ]
    gap = 0.16
    pw = (CONTENT_W - 3 * gap) / 4
    ph = CONTENT_H - 0.05
    for i, (num, title, body, tag) in enumerate(items):
        x = MARGIN + i * (pw + gap)
        card(slide, x, CONTENT_TOP, pw, ph)
        add_rect(slide, x, CONTENT_TOP, pw, 0.08, ACCENT)
        numbered_badge(slide, x + 0.2, CONTENT_TOP + 0.3, num, 0.42, ACCENT, 14)
        add_text_box(slide, x + 0.2, CONTENT_TOP + 0.95, pw - 0.4, 0.75, title, 15, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.2, CONTENT_TOP + 1.85, pw - 0.4, 2.8, body, 11.5, BODY)
        add_rect(slide, x + 0.2, CONTENT_TOP + ph - 0.55, pw - 0.4, 0.38, ACCENT_SOFT)
        add_text_box(slide, x + 0.25, CONTENT_TOP + ph - 0.55, pw - 0.5, 0.38, tag, 9.5, ACCENT_DARK, True,
                     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def slide_16(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Financials — Illustrative",
                   "Five-Year Revenue Model (Illustrative Projections)", 16,
                   "Figures are illustrative projections for investor discussion purposes only. Actual results will vary materially.")
    chart_data = CategoryChartData()
    chart_data.categories = ["Yr 1", "Yr 2", "Yr 3", "Yr 4", "Yr 5"]
    chart_data.add_series("Contract Services ($M)", (0.5, 1.2, 3.0, 7.0, 14.0))
    chart_data.add_series("Grant Revenue ($M)", (0.3, 1.0, 1.5, 2.0, 2.5))
    chart_data.add_series("IP Licensing ($M)", (0.0, 0.2, 1.6, 4.5, 7.7))
    chart_h = CONTENT_H - 0.1
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_STACKED, Inches(MARGIN), Inches(CONTENT_TOP),
        Inches(7.15), Inches(chart_h), chart_data
    ).chart
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    colors = [CHART_1, CHART_2, CHART_3]
    for s, c in zip(chart.series, colors):
        s.format.fill.solid()
        s.format.fill.fore_color.rgb = c

    rows = [
        ("Year", "Total Rev.", "Key Driver"),
        ("Year 1", "$0.8M", "First service contracts + NASA grant"),
        ("Year 2", "$2.4M", "Expanded services + IP licensing launch"),
        ("Year 3", "$6.1M", "Cardiac/ortho milestones + Starlab start"),
        ("Year 4", "$13.5M", "IP licensing deals + multi-station ops"),
        ("Year 5", "$24.2M", "Full vertical scale-up across four verticals"),
    ]
    tw = CONTENT_W - 7.3
    table_shape = slide.shapes.add_table(6, 3, Inches(MARGIN + 7.3), Inches(CONTENT_TOP),
                                        Inches(tw), Inches(chart_h))
    table = table_shape.table
    table.columns[0].width = Inches(1.05)
    table.columns[1].width = Inches(1.15)
    table.columns[2].width = Inches(tw - 2.2)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            table.cell(r, c).text = val
    style_table(table)


def slide_17(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "The Ask",
                   "An $18M Series A to Reach First Commercial Flight", 17)
    top_h = 2.15
    card(slide, MARGIN, CONTENT_TOP, 4.35, top_h, fill=ACCENT_SOFT)
    accent_bar(slide, MARGIN, CONTENT_TOP, top_h, 0.08)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + 0.2, 3.8, 0.25, "ROUND SIZE", 11, ACCENT, True)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + 0.5, 3.8, 0.65, "$18M", 40, NAVY, True, FONT_DISPLAY)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + 1.25, 3.8, 0.35, "Series A Preferred Equity", 14, BODY, True)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + 1.65, 3.8, 0.35,
                 "Allocation figures are illustrative and subject to diligence.", 10, MUTED)

    kpis = [
        ("~30 months", "Runway to next value-inflection milestone"),
        ("2", "Product verticals reaching clinical pilot stage"),
        ("3–5", "Targeted commercial flight missions secured"),
    ]
    kw = (CONTENT_W - 4.35 - 0.3) / 3
    for i, (v, t) in enumerate(kpis):
        x = MARGIN + 4.5 + i * (kw + 0.1)
        card(slide, x, CONTENT_TOP, kw, top_h)
        accent_bar(slide, x, CONTENT_TOP, top_h, 0.07)
        add_text_box(slide, x + 0.15, CONTENT_TOP + 0.4, kw - 0.3, 0.55, v, 22, ACCENT, True, FONT_DISPLAY,
                     align=PP_ALIGN.CENTER)
        add_text_box(slide, x + 0.15, CONTENT_TOP + 1.15, kw - 0.3, 0.8, t, 12, BODY,
                     align=PP_ALIGN.CENTER)

    by = CONTENT_TOP + top_h + 0.14
    bh = CONTENT_BOTTOM - by
    add_text_box(slide, MARGIN, by, 5.5, 0.28, "Use of Funds", 14, NAVY, True, FONT_DISPLAY)
    chart_data = CategoryChartData()
    chart_data.categories = ["R&D / Payload Dev", "Flight & Integration", "Regulatory & Clinical",
                             "Team Build-Out", "G&A"]
    chart_data.add_series("Use of Funds", (40.0, 25.0, 15.0, 15.0, 5.0))
    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.DOUGHNUT, Inches(MARGIN), Inches(by + 0.3),
        Inches(3.5), Inches(bh - 0.35), chart_data
    ).chart
    chart.has_legend = False
    colors = [CHART_1, CHART_2, CHART_3, CHART_4, CHART_5]
    for s in chart.series:
        try:
            for i, pt in enumerate(s.points):
                pt.format.fill.solid()
                pt.format.fill.fore_color.rgb = colors[i]
        except Exception:
            s.format.fill.solid()
            s.format.fill.fore_color.rgb = ACCENT

    legend = [
        ("40%", "R&D / Payload Development"),
        ("25%", "Flight & Integration Costs"),
        ("15%", "Regulatory & Clinical Prep"),
        ("15%", "Team Build-Out"),
        ("5%", "G&A"),
    ]
    for i, (pct, name) in enumerate(legend):
        y = by + 0.45 + i * 0.48
        swatch = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(MARGIN + 3.6), Inches(y + 0.05),
                                       Inches(0.22), Inches(0.22))
        swatch.fill.solid()
        swatch.fill.fore_color.rgb = colors[i]
        swatch.line.fill.background()
        _strip_shadow(swatch)
        add_text_box(slide, MARGIN + 3.95, y, 0.7, 0.3, pct, 12, NAVY, True)
        add_text_box(slide, MARGIN + 4.6, y, 2.7, 0.3, name, 12, BODY)

    card(slide, 8.0, by, CONTENT_W - (8.0 - MARGIN), bh)
    accent_bar(slide, 8.0, by, bh, 0.08)
    add_text_box(slide, 8.3, by + 0.2, 4.5, 0.35, "Key Use-of-Funds Milestones", 14, NAVY, True, FONT_DISPLAY)
    add_bullets(slide, 8.3, by + 0.7, 4.5, bh - 0.9, [
        "Flight-qualify modular payload hardware",
        "Complete 3–5 commissioned ISS missions",
        "File first NASA InSPA Phase 1 grant",
        "Secure Starlab payload reservation",
        "Build core scientific & engineering team",
    ], 13, BODY, 10)


def slide_18(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Risk Factors", "Key Risks and Mitigating Strategies", 18)
    rows = [
        ("Risk Category", "Risk", "Mitigant"),
        ("Technical",
         "Scaling from validated tissue patches to full vascularised organs remains scientifically unsolved industry-wide.",
         "Phased roadmap monetises near-term, lower-complexity verticals (wound care, orthopedic) before attempting whole-organ printing."),
        ("Regulatory",
         "FDA approval pathways for novel tissue-engineered products are long, costly, and precedent-limited.",
         "Symvess clearance (Dec 2024) establishes the first direct precedent. Samuel engages regulatory counsel from Year 1 and targets lower-risk 510(k) products first."),
        ("Platform / Access",
         "ISS retirement (~2030) creates dependency risk if commercial successor stations are delayed.",
         "Station-agnostic hardware design supports ISS, Starlab, Axiom, Orbital Reef, and Vast simultaneously, diversifying single-platform exposure."),
        ("Capital Intensity",
         "Flight missions, payload hardware, and clinical trials require sustained, significant capital investment.",
         "Blended funding strategy combines equity with non-dilutive NASA/ESA/DoD grants, following the precedent set by LambdaVision and Redwire."),
        ("Competitive",
         "Incumbent operators (Redwire, LambdaVision) hold first-mover relationships with NASA and station operators.",
         "Samuel differentiates on multi-vertical platform breadth and station-agnostic flexibility rather than competing head-on for any single contract."),
    ]
    table_shape = slide.shapes.add_table(6, 3, Inches(MARGIN), Inches(CONTENT_TOP),
                                        Inches(CONTENT_W), Inches(5.5))
    table = table_shape.table
    table.columns[0].width = Inches(1.8)
    table.columns[1].width = Inches(5.2)
    table.columns[2].width = Inches(CONTENT_W - 7.0)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = val
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    style_table(table)


def slide_19(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Organisation", "Team & Scientific Advisory Approach", 19)
    add_text_box(slide, MARGIN, CONTENT_TOP, CONTENT_W, 0.55,
                 "Samuel Space is building a founding team and Scientific Advisory Board around the specific expertise this platform requires. Roles below reflect the Company's hiring plan and target advisory profile for the Series A period.",
                 12, BODY)
    roles = [
        ("01", "Founder / Chief Executive Officer",
         "Commercial space or biotech venture leadership; track record building regulated hardware or therapeutics businesses from seed to scale."),
        ("02", "Chief Scientific Officer",
         "PhD-level expertise in tissue engineering, stem cell biology, or regenerative medicine; prior peer-reviewed bioprinting research."),
        ("03", "VP, Payload Engineering",
         "Aerospace hardware engineering background; experience flight-qualifying instrumentation for ISS or commercial LEO payloads."),
        ("04", "Head of Regulatory Affairs",
         "FDA medical device / biologics regulatory pathway expertise; experience navigating 510(k) and PMA submissions."),
    ]
    gap = 0.16
    cw = (CONTENT_W - gap) / 2
    ch = 1.45
    for i, (num, title, body) in enumerate(roles):
        col, row = i % 2, i // 2
        x = MARGIN + col * (cw + gap)
        y = CONTENT_TOP + 0.7 + row * (ch + 0.14)
        card(slide, x, y, cw, ch)
        accent_bar(slide, x, y, ch, 0.08)
        add_text_box(slide, x + 0.25, y + 0.2, 0.6, 0.3, num, 14, ACCENT, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.9, y + 0.22, cw - 1.15, 0.35, title, 14, NAVY, True, FONT_DISPLAY)
        add_text_box(slide, x + 0.25, y + 0.7, cw - 0.5, 0.6, body, 11.5, BODY)

    # SAB
    y = CONTENT_TOP + 0.7 + 2 * (ch + 0.14) + 0.05
    card(slide, MARGIN, y, CONTENT_W, 1.55, fill=ACCENT_SOFT)
    accent_bar(slide, MARGIN, y, 1.55, 0.08)
    add_text_box(slide, MARGIN + 0.3, y + 0.15, CONTENT_W - 0.5, 0.3,
                 "Scientific Advisory Board — In Formation", 14, NAVY, True, FONT_DISPLAY)
    add_text_box(slide, MARGIN + 0.3, y + 0.55, CONTENT_W - 0.5, 0.85,
                 "Samuel is actively recruiting an SAB spanning the disciplines proven essential by the field's existing leaders: stem-cell and cardiac tissue biology (in the tradition of Harvard's bioengineering research community), microgravity payload design (as demonstrated by Redwire's BFF engineering team), and protein-based biomanufacturing (as practiced by LambdaVision and Space Tango). Target SAB appointments will be announced as the round closes.",
                 11.5, BODY)


def slide_20(prs, title_img):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_bg(slide, WHITE)
    add_rect(slide, 0, 0, W, 0.08, ACCENT)
    slide.shapes.add_picture(str(title_img), Inches(0), Inches(0.08), Inches(4.9), Inches(7.34))
    add_rect(slide, 4.82, 0.08, 0.08, 7.34, ACCENT)
    add_text_box(slide, 5.3, 0.7, 7.5, 0.35, "THE OPPORTUNITY", 13, ACCENT, True)
    add_text_box(slide, 5.3, 1.35, 7.5, 1.7, "", 32, NAVY, True, FONT_DISPLAY,
                 lines=[
                     ("Every organ printed in orbit", 32, NAVY, True, FONT_DISPLAY),
                     ("is a life that didn't have to wait.", 32, NAVY, True, FONT_DISPLAY),
                 ])
    add_line(slide, 5.3, 3.3, 3.4, ACCENT, 3)
    add_text_box(slide, 5.3, 3.65, 7.4, 1.7,
                 "Samuel Space stands at the intersection of two converging trends: a mature, validated microgravity bioprinting science base, and the arrival of commercial space stations built for exactly this kind of work. The science already works. What's needed now is a company built to commercialise it.",
                 15, BODY)
    add_rect(slide, 5.3, 5.7, 7.4, 1.15, PANEL)
    accent_bar(slide, 5.3, 5.7, 1.15, 0.08)
    add_text_box(slide, 5.6, 5.9, 6.9, 0.35, "Samuel Space Organ and Tissue Engineering", 14, NAVY, True)
    add_text_box(slide, 5.6, 6.3, 6.9, 0.35, "investors@samuelspace.com  (placeholder)", 13, MUTED)
    page_number(slide, 20)


def slide_21(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    section_header(slide, "Appendix", "Sources & References", 21)
    left_items = [
        'NASA — "3D Bioprinting" research overview, iss-research program page',
        "NASA — BioFabrication Facility mission page, Redwire Corporation",
        "Redwire Corporation — Newsroom: BFF-Meniscus first human knee meniscus printed in orbit, Sept 2023",
        "Redwire Corporation — BFF-Cardiac investigation announcements and ISS return coverage, 2024",
        "Redwire Corporation — NASA IDIQ contract award coverage, Sept 2025",
        "ESA / DLR — Bioprint FirstAid investigation overview, Cosmic Kiss mission materials",
        "LambdaVision Inc. — Press releases: Space Tango partnership, NASA InSPA Phase 2 award, seed financing, 2025–2026",
    ]
    right_items = [
        "ROSCOSMOS — 3D MBP / Organ.Aut magnetic bioprinter mission summaries, 2018–2020",
        "American Heart Association — 2025 Heart Disease & Stroke Statistics Update",
        "OPTN / organdonor.gov (HRSA) — U.S. transplant waiting list data, 2025",
        "World Health Organization — Global Observatory on Donation and Transplantation",
        "Mordor Intelligence — 3D Bioprinting Market Report, updated January 2026",
        "U.S. FDA — Symvess (acellular tissue-engineered vessel) clearance announcement, Dec 2024",
        "NASA Commercial LEO Destinations (CLD) program — Starlab, Axiom, Orbital Reef, Vast program pages",
    ]
    gap = 0.2
    cw = (CONTENT_W - gap) / 2
    ch = 3.85
    card(slide, MARGIN, CONTENT_TOP, cw, ch)
    accent_bar(slide, MARGIN, CONTENT_TOP, ch, 0.08)
    add_text_box(slide, MARGIN + 0.25, CONTENT_TOP + 0.18, cw - 0.4, 0.35,
                 "Bioprinting Science & Technology", 13, NAVY, True, FONT_DISPLAY)
    add_bullets(slide, MARGIN + 0.25, CONTENT_TOP + 0.6, cw - 0.45, 3.1, left_items, 10.5, BODY, 5)

    card(slide, MARGIN + cw + gap, CONTENT_TOP, cw, ch)
    accent_bar(slide, MARGIN + cw + gap, CONTENT_TOP, ch, 0.08)
    add_text_box(slide, MARGIN + cw + gap + 0.25, CONTENT_TOP + 0.18, cw - 0.4, 0.35,
                 "Market, Regulatory & Industry Data", 13, NAVY, True, FONT_DISPLAY)
    add_bullets(slide, MARGIN + cw + gap + 0.25, CONTENT_TOP + 0.6, cw - 0.45, 3.1, right_items, 10.5, BODY, 5)

    card(slide, MARGIN, CONTENT_TOP + ch + 0.15, CONTENT_W, 1.45, fill=PANEL)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + ch + 0.3, CONTENT_W - 0.5, 0.3,
                 "A Note on Methodology", 12, NAVY, True, FONT_DISPLAY)
    add_text_box(slide, MARGIN + 0.3, CONTENT_TOP + ch + 0.65, CONTENT_W - 0.5, 0.7,
                 "All third-party scientific milestones, dates, and figures in this presentation are drawn from the public sources listed above and were current as of the citation date shown. Samuel Space Organ and Tissue Engineering is an independent platform company; references to NASA, Redwire, LambdaVision, ROSCOSMOS, and ESA/DLR describe industry validation of the underlying science, not partnerships, endorsements, or affiliations with Samuel Space.",
                 11, BODY)


def blank_layout(prs):
    """Ensure a blank layout exists; fall back to index 6 or 0."""
    for i, layout in enumerate(prs.slide_layouts):
        name = (layout.name or "").lower()
        if "blank" in name:
            return layout
    # Prefer last layout (often blank) else first
    if len(prs.slide_layouts) > 6:
        return prs.slide_layouts[6]
    return prs.slide_layouts[0]


def main():
    ASSET_DIR.mkdir(parents=True, exist_ok=True)
    earth_img, orbit_img = gen_earth_vs_orbit()
    schematic = gen_bioprint_schematic()
    title_img = gen_title_panel()

    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)

    # python-pptx creates default layouts; use blank-ish layout by clearing via index
    # We'll monkey-patch add_slide callers by setting a consistent blank layout helper
    blank = blank_layout(prs)

    def add():
        return prs.slides.add_slide(blank)

    # Rebuild helpers to use blank layout — patch slide_* to use our add
    # Simpler: temporarily replace prs.slide_layouts[6] usage by wrapping
    original_add = prs.slides.add_slide

    def add_slide_blank(_layout):
        return original_add(blank)

    prs.slides.add_slide = add_slide_blank  # type: ignore

    slide_01(prs, title_img)
    slide_02(prs)
    slide_03(prs)
    slide_04(prs)
    slide_05(prs, earth_img, orbit_img)
    slide_06(prs, schematic)
    slide_07(prs)
    slide_08(prs)
    slide_09(prs)
    slide_10(prs)
    slide_11(prs)
    slide_12(prs)
    slide_13(prs)
    slide_14(prs)
    slide_15(prs)
    slide_16(prs)
    slide_17(prs)
    slide_18(prs)
    slide_19(prs)
    slide_20(prs, title_img)
    slide_21(prs)

    # Fix competitive matrix z-order issue: rebuild is fine as-is
    prs.save(str(OUT))
    print(f"Wrote {OUT} with {len(prs.slides)} slides")


if __name__ == "__main__":
    main()
