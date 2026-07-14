from __future__ import annotations

from math import cos, pi, sin
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


OUTPUT = Path("Samuel Space Organ and Tissue Engineering - Investor Presentation (2).pptx")

SLIDE_W = 13.333
SLIDE_H = 7.5
MARGIN = 0.55
CONTENT_W = 12.23

FONT = "Lato"
FONT_DISPLAY = "Lato"

BG = "F7F9F7"
WHITE = "FFFFFF"
NAVY = "10243D"
INK = "23384D"
SLATE = "5E7082"
MUTED = "8492A0"
LINE = "DCE4E7"
TEAL = "00A58E"
TEAL_DARK = "087C71"
MINT = "DDF4EE"
MINT_2 = "ECF8F4"
BLUE = "E8F0F7"
BLUE_2 = "F1F5F8"
GOLD = "E5A647"
GOLD_PALE = "FAF0DF"
CORAL = "D87566"
CORAL_PALE = "FAECE9"
LAVENDER = "EEEAF7"


def I(value: float):
    return Inches(value)


def C(value: str):
    value = value.lstrip("#")
    return RGBColor(int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16))


prs = Presentation()
prs.slide_width = I(SLIDE_W)
prs.slide_height = I(SLIDE_H)
blank = prs.slide_layouts[6]
prs.core_properties.title = "Samuel Space Organ and Tissue Engineering — Investor Presentation"
prs.core_properties.subject = "Series A financing presentation"
prs.core_properties.author = "Samuel Space Organ and Tissue Engineering"
prs.core_properties.keywords = "microgravity, bioprinting, tissue engineering, investor presentation"
prs.core_properties.comments = "Redesigned for a professional, light-background investor presentation."


def set_background(slide, color: str = BG):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = C(color)


def add_shape(
    slide,
    shape_type,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str | None = None,
    line: str | None = None,
    line_width: float = 0.75,
    radius: bool = False,
):
    if radius:
        shape_type = MSO_SHAPE.ROUNDED_RECTANGLE
    shape = slide.shapes.add_shape(shape_type, I(x), I(y), I(w), I(h))
    if fill is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = C(fill)
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = C(line)
        shape.line.width = Pt(line_width)
    return shape


def add_line(
    slide,
    x1: float,
    y1: float,
    x2: float,
    y2: float,
    color: str = LINE,
    width: float = 1.0,
):
    line = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, I(x1), I(y1), I(x2), I(y2)
    )
    line.line.color.rgb = C(color)
    line.line.width = Pt(width)
    return line


def add_text(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float,
    size: float = 12,
    color: str = INK,
    bold: bool = False,
    font: str = FONT,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    italic: bool = False,
    margin: float = 0,
    line_spacing: float = 1.0,
):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = C(color)
    return box


def add_rich_text(
    slide,
    runs: list[dict],
    x: float,
    y: float,
    w: float,
    h: float,
    align=PP_ALIGN.LEFT,
    valign=MSO_ANCHOR.TOP,
    margin: float = 0,
):
    box = slide.shapes.add_textbox(I(x), I(y), I(w), I(h))
    tf = box.text_frame
    tf.clear()
    tf.margin_left = I(margin)
    tf.margin_right = I(margin)
    tf.margin_top = I(margin)
    tf.margin_bottom = I(margin)
    tf.word_wrap = True
    tf.vertical_anchor = valign
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    for spec in runs:
        run = p.add_run()
        run.text = spec["text"]
        run.font.name = spec.get("font", FONT)
        run.font.size = Pt(spec.get("size", 12))
        run.font.bold = spec.get("bold", False)
        run.font.italic = spec.get("italic", False)
        run.font.color.rgb = C(spec.get("color", INK))
    return box


def card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str = WHITE,
    line: str = LINE,
    line_width: float = 0.75,
):
    return add_shape(
        slide,
        MSO_SHAPE.ROUNDED_RECTANGLE,
        x,
        y,
        w,
        h,
        fill,
        line,
        line_width,
        radius=True,
    )


def pill(
    slide,
    text: str,
    x: float,
    y: float,
    w: float,
    h: float = 0.32,
    fill: str = MINT,
    color: str = TEAL_DARK,
    size: float = 8.5,
):
    card(slide, x, y, w, h, fill=fill, line=fill)
    add_text(
        slide,
        text,
        x + 0.06,
        y,
        w - 0.12,
        h,
        size=size,
        color=color,
        bold=True,
        align=PP_ALIGN.CENTER,
        valign=MSO_ANCHOR.MIDDLE,
    )


def add_footer(slide, number: int):
    add_line(slide, MARGIN, 7.1, 12.78, 7.1, LINE, 0.7)
    add_text(slide, "CONFIDENTIAL", MARGIN, 7.18, 1.6, 0.16, 7, MUTED, True)
    add_text(
        slide,
        "SAMUEL SPACE",
        5.55,
        7.18,
        2.25,
        0.16,
        7,
        MUTED,
        True,
        align=PP_ALIGN.CENTER,
    )
    add_text(
        slide,
        f"{number:02d} / 21",
        11.82,
        7.18,
        0.96,
        0.16,
        7,
        MUTED,
        True,
        align=PP_ALIGN.RIGHT,
    )


def base_slide(section: str, title: str, number: int, title_size: float = 24):
    slide = prs.slides.add_slide(blank)
    set_background(slide)
    add_shape(slide, MSO_SHAPE.RECTANGLE, MARGIN, 0.27, 0.08, 0.24, TEAL, TEAL)
    add_text(slide, section.upper(), 0.71, 0.26, 5.6, 0.23, 8.5, TEAL_DARK, True)
    add_text(
        slide,
        "INVESTOR PRESENTATION  /  JUNE 2026",
        9.12,
        0.26,
        3.66,
        0.23,
        7.5,
        MUTED,
        True,
        align=PP_ALIGN.RIGHT,
    )
    add_text(
        slide,
        title,
        MARGIN,
        0.62,
        12.23,
        0.56,
        title_size,
        NAVY,
        True,
        FONT_DISPLAY,
        valign=MSO_ANCHOR.MIDDLE,
    )
    add_line(slide, MARGIN, 1.34, 12.78, 1.34, LINE, 0.8)
    add_footer(slide, number)
    return slide


def add_bullet_rows(
    slide,
    items: list[str],
    x: float,
    y: float,
    w: float,
    item_h: float,
    size: float = 10,
    color: str = INK,
    marker: str = TEAL,
):
    for idx, text in enumerate(items):
        yy = y + idx * item_h
        add_shape(
            slide,
            MSO_SHAPE.ROUNDED_RECTANGLE,
            x,
            yy + 0.12,
            0.08,
            0.08,
            marker,
            marker,
            radius=True,
        )
        add_text(
            slide,
            text,
            x + 0.2,
            yy,
            w - 0.2,
            item_h,
            size,
            color,
            valign=MSO_ANCHOR.TOP,
            line_spacing=0.95,
        )


def metric_card(
    slide,
    x: float,
    y: float,
    w: float,
    h: float,
    value: str,
    label: str,
    source: str | None = None,
    accent: str = TEAL,
):
    card(slide, x, y, w, h, WHITE, LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.08, h, accent, accent)
    add_text(slide, value, x + 0.22, y + 0.16, w - 0.42, 0.48, 22, NAVY, True)
    add_text(slide, label, x + 0.22, y + 0.67, w - 0.42, 0.36, 9, SLATE, True)
    if source:
        add_text(slide, source, x + 0.22, y + h - 0.28, w - 0.42, 0.16, 6.8, MUTED)


def draw_cell_cluster(slide, x: float, y: float, w: float, h: float, microgravity: bool):
    if microgravity:
        add_shape(slide, MSO_SHAPE.OVAL, x + 0.62, y + 0.12, w - 1.24, h - 0.24, MINT_2, TEAL, 1.2)
        cells = [
            (0.28, 0.27, 0.30),
            (0.48, 0.20, 0.24),
            (0.66, 0.34, 0.28),
            (0.40, 0.45, 0.34),
            (0.59, 0.53, 0.24),
            (0.29, 0.62, 0.26),
            (0.71, 0.65, 0.30),
        ]
        for idx, (rx, ry, rs) in enumerate(cells):
            size = rs * min(w, h)
            fill = TEAL if idx % 3 == 0 else ("58C9B7" if idx % 3 == 1 else "A9E3D7")
            add_shape(
                slide,
                MSO_SHAPE.OVAL,
                x + rx * w - size / 2,
                y + ry * h - size / 2,
                size,
                size,
                fill,
                WHITE,
                0.7,
            )
        for angle in (0, pi / 2):
            x1 = x + w / 2 + cos(angle) * w * 0.38
            y1 = y + h / 2 + sin(angle) * h * 0.35
            x2 = x + w / 2 - cos(angle) * w * 0.38
            y2 = y + h / 2 - sin(angle) * h * 0.35
            add_line(slide, x1, y1, x2, y2, "9FDCCF", 0.7)
    else:
        for i in range(8):
            xx = x + 0.35 + i * (w - 0.7) / 7
            add_line(slide, xx, y + h * 0.62, xx, y + h * 0.96, "AEBCC6", 0.7)
        for j in range(4):
            yy = y + h * 0.62 + j * h * 0.11
            add_line(slide, x + 0.25, yy, x + w - 0.25, yy, "AEBCC6", 0.7)
        cells = [
            (0.18, 0.41, 0.25),
            (0.32, 0.48, 0.31),
            (0.48, 0.53, 0.27),
            (0.61, 0.46, 0.33),
            (0.77, 0.55, 0.26),
            (0.88, 0.46, 0.22),
        ]
        for idx, (rx, ry, rs) in enumerate(cells):
            ww = rs * h
            hh = ww * 0.66
            add_shape(
                slide,
                MSO_SHAPE.OVAL,
                x + rx * w - ww / 2,
                y + ry * h - hh / 2,
                ww,
                hh,
                "8CCFC3" if idx % 2 else "55B9AA",
                WHITE,
                0.7,
            )
        add_line(slide, x + 0.22, y + h * 0.6, x + w - 0.22, y + h * 0.6, CORAL, 1.5)
        for xx in (x + w * 0.35, x + w * 0.7):
            add_line(slide, xx, y + 0.1, xx, y + h * 0.3, CORAL, 1.1)
            add_shape(slide, MSO_SHAPE.CHEVRON, xx - 0.07, y + h * 0.27, 0.14, 0.16, CORAL, CORAL)


def draw_hardware_schematic(slide, x: float, y: float, w: float, h: float):
    card(slide, x, y, w, h, BLUE_2, LINE)
    add_text(slide, "FLIGHT PAYLOAD  /  FUNCTIONAL SCHEMATIC", x + 0.22, y + 0.18, w - 0.44, 0.2, 7.5, SLATE, True)
    rack_x, rack_y = x + 0.45, y + 0.62
    rack_w, rack_h = w - 0.9, h - 0.98
    add_shape(slide, MSO_SHAPE.RECTANGLE, rack_x, rack_y, rack_w, rack_h, WHITE, NAVY, 1.2)
    add_shape(slide, MSO_SHAPE.RECTANGLE, rack_x + 0.18, rack_y + 0.18, rack_w - 0.36, 0.38, BLUE, LINE)
    add_text(slide, "CONTROL + TELEMETRY", rack_x + 0.28, rack_y + 0.25, rack_w - 0.56, 0.16, 6.8, NAVY, True, align=PP_ALIGN.CENTER)
    chamber_x = rack_x + 0.46
    chamber_y = rack_y + 0.79
    chamber_w = rack_w - 0.92
    chamber_h = rack_h - 1.03
    add_shape(slide, MSO_SHAPE.RECTANGLE, chamber_x, chamber_y, chamber_w, chamber_h, MINT_2, TEAL_DARK, 1.0)
    add_shape(slide, MSO_SHAPE.RECTANGLE, chamber_x + 0.56, chamber_y + 0.16, chamber_w - 1.12, 0.17, NAVY, NAVY)
    add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, chamber_x + chamber_w / 2 - 0.12, chamber_y + 0.31, 0.24, 0.28, TEAL, TEAL)
    add_line(slide, chamber_x + chamber_w / 2, chamber_y + 0.56, chamber_x + chamber_w / 2, chamber_y + 0.86, TEAL_DARK, 1.2)
    add_shape(slide, MSO_SHAPE.OVAL, chamber_x + chamber_w / 2 - 0.22, chamber_y + 0.83, 0.44, 0.18, "A7E2D6", TEAL_DARK, 0.8)
    add_shape(slide, MSO_SHAPE.RECTANGLE, chamber_x + 0.24, chamber_y + chamber_h - 0.42, chamber_w - 0.48, 0.18, NAVY, NAVY)
    add_text(slide, "SEALED BIOPRINTING CHAMBER", chamber_x + 0.12, chamber_y + chamber_h - 0.2, chamber_w - 0.24, 0.13, 6.3, SLATE, True, align=PP_ALIGN.CENTER)


def section_bar(slide, text: str, x: float, y: float, w: float, fill: str = NAVY):
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, 0.36, fill, fill, radius=True)
    add_text(slide, text, x + 0.12, y, w - 0.24, 0.36, 8, WHITE, True, valign=MSO_ANCHOR.MIDDLE)


# Slide 1 — Cover
slide = prs.slides.add_slide(blank)
set_background(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, SLIDE_H, TEAL, TEAL)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.16, 0, 0.05, SLIDE_H, NAVY, NAVY)
add_text(slide, "CONFIDENTIAL  /  INVESTOR PRESENTATION", 0.72, 0.6, 5.8, 0.26, 8.5, TEAL_DARK, True)
add_text(
    slide,
    "Samuel Space\nOrgan and Tissue\nEngineering",
    0.72,
    1.22,
    6.75,
    2.15,
    29,
    NAVY,
    True,
    FONT_DISPLAY,
    line_spacing=0.91,
)
add_text(
    slide,
    "Manufacturing the next generation of human tissue — in orbit.",
    0.74,
    3.62,
    6.15,
    0.62,
    15,
    INK,
    False,
    line_spacing=1.0,
)
add_line(slide, 0.74, 4.56, 6.78, 4.56, LINE, 1.0)
pill(slide, "SERIES A FINANCING", 0.74, 5.0, 2.02, 0.38, NAVY, WHITE, 8.2)
pill(slide, "JUNE 2026", 2.92, 5.0, 1.34, 0.38, MINT, TEAL_DARK, 8.2)
add_text(
    slide,
    "A station-agnostic platform for microgravity bioprinting.",
    0.74,
    5.65,
    5.9,
    0.32,
    10,
    SLATE,
)
add_text(slide, "SAMUEL SPACE", 0.74, 6.78, 2.2, 0.2, 7.5, MUTED, True)
add_text(slide, "01 / 21", 5.78, 6.78, 0.72, 0.2, 7.5, MUTED, True, align=PP_ALIGN.RIGHT)

# Cover visual
card(slide, 7.58, 0.38, 5.2, 6.73, MINT_2, "BFDCD5", 0.9)
add_text(slide, "MICROGRAVITY TISSUE PLATFORM", 7.96, 0.78, 4.42, 0.22, 7.5, TEAL_DARK, True, align=PP_ALIGN.RIGHT)
add_shape(slide, MSO_SHAPE.OVAL, 8.48, 1.5, 3.72, 3.72, WHITE, "A9D8CE", 1.0)
add_shape(slide, MSO_SHAPE.OVAL, 8.85, 1.87, 2.98, 2.98, MINT, "8DCEC0", 1.0)
add_shape(slide, MSO_SHAPE.OVAL, 9.3, 2.32, 2.08, 2.08, "B8E8DD", TEAL, 1.3)
for angle in range(0, 360, 45):
    a = angle * pi / 180
    cx = 10.34 + cos(a) * 1.25
    cy = 3.36 + sin(a) * 1.25
    add_shape(slide, MSO_SHAPE.OVAL, cx - 0.09, cy - 0.09, 0.18, 0.18, TEAL, WHITE, 0.6)
for angle in range(0, 360, 60):
    a = angle * pi / 180
    cx = 10.34 + cos(a) * 0.67
    cy = 3.36 + sin(a) * 0.67
    add_shape(slide, MSO_SHAPE.OVAL, cx - 0.12, cy - 0.12, 0.24, 0.24, "48BDAE", WHITE, 0.6)
add_shape(slide, MSO_SHAPE.OVAL, 10.04, 3.06, 0.6, 0.6, TEAL_DARK, WHITE, 0.8)
add_line(slide, 8.05, 5.58, 12.3, 5.58, "B6DAD2", 0.8)
add_text(slide, "ORBITAL ADVANTAGE", 8.05, 5.8, 1.8, 0.2, 7.2, TEAL_DARK, True)
add_text(slide, "Scaffold-free 3D geometry", 8.05, 6.12, 2.08, 0.42, 10, NAVY, True)
add_text(slide, "PLATFORM MODEL", 10.38, 5.8, 1.55, 0.2, 7.2, TEAL_DARK, True)
add_text(slide, "One payload. Many stations.", 10.38, 6.12, 1.95, 0.42, 10, NAVY, True)


# Slide 2 — Notice
slide = base_slide("Notice", "Confidential — Forward-Looking Statement", 2)
card(slide, 0.55, 1.56, 3.05, 5.35, NAVY, NAVY)
add_text(slide, "FOR DISCUSSION\nONLY", 0.86, 1.94, 2.4, 0.76, 17, WHITE, True, line_spacing=0.9)
add_text(slide, "This material is:", 0.86, 3.0, 2.2, 0.22, 8.5, "A9B9C8", True)
for idx, item in enumerate(("Confidential", "Non-binding", "Subject to diligence")):
    yy = 3.42 + idx * 0.68
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.86, yy + 0.05, 0.12, 0.12, TEAL, TEAL)
    add_text(slide, item, 1.14, yy, 1.95, 0.24, 10.5, WHITE, True)
add_line(slide, 0.86, 5.62, 3.26, 5.62, "324860", 0.8)
add_text(
    slide,
    "No securities are offered by this presentation.",
    0.86,
    5.88,
    2.24,
    0.62,
    9.2,
    "CED8E0",
)
add_text(slide, "SAMUEL SPACE", 0.86, 6.55, 2.0, 0.18, 7.5, TEAL, True)

notice_sections = [
    (
        "01",
        "Purpose and confidentiality",
        'This presentation has been prepared by Samuel Space Organ and Tissue Engineering ("the Company") solely to assist prospective investors in evaluating a potential investment. It is confidential and may not be reproduced or distributed, in whole or in part, without the Company\'s prior written consent.',
    ),
    (
        "02",
        "Forward-looking statements",
        "This presentation contains forward-looking statements, including projections of financial and operating performance, market size, regulatory timelines, and technology milestones. These statements reflect current expectations and assumptions and are subject to significant risks and uncertainties; actual results may differ materially. Nothing herein constitutes an offer to sell, or a solicitation of an offer to buy, any securities, and any such offer will be made only pursuant to definitive offering documents.",
    ),
    (
        "03",
        "Third-party information",
        "Certain industry, market, and scientific data in this presentation has been compiled from third-party sources — including government agencies, peer-reviewed publications, and public company disclosures — as cited in the Appendix. The Company has not independently verified all such data and makes no representation as to its completeness or accuracy.",
    ),
]
for idx, (num, heading, body) in enumerate(notice_sections):
    yy = 1.56 + idx * 1.78
    card(slide, 3.83, yy, 8.95, 1.57, WHITE, LINE)
    pill(slide, num, 4.08, yy + 0.22, 0.48, 0.32, MINT, TEAL_DARK, 8)
    add_text(slide, heading, 4.77, yy + 0.18, 7.58, 0.28, 11.5, NAVY, True)
    add_text(slide, body, 4.08, yy + 0.62, 8.28, 0.75, 8.8, SLATE, line_spacing=0.96)


# Slide 3 — Executive summary
slide = base_slide("Executive Summary", "Translating Proven Orbital Science Into a Commercial Tissue Platform", 3, 22.5)
metrics = [
    ("103,223", "People on the U.S. transplant waitlist", "OPTN, 2025"),
    ("13 / day", "Americans die waiting for a donor organ", "U.S. organ donation data"),
    ("$1.7B to $3.5B+", "Global bioprinting market, 2025 to 2030", "~16% CAGR"),
    ("8 years", "Continuous in-orbit bioprinting validation", "2018–2026"),
]
for idx, (value, label, source) in enumerate(metrics):
    metric_card(slide, 0.55 + idx * 3.08, 1.57, 2.87, 1.33, value, label, source)

pillars = [
    (
        "01",
        "THE PROBLEM",
        "Organ and tissue shortages persist globally. Earth-based bioprinting hits a physical ceiling: gravity causes tissues to collapse without a temporary scaffold, limiting achievable geometry.",
        CORAL_PALE,
        CORAL,
    ),
    (
        "02",
        "THE UNLOCK",
        "Microgravity eliminates scaffold dependency. NASA, Redwire, ESA, LambdaVision and ROSCOSMOS have validated the science aboard the ISS continuously since 2018.",
        MINT_2,
        TEAL,
    ),
    (
        "03",
        "THE MARKET",
        "A $1.7B market growing at ~16% CAGR, the first FDA-cleared tissue-engineered vascular graft (Symvess, Dec 2024), and four commercial stations entering service from 2026.",
        BLUE_2,
        "5E91B5",
    ),
    (
        "04",
        "THE ASK",
        "An $18M Series A to build Samuel's first flight-qualified payload platform, commercialise wound-care and orthopedic verticals, and file NASA InSPA Phase 1 grants.",
        NAVY,
        TEAL,
    ),
]
for idx, (num, heading, body, fill, accent) in enumerate(pillars):
    x = 0.55 + idx * 3.08
    card(slide, x, 3.12, 2.87, 3.82, fill, fill if fill == NAVY else LINE)
    pill(slide, num, x + 0.22, 3.36, 0.5, 0.32, accent, WHITE, 8)
    heading_color = WHITE if fill == NAVY else NAVY
    body_color = "CAD6DF" if fill == NAVY else SLATE
    add_text(slide, heading, x + 0.22, 3.91, 2.4, 0.3, 9.2, accent if fill != NAVY else TEAL, True)
    add_text(slide, heading.title(), x + 0.22, 4.34, 2.38, 0.44, 16, heading_color, True)
    add_text(slide, body, x + 0.22, 4.94, 2.4, 1.62, 9.5, body_color, line_spacing=0.98)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.22, 6.68, 0.62, 0.05, accent, accent)


# Slide 4 — The problem
slide = base_slide("The Problem", "A Global Organ & Tissue Crisis", 4)
card(slide, 0.55, 1.56, 7.12, 4.86, WHITE, LINE)
add_text(slide, "U.S. TRANSPLANT WAITLIST BY ORGAN  /  2025", 0.83, 1.83, 5.7, 0.24, 8.2, NAVY, True)
add_text(slide, "Share of 103,223 people waiting", 0.83, 2.14, 5.7, 0.2, 7.6, MUTED)
organs = [("Kidney", 87.0), ("Liver", 9.0), ("Heart", 2.0), ("Lung", 1.0), ("Pancreas", 0.5), ("Other", 0.5)]
for idx, (label, value) in enumerate(organs):
    yy = 2.55 + idx * 0.56
    add_text(slide, label, 0.83, yy, 0.88, 0.22, 8.8, SLATE, True)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, yy + 0.03, 4.85, 0.16, "EDF1F3", "EDF1F3", radius=True)
    bar_w = max(0.06, 4.85 * value / 87.0)
    add_shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, 1.8, yy + 0.03, bar_w, 0.16, TEAL if idx == 0 else "73CABB", TEAL if idx == 0 else "73CABB", radius=True)
    add_text(slide, f"{value:g}%", 6.75, yy, 0.52, 0.22, 8.8, NAVY, True, align=PP_ALIGN.RIGHT)
add_text(slide, "Source: OPTN / organdonor.gov, 2025", 0.83, 6.05, 4.0, 0.16, 6.7, MUTED)

stats = [
    ("941,652", "CVD deaths in the U.S. in 2022 — still the #1 cause of death, killing more people than all cancers combined"),
    ("~1 per 34 sec", "A cardiovascular death occurs in the U.S. on average every 34 seconds — roughly 2,500 per day"),
    ("10% met", "Of estimated global transplant need: ~130,000 procedures annually versus demand above 1.3M"),
]
for idx, (value, body) in enumerate(stats):
    yy = 1.56 + idx * 1.62
    card(slide, 7.89, yy, 4.89, 1.4, BLUE_2 if idx != 1 else MINT_2, LINE)
    add_text(slide, value, 8.16, yy + 0.2, 1.5, 0.43, 17, TEAL_DARK, True)
    add_text(slide, body, 9.78, yy + 0.19, 2.7, 0.88, 8.8, SLATE, line_spacing=0.95)
add_text(
    slide,
    "Source: AHA 2025 Heart Disease & Stroke Statistics; WHO Global Observatory on Donation and Transplantation.",
    8.02,
    6.1,
    4.5,
    0.22,
    6.4,
    MUTED,
)
card(slide, 0.55, 6.55, 12.23, 0.42, NAVY, NAVY)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.79, 6.67, 0.22, 0.06, TEAL, TEAL)
add_text(
    slide,
    "Patient-specific tissue can eliminate cross-reactive rejection — reducing both clinical burden and lifetime immunosuppressive drug cost.",
    1.2,
    6.56,
    11.2,
    0.4,
    9.5,
    WHITE,
    True,
    valign=MSO_ANCHOR.MIDDLE,
)


# Slide 5 — Core constraint
slide = base_slide("The Core Constraint", "Why Earth-Based Bioprinting Has a Hard Ceiling", 5)
card(slide, 0.55, 1.54, 12.23, 0.68, BLUE_2, LINE)
add_text(
    slide,
    "Gravity forces printed tissue to collapse under its own weight. Scaffolds add foreign material, increase process complexity, and limit geometry. Microgravity removes that constraint.",
    0.82,
    1.55,
    11.69,
    0.66,
    10.8,
    INK,
    valign=MSO_ANCHOR.MIDDLE,
)

comparison = [
    (
        0.55,
        "EARTH GRAVITY",
        "Scaffold-dependent",
        [
            "Cells slump and fuse unevenly under load",
            "Scaffold prints first, then must be dissolved",
            "Complex hollow geometries remain difficult",
        ],
        False,
        CORAL,
        CORAL_PALE,
    ),
    (
        6.78,
        "MICROGRAVITY  /  ISS ORBIT",
        "True 3D self-support",
        [
            "Tissue self-supports without a scaffold",
            "Cell distribution and integrity improve",
            "Vascular and hollow forms become achievable",
        ],
        True,
        TEAL,
        MINT_2,
    ),
]
for x, label, sub, bullets, micro, accent, fill in comparison:
    card(slide, x, 2.42, 6.0, 4.55, fill, accent, 1.0)
    pill(slide, label, x + 0.25, 2.67, 2.2 if micro else 1.55, 0.34, accent, WHITE, 7.5)
    add_text(slide, sub, x + 0.25, 3.16, 3.0, 0.3, 14, NAVY, True)
    draw_cell_cluster(slide, x + 0.42, 3.55, 2.48, 2.2, micro)
    add_bullet_rows(slide, bullets, x + 3.18, 3.6, 2.5, 0.73, 8.8, INK, accent)
    add_text(
        slide,
        "PHYSICAL CONSTRAINT" if not micro else "ORBITAL ADVANTAGE",
        x + 3.18,
        6.2,
        2.3,
        0.2,
        7.2,
        accent,
        True,
    )
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 3.18, 6.52, 2.05, 0.05, accent, accent)


# Slide 6 — Technology
slide = base_slide("The Technology", "How Space-Based Bioprinting Works: Five-Step Process", 6)
steps = [
    ("01", "Cell Harvest", "Biopsy yields patient-specific cells, reducing immunological rejection risk."),
    ("02", "Stem Cell Reprogramming", "The Nobel-winning iPSC method converts adult cells into pluripotent stem cells."),
    ("03", "Bioink Formulation", "Cells are suspended in hydrogel tuned for viscosity and viability."),
    ("04", "In-Orbit Printing", "Layer-by-layer deposition in microgravity requires no scaffold."),
    ("05", "Tissue Maturation", "Bioreactor conditioning differentiates cells into target tissue."),
]
step_w = 2.36
for idx, (num, title, body) in enumerate(steps):
    x = 0.55 + idx * 2.47
    card(slide, x, 1.57, step_w, 1.76, WHITE if idx != 3 else MINT_2, TEAL if idx == 3 else LINE, 0.9)
    pill(slide, num, x + 0.18, 1.76, 0.46, 0.3, TEAL if idx == 3 else NAVY, WHITE, 7.5)
    add_text(slide, title, x + 0.18, 2.18, step_w - 0.36, 0.36, 11.2, NAVY, True)
    add_text(slide, body, x + 0.18, 2.62, step_w - 0.36, 0.52, 8.1, SLATE, line_spacing=0.94)
    if idx < 4:
        add_line(slide, x + step_w, 2.45, x + 2.47, 2.45, TEAL, 1.4)

draw_hardware_schematic(slide, 0.55, 3.58, 4.33, 3.39)
card(slide, 5.1, 3.58, 7.68, 3.39, NAVY, NAVY)
pill(slide, "WHY ORBIT?", 5.4, 3.86, 1.28, 0.34, TEAL, WHITE, 7.5)
add_text(
    slide,
    "Near-weightlessness lets tissue grow in true three dimensions.",
    5.4,
    4.4,
    6.9,
    0.5,
    16,
    WHITE,
    True,
)
add_text(
    slide,
    "The principle has been validated continuously since 2018 across NASA, Redwire Corporation, ROSCOSMOS, LambdaVision, and ESA/DLR — eight years of confirmation by independent organisations.",
    5.4,
    5.02,
    4.78,
    1.08,
    9.6,
    "C7D2DC",
    line_spacing=0.98,
)
add_shape(slide, MSO_SHAPE.RECTANGLE, 10.43, 4.98, 0.05, 1.42, "365068", "365068")
add_text(slide, "8", 10.78, 4.92, 1.2, 0.66, 28, TEAL, True)
add_text(slide, "YEARS OF\nVALIDATION", 10.8, 5.62, 1.42, 0.54, 8, WHITE, True, line_spacing=0.9)
add_line(slide, 5.4, 6.35, 12.42, 6.35, "365068", 0.8)
add_text(
    slide,
    "Adult cardiac muscle cannot self-regenerate; scar tissue blocks electrical signals — making cardiac Samuel's highest-value target vertical.",
    5.4,
    6.51,
    6.95,
    0.3,
    7.8,
    "AEBCC8",
    italic=True,
)


# Slide 7 — Product strategy
slide = base_slide("Product Strategy", "Four High-Value Clinical Verticals", 7)
verticals = [
    (
        "01",
        "Cardiac Tissue",
        "CVD kills ~2,500 Americans daily. Adult heart tissue cannot regenerate. Redwire returned the first live human cardiac tissue from orbit in May 2024 under BFF-Cardiac.",
        "PROVEN BY  /  Redwire BFF-Cardiac, ISS, May 2024",
        CORAL_PALE,
        CORAL,
    ),
    (
        "02",
        "Orthopedic Grafts",
        "Meniscal tears are among the most common U.S. military injuries. Redwire printed the first human knee meniscus in orbit in September 2023, establishing the foundational protocol.",
        "PROVEN BY  /  Redwire BFF-Meniscus-2, ISS, Sept 2023",
        BLUE_2,
        "5E91B5",
    ),
    (
        "03",
        "Artificial Retinas",
        "30 million people worldwide suffer from degenerative retinal disease. LambdaVision has produced consistent 200-layer protein retina films in microgravity across 9 ISS missions.",
        "PROVEN BY  /  LambdaVision + Space Tango, 9 ISS missions",
        LAVENDER,
        "7D6CB0",
    ),
    (
        "04",
        "Wound Care Patches",
        "ESA and DLR demonstrated a handheld bioprinter creating customised skin patches aboard the ISS during the Cosmic Kiss mission. Lowest regulatory burden — Samuel's Phase 1 entry point.",
        "PROVEN BY  /  ESA/DLR Bioprint FirstAid, ISS, Dec 2021",
        MINT_2,
        TEAL,
    ),
]
for idx, (num, title, body, proof, fill, accent) in enumerate(verticals):
    col = idx % 2
    row = idx // 2
    x = 0.55 + col * 6.17
    y = 1.56 + row * 2.73
    card(slide, x, y, 6.0, 2.55, fill, LINE)
    pill(slide, num, x + 0.24, y + 0.23, 0.5, 0.32, accent, WHITE, 8)
    add_text(slide, title, x + 0.94, y + 0.18, 4.7, 0.38, 15, NAVY, True)
    add_text(slide, body, x + 0.24, y + 0.78, 5.5, 0.96, 9.4, SLATE, line_spacing=0.98)
    add_line(slide, x + 0.24, y + 1.92, x + 5.76, y + 1.92, "CBD7DC", 0.7)
    add_text(slide, proof, x + 0.24, y + 2.08, 5.5, 0.22, 7.1, accent, True)


# Slide 8 — Validation timeline
slide = base_slide("Third-Party Validation", "Eight Years of Proven In-Orbit Bioprinting", 8)
card(slide, 0.55, 1.54, 12.23, 0.52, NAVY, NAVY)
add_text(
    slide,
    "Independent milestones validate the market premise and demonstrate commercial feasibility — none were announced by Samuel Space.",
    0.78,
    1.54,
    11.77,
    0.52,
    9,
    WHITE,
    valign=MSO_ANCHOR.MIDDLE,
)
events = [
    ("2018", "ROSCOSMOS Organ.Aut", "First magnetic bioprinter reaches ISS; constructs prove the approach."),
    ("2019", "Redwire BFF launched", "BioFabrication Facility reaches ISS; later recognised by Popular Science."),
    ("DEC '21", "ESA Bioprint FirstAid", "Handheld skin-patch bioprinter demonstrated on Cosmic Kiss."),
    ("SEP '23", "First knee meniscus", "First human knee meniscus is printed in orbit."),
    ("MAY '24", "First cardiac tissue", "First live 3D-bioprinted human cardiac tissue returns."),
    ("SEP '25", "Grants + contracts", "InSPA Phase 2 and a $25M NASA IDIQ expand the field."),
    ("2026", "Commercial pivot", "$7M seed and payload reservations signal commercial transition."),
]
add_line(slide, 1.0, 4.28, 12.35, 4.28, TEAL_DARK, 1.5)
for idx, (date, title, body) in enumerate(events):
    cx = 1.15 + idx * 1.82
    top = idx % 2 == 0
    add_shape(slide, MSO_SHAPE.OVAL, cx - 0.09, 4.19, 0.18, 0.18, TEAL, WHITE, 0.8)
    if top:
        add_line(slide, cx, 4.19, cx, 3.83, TEAL_DARK, 0.9)
        yy = 2.27
    else:
        add_line(slide, cx, 4.37, cx, 4.73, TEAL_DARK, 0.9)
        yy = 4.73
    card(slide, cx - 0.73, yy, 1.46, 1.38, WHITE, LINE)
    pill(slide, date, cx - 0.57, yy + 0.15, 1.14, 0.29, MINT, TEAL_DARK, 7.1)
    add_text(slide, title, cx - 0.59, yy + 0.55, 1.18, 0.31, 8.2, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, body, cx - 0.62, yy + 0.9, 1.24, 0.4, 7.0, SLATE, align=PP_ALIGN.CENTER, line_spacing=0.9)
add_text(slide, "SCIENCE VALIDATED", 0.68, 6.75, 2.1, 0.16, 7.1, MUTED, True)
add_text(slide, "COMMERCIAL ACCESS EMERGING", 10.08, 6.75, 2.54, 0.16, 7.1, TEAL_DARK, True, align=PP_ALIGN.RIGHT)


# Slide 9 — Market opportunity
slide = base_slide("Market Opportunity", "A Rapidly Expanding Global Addressable Market", 9)
card(slide, 0.55, 1.56, 7.16, 5.37, WHITE, LINE)
add_text(slide, "GLOBAL 3D BIOPRINTING MARKET  /  USD $B", 0.83, 1.84, 5.8, 0.22, 8.2, NAVY, True)
add_text(slide, "~15.9% CAGR", 5.92, 1.84, 1.42, 0.22, 8.2, TEAL_DARK, True, align=PP_ALIGN.RIGHT)
years = ["2023", "2024", "2025E", "2026E", "2027E", "2028E", "2029E", "2030E"]
values = [1.28, 1.46, 1.70, 1.97, 2.28, 2.64, 3.06, 3.49]
chart_x, chart_y, chart_w, chart_h = 1.05, 2.43, 6.18, 3.55
for grid_value in (1, 2, 3):
    yy = chart_y + chart_h - grid_value / 3.6 * chart_h
    add_line(slide, chart_x, yy, chart_x + chart_w, yy, "E7ECEE", 0.7)
    add_text(slide, f"${grid_value}B", 0.68, yy - 0.1, 0.3, 0.18, 6.8, MUTED, align=PP_ALIGN.RIGHT)
bar_w = 0.48
gap = (chart_w - len(values) * bar_w) / (len(values) - 1)
for idx, (year, value) in enumerate(zip(years, values)):
    xx = chart_x + idx * (bar_w + gap)
    hh = value / 3.6 * chart_h
    fill = TEAL if idx >= 2 else "8FD7CA"
    add_shape(slide, MSO_SHAPE.RECTANGLE, xx, chart_y + chart_h - hh, bar_w, hh, fill, fill)
    add_text(slide, f"{value:.2f}", xx - 0.1, chart_y + chart_h - hh - 0.24, bar_w + 0.2, 0.18, 6.8, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, year, xx - 0.12, chart_y + chart_h + 0.12, bar_w + 0.24, 0.18, 6.8, SLATE, True, align=PP_ALIGN.CENTER)
add_text(slide, "Source: Mordor Intelligence, January 2026", 0.83, 6.58, 4.4, 0.16, 6.7, MUTED)

market_stats = [
    ("$3.5B+", "Projected global 3D bioprinting market by 2030 at ~15.9% CAGR", MINT_2, TEAL),
    ("DEC 2024", "FDA cleared Symvess — first tissue-engineered vascular graft — setting a direct precedent", BLUE_2, "5E91B5"),
    ("$420M+", "Committed by pharma companies to commercial-station research through 2030 (industry research)", GOLD_PALE, GOLD),
    ("$200–400M", "Samuel's estimated initial SAM by 2030 across cardiac, orthopedic, and retinal constructs", NAVY, TEAL),
]
for idx, (value, body, fill, accent) in enumerate(market_stats):
    yy = 1.56 + idx * 1.37
    card(slide, 7.92, yy, 4.86, 1.19, fill, fill if fill == NAVY else LINE)
    add_text(slide, value, 8.18, yy + 0.18, 1.42, 0.32, 14.5, WHITE if fill == NAVY else NAVY, True)
    add_text(slide, body, 9.78, yy + 0.17, 2.7, 0.71, 8.3, "C9D4DD" if fill == NAVY else SLATE, line_spacing=0.95)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 8.18, yy + 0.85, 0.58, 0.04, accent, accent)


# Slide 10 — Market timing
slide = base_slide("Market Timing", "Four Converging Tailwinds Create a Unique Entry Window", 10)
tailwinds = [
    (
        "~89%",
        "Launch Costs Have Collapsed",
        "Reduction in launch cost/kg since 2010",
        "Cost-to-orbit has fallen from ~$54,000/kg in the Shuttle era to ~$6,000/kg on Falcon 9, bringing research payloads into commercial range.",
        MINT_2,
        TEAL,
    ),
    (
        "4",
        "Commercial Stations Are Coming",
        "NASA-funded successor programs",
        "ISS retires around 2030. Starlab, Axiom Station, Orbital Reef, and Vast Haven-1 are creating a diversified commercial LEO ecosystem.",
        BLUE_2,
        "5E91B5",
    ),
    (
        "DEC 2024",
        "Regulatory Precedent Is Set",
        "First FDA-cleared tissue-engineered vascular graft",
        "Symvess established the first direct precedent for clinical tissue engineering and validates an actionable FDA review pathway.",
        GOLD_PALE,
        GOLD,
    ),
    (
        "$329M",
        "Capital Is Already Flowing",
        "Raised by in-space manufacturing comparator Varda",
        "Varda has raised $329M total. LambdaVision closed a $7M seed after securing more than $15M in non-dilutive grants.",
        LAVENDER,
        "7D6CB0",
    ),
]
for idx, (metric, title, sub, body, fill, accent) in enumerate(tailwinds):
    col = idx % 2
    row = idx // 2
    x = 0.55 + col * 6.17
    y = 1.56 + row * 2.73
    card(slide, x, y, 6.0, 2.55, fill, LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, y, 0.09, 2.55, accent, accent)
    add_text(slide, metric, x + 0.3, y + 0.2, 1.75, 0.47, 20, NAVY, True)
    add_text(slide, sub.upper(), x + 2.15, y + 0.25, 3.5, 0.4, 7.3, accent, True, align=PP_ALIGN.RIGHT)
    add_line(slide, x + 0.3, y + 0.84, x + 5.7, y + 0.84, "CBD8DD", 0.7)
    add_text(slide, title, x + 0.3, y + 1.08, 5.36, 0.34, 14, NAVY, True)
    add_text(slide, body, x + 0.3, y + 1.58, 5.36, 0.68, 9.1, SLATE, line_spacing=0.97)


# Slide 11 — Platform
slide = base_slide("Our Platform", "Station-Agnostic Payload Platform", 11)
add_text(
    slide,
    "One flight-qualified hardware architecture, one growing bioink library, and one protocol stack — deployable across the commercial LEO ecosystem.",
    0.55,
    1.53,
    7.15,
    0.48,
    10.2,
    SLATE,
)
features = [
    ("01", "Modular Payload Design", "Standard double-locker form factor designed for ISS, Starlab, and Axiom interfaces — one build, many stations."),
    ("02", "Proprietary Cell-Line IP", "Microgravity-optimised bioinks and print protocols compound into the platform's core defensible asset."),
    ("03", "Service-First Revenue", "Fee-for-service missions create near-term revenue while each flight expands the reusable IP base."),
]
for idx, (num, title, body) in enumerate(features):
    yy = 2.15 + idx * 1.27
    card(slide, 0.55, yy, 7.12, 1.08, WHITE, LINE)
    pill(slide, num, 0.8, yy + 0.24, 0.48, 0.31, NAVY if idx != 1 else TEAL, WHITE, 7.4)
    add_text(slide, title, 1.53, yy + 0.18, 2.25, 0.3, 11.5, NAVY, True)
    add_text(slide, body, 3.76, yy + 0.17, 3.58, 0.61, 8.8, SLATE, line_spacing=0.95)
card(slide, 0.55, 6.1, 7.12, 0.87, NAVY, NAVY)
add_text(slide, "ENGINEERED FOR COMPATIBILITY", 0.83, 6.31, 2.2, 0.2, 7.5, TEAL, True)
add_text(slide, "ISS today. Commercial stations tomorrow.", 3.17, 6.2, 4.07, 0.42, 12, WHITE, True, valign=MSO_ANCHOR.MIDDLE)

card(slide, 7.91, 1.53, 4.87, 5.44, BLUE_2, LINE)
add_text(slide, "STATION ECOSYSTEM COMPATIBILITY", 8.2, 1.82, 4.28, 0.22, 8, NAVY, True)
card(slide, 8.2, 2.22, 4.28, 0.68, NAVY, NAVY)
add_text(slide, "SAMUEL SPACE  /  PAYLOAD PLATFORM", 8.45, 2.22, 3.78, 0.68, 10, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
stations = [
    ("ISS", "Redwire BFF — active through ~2030"),
    ("Starlab", "Voyager/Airbus JV — targeting 2029"),
    ("Axiom Station", "First commercial modules ~2026–2028"),
    ("Orbital Reef", "Blue Origin/Sierra Space — CLD-funded"),
    ("Vast Haven-1", "Dedicated commercial station — 2026 target"),
]
for idx, (name, detail) in enumerate(stations):
    yy = 3.12 + idx * 0.66
    add_line(slide, 8.54, yy - 0.22, 8.54, yy + 0.22, TEAL, 1.0)
    add_shape(slide, MSO_SHAPE.OVAL, 8.44, yy - 0.1, 0.2, 0.2, TEAL, WHITE, 0.6)
    add_text(slide, name, 8.82, yy - 0.14, 1.28, 0.22, 9, NAVY, True)
    add_text(slide, detail, 10.12, yy - 0.14, 2.12, 0.32, 7.4, SLATE)
add_text(
    slide,
    "Market context only — station names do not indicate confirmed Samuel partnerships.",
    8.2,
    6.62,
    4.25,
    0.22,
    6.5,
    MUTED,
    italic=True,
)


# Slide 12 — Roadmap
slide = base_slide("Execution Plan", "Four-Phase Technology & Commercialisation Roadmap", 12)
phases = [
    ("PHASE 1", "YEARS 1–2", "Platform +\nWound Care", ["Flight-qualify payload hardware", "Commercial wound-care trials", "NASA InSPA Phase 1 grants", "Research models for pharma CROs"], "TRL 4–6", MINT_2, TEAL),
    ("PHASE 2", "YEARS 3–5", "Orthopedic +\nCardiac", ["510(k) wound-care filing", "Orthopedic graft program", "Cardiac patch preclinical work", "Series B financing"], "TRL 6–8", BLUE_2, "5E91B5"),
    ("PHASE 3", "YEARS 6–9", "Retinal +\nVascular", ["Retinal construct trials", "FDA PMA pathway engagement", "Expand to commercial stations", "Series C / pre-IPO financing"], "TRL 7–9", LAVENDER, "7D6CB0"),
    ("PHASE 4", "YEAR 10+", "Whole-Organ\nPlatform", ["Vascularised organ constructs", "Commercial station scale-up", "Platform licensing to pharma", "IPO or strategic exit"], "TRL 9+", NAVY, TEAL),
]
for idx, (phase, years, title, bullets, trl, fill, accent) in enumerate(phases):
    x = 0.55 + idx * 3.08
    card(slide, x, 1.56, 2.87, 5.41, fill, fill if fill == NAVY else LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x, 1.56, 2.87, 0.08, accent, accent)
    add_text(slide, phase, x + 0.23, 1.84, 1.3, 0.22, 8, accent, True)
    add_text(slide, years, x + 1.47, 1.84, 1.16, 0.22, 7.4, "B9C7D2" if fill == NAVY else MUTED, True, align=PP_ALIGN.RIGHT)
    add_text(slide, title, x + 0.23, 2.3, 2.35, 0.73, 16, WHITE if fill == NAVY else NAVY, True, line_spacing=0.9)
    add_line(slide, x + 0.23, 3.18, x + 2.64, 3.18, "3B5064" if fill == NAVY else "CBD7DC", 0.7)
    add_bullet_rows(slide, bullets, x + 0.24, 3.48, 2.34, 0.59, 8.5, "CAD5DE" if fill == NAVY else SLATE, accent)
    card(slide, x + 0.23, 6.32, 2.41, 0.39, accent if fill != NAVY else TEAL, accent if fill != NAVY else TEAL)
    add_text(slide, trl, x + 0.23, 6.32, 2.41, 0.39, 8.2, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)


# Slide 13 — Competitive landscape
slide = base_slide("Competitive Landscape", "Samuel Space Targets the Space-Based Clinical Quadrant", 13)
matrix_x, matrix_y, matrix_w, matrix_h = 1.35, 1.62, 11.0, 5.25
card(slide, matrix_x, matrix_y, matrix_w, matrix_h, WHITE, LINE)
add_shape(slide, MSO_SHAPE.RECTANGLE, matrix_x, matrix_y, matrix_w / 2, matrix_h / 2, BLUE_2, BLUE_2)
add_shape(slide, MSO_SHAPE.RECTANGLE, matrix_x + matrix_w / 2, matrix_y, matrix_w / 2, matrix_h / 2, MINT_2, MINT_2)
add_shape(slide, MSO_SHAPE.RECTANGLE, matrix_x, matrix_y + matrix_h / 2, matrix_w / 2, matrix_h / 2, "F5F6F7", "F5F6F7")
add_shape(slide, MSO_SHAPE.RECTANGLE, matrix_x + matrix_w / 2, matrix_y + matrix_h / 2, matrix_w / 2, matrix_h / 2, BLUE_2, BLUE_2)
add_line(slide, matrix_x + matrix_w / 2, matrix_y, matrix_x + matrix_w / 2, matrix_y + matrix_h, "B7C6CE", 1.0)
add_line(slide, matrix_x, matrix_y + matrix_h / 2, matrix_x + matrix_w, matrix_y + matrix_h / 2, "B7C6CE", 1.0)
pill(slide, "EARTH-BASED", matrix_x + 0.25, matrix_y + 0.2, 1.4, 0.3, WHITE, SLATE, 7.4)
pill(slide, "SPACE-BASED", matrix_x + matrix_w - 1.65, matrix_y + 0.2, 1.4, 0.3, TEAL, WHITE, 7.4)
add_text(slide, "CLINICAL / COMMERCIAL", matrix_x + 0.25, matrix_y + 0.72, 1.7, 0.18, 6.6, SLATE, True)
add_text(slide, "EARLY RESEARCH", matrix_x + 0.25, matrix_y + matrix_h / 2 + 0.28, 1.35, 0.18, 6.6, SLATE, True)
points = [
    ("Cellink / BICO", 3.05, 3.25, NAVY),
    ("Organovo", 3.3, 4.82, NAVY),
    ("3D Systems (Med)", 4.42, 2.72, NAVY),
    ("Redwire BFF", 8.95, 4.02, "5E91B5"),
    ("LambdaVision", 10.28, 3.12, "7D6CB0"),
    ("ROSCOSMOS / 3D BPS", 8.05, 5.45, "5E91B5"),
]
for name, x, y, color in points:
    add_shape(slide, MSO_SHAPE.OVAL, x - 0.09, y - 0.09, 0.18, 0.18, color, WHITE, 0.7)
    add_text(slide, name, x - 0.65, y + 0.15, 1.3, 0.24, 7.4, INK, True, align=PP_ALIGN.CENTER)
entry_x, entry_y = 8.34, 4.78
target_x, target_y = 10.42, 2.48
add_shape(slide, MSO_SHAPE.OVAL, entry_x - 0.08, entry_y - 0.08, 0.16, 0.16, WHITE, TEAL, 1.1)
add_text(slide, "Entry", entry_x - 0.45, entry_y + 0.14, 0.9, 0.2, 7.3, TEAL_DARK, True, align=PP_ALIGN.CENTER)
add_line(slide, entry_x + 0.12, entry_y - 0.12, target_x - 0.17, target_y + 0.2, TEAL, 1.7)
add_shape(slide, MSO_SHAPE.ISOSCELES_TRIANGLE, target_x - 0.18, target_y + 0.11, 0.18, 0.18, TEAL, TEAL).rotation = 45
add_shape(slide, MSO_SHAPE.OVAL, target_x - 0.18, target_y - 0.18, 0.36, 0.36, WHITE, TEAL, 2.2)
add_shape(slide, MSO_SHAPE.OVAL, target_x - 0.09, target_y - 0.09, 0.18, 0.18, TEAL, WHITE, 0.6)
card(slide, 10.7, 2.22, 1.32, 0.64, NAVY, NAVY)
add_text(slide, "SAMUEL SPACE\nTARGET", 10.83, 2.26, 1.06, 0.54, 7.4, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE, line_spacing=0.9)
add_text(slide, "Positions are illustrative, based on public technology and commercialisation-stage disclosures.", 1.52, 6.58, 7.9, 0.16, 6.6, MUTED, italic=True)


# Slide 14 — Revenue model
slide = base_slide("Revenue Model", "Three Complementary Revenue Streams", 14)
streams = [
    ("01", "Contract Bioprinting Services", "Fee-for-service orbital bioprinting missions for pharma, academic medical centres, and government research agencies. Revenue begins in Year 1.", "NEAR-TERM REVENUE", MINT_2, TEAL),
    ("02", "Non-Dilutive Grant Funding", "NASA InSPA, ESA BIC, DoD SBIR/STTR, and NEI mechanisms can finance development without equity dilution. Samuel targets $4–8M across Years 1–4.", "REDUCES EQUITY BURN", BLUE_2, "5E91B5"),
    ("03", "Tissue IP Licensing + Royalties", "As bioink formulations and protocols demonstrate clinical-grade output, Samuel licenses IP to pharmaceutical and device manufacturers from Year 3 onward.", "LONG-TERM VALUE DRIVER", NAVY, TEAL),
]
for idx, (num, title, body, timing, fill, accent) in enumerate(streams):
    y = 1.57 + idx * 1.78
    card(slide, 0.55, y, 12.23, 1.55, fill, fill if fill == NAVY else LINE)
    add_text(slide, num, 0.85, y + 0.21, 0.82, 0.55, 22, accent, True)
    add_text(slide, f"REVENUE STREAM {idx + 1}", 1.88, y + 0.2, 2.2, 0.2, 7.4, accent, True)
    add_text(slide, title, 1.88, y + 0.52, 3.55, 0.36, 14, WHITE if fill == NAVY else NAVY, True)
    add_text(slide, body, 5.65, y + 0.29, 4.78, 0.83, 9, "CAD5DE" if fill == NAVY else SLATE, line_spacing=0.96)
    card(slide, 10.72, y + 0.42, 1.75, 0.55, accent, accent)
    add_text(slide, timing, 10.82, y + 0.42, 1.55, 0.55, 7.2, WHITE, True, align=PP_ALIGN.CENTER, valign=MSO_ANCHOR.MIDDLE)
add_line(slide, 1.26, 6.88, 11.8, 6.88, LINE, 0.8)
add_text(slide, "SERVICES", 0.64, 6.76, 1.25, 0.18, 7, TEAL_DARK, True)
add_text(slide, "GRANTS", 5.84, 6.76, 1.25, 0.18, 7, "5E91B5", True, align=PP_ALIGN.CENTER)
add_text(slide, "SCALABLE IP", 11.18, 6.76, 1.25, 0.18, 7, TEAL_DARK, True, align=PP_ALIGN.RIGHT)


# Slide 15 — Go-to-market
slide = base_slide("Go-to-Market", "Sequenced Market Entry Across Four Partner Categories", 15)
partners = [
    ("01", "Commercial Station Operators", "Secure research-service agreements and payload manifesting rights.", "Redwire, Starlab, Axiom, Vast", "ACCESS", MINT_2, TEAL),
    ("02", "Payload Integrators + CROs", "Outsource environmental control, crew interface, and downlink to move faster and preserve capital.", "Space Tango, Nanoracks / Voyager", "INTEGRATE", BLUE_2, "5E91B5"),
    ("03", "Pharma + Academic Medical Centres", "Win contract-service demand for orbital tissue models and drug-screening programs.", "Pharma R&D, NIH-funded labs", "MONETISE", GOLD_PALE, GOLD),
    ("04", "Government + Grant Agencies", "Use non-dilutive funding to de-risk development and validate technology for commercial buyers.", "NASA, ESA, NIH NEI, DoD MTEC", "DE-RISK", LAVENDER, "7D6CB0"),
]
for idx, (num, title, body, targets, stage, fill, accent) in enumerate(partners):
    x = 0.55 + idx * 3.08
    card(slide, x, 1.56, 2.87, 5.41, fill, LINE)
    pill(slide, num, x + 0.22, 1.8, 0.5, 0.31, accent, WHITE, 7.6)
    add_text(slide, stage, x + 1.12, 1.84, 1.52, 0.2, 7.2, accent, True, align=PP_ALIGN.RIGHT)
    add_text(slide, title, x + 0.22, 2.39, 2.42, 0.72, 14, NAVY, True, line_spacing=0.92)
    add_line(slide, x + 0.22, 3.25, x + 2.64, 3.25, "C7D4DA", 0.7)
    add_text(slide, body, x + 0.22, 3.54, 2.42, 1.22, 9.1, SLATE, line_spacing=0.98)
    add_text(slide, "TARGET SET", x + 0.22, 5.28, 2.2, 0.2, 7.2, accent, True)
    add_text(slide, targets, x + 0.22, 5.66, 2.42, 0.55, 8.7, NAVY, True)
    add_shape(slide, MSO_SHAPE.RECTANGLE, x + 0.22, 6.52, 2.42, 0.06, accent, accent)


# Slide 16 — Financials
slide = base_slide("Financials — Illustrative", "Five-Year Revenue Model (Illustrative Projections)", 16, 22.5)
add_text(slide, "Illustrative projections for investor discussion only; actual results may vary materially.", 0.55, 1.48, 7.0, 0.2, 7.5, MUTED, italic=True)
card(slide, 0.55, 1.82, 7.18, 5.15, WHITE, LINE)
add_text(slide, "REVENUE BUILD  /  USD $M", 0.84, 2.08, 4.2, 0.22, 8.2, NAVY, True)
series = [
    ("Contract services", [0.5, 1.2, 3.0, 7.0, 14.0], TEAL),
    ("Grant revenue", [0.3, 1.0, 1.5, 2.0, 2.5], "5E91B5"),
    ("IP licensing", [0.0, 0.2, 1.6, 4.5, 7.7], GOLD),
]
legend_x = 4.35
for idx, (name, _, color) in enumerate(series):
    add_shape(slide, MSO_SHAPE.RECTANGLE, legend_x + idx * 1.0, 2.11, 0.12, 0.12, color, color)
    add_text(slide, name.split()[0], legend_x + 0.18 + idx * 1.0, 2.07, 0.78, 0.2, 6.4, SLATE)
plot_x, plot_y, plot_w, plot_h = 1.17, 2.66, 5.95, 3.45
for grid_val in (5, 10, 15, 20, 25):
    yy = plot_y + plot_h - grid_val / 25 * plot_h
    add_line(slide, plot_x, yy, plot_x + plot_w, yy, "E6EBED", 0.7)
    add_text(slide, str(grid_val), 0.76, yy - 0.1, 0.3, 0.18, 6.5, MUTED, align=PP_ALIGN.RIGHT)
bar_w = 0.68
bar_gap = (plot_w - 5 * bar_w) / 4
for year_idx in range(5):
    x = plot_x + year_idx * (bar_w + bar_gap)
    current_y = plot_y + plot_h
    total = 0
    for _, vals, color in series:
        value = vals[year_idx]
        hh = value / 25 * plot_h
        if hh > 0:
            add_shape(slide, MSO_SHAPE.RECTANGLE, x, current_y - hh, bar_w, hh, color, WHITE, 0.4)
            current_y -= hh
        total += value
    add_text(slide, f"${total:.1f}", x - 0.1, current_y - 0.25, bar_w + 0.2, 0.18, 7.3, NAVY, True, align=PP_ALIGN.CENTER)
    add_text(slide, f"YR {year_idx + 1}", x - 0.08, plot_y + plot_h + 0.12, bar_w + 0.16, 0.18, 6.8, SLATE, True, align=PP_ALIGN.CENTER)
add_text(slide, "Illustrative revenue, not contracted backlog", 0.84, 6.59, 3.4, 0.16, 6.5, MUTED)

card(slide, 7.94, 1.82, 4.84, 5.15, BLUE_2, LINE)
add_text(slide, "ANNUAL REVENUE DETAIL", 8.2, 2.08, 4.3, 0.22, 8.2, NAVY, True)
cols = [("YEAR", 0.65), ("TOTAL", 0.82), ("KEY DRIVER", 2.75)]
cx = 8.2
for label, width in cols:
    add_text(slide, label, cx, 2.48, width, 0.2, 7.2, MUTED, True)
    cx += width + 0.08
add_line(slide, 8.2, 2.78, 12.48, 2.78, "BFCED5", 0.8)
financial_rows = [
    ("Y1", "$0.8M", "First service contracts + NASA grant"),
    ("Y2", "$2.4M", "Expanded services + licensing launch"),
    ("Y3", "$6.1M", "Cardiac / ortho milestones + Starlab start"),
    ("Y4", "$13.5M", "Licensing deals + multi-station operations"),
    ("Y5", "$24.2M", "Scale-up across four product verticals"),
]
for idx, (year, total, driver) in enumerate(financial_rows):
    yy = 2.94 + idx * 0.72
    if idx == 4:
        card(slide, 8.13, yy - 0.08, 4.4, 0.61, NAVY, NAVY)
    color = WHITE if idx == 4 else NAVY
    body_color = "CED8E0" if idx == 4 else SLATE
    add_text(slide, year, 8.2, yy, 0.58, 0.25, 8.2, color, True)
    add_text(slide, total, 8.92, yy, 0.78, 0.25, 9.5, TEAL if idx == 4 else TEAL_DARK, True)
    add_text(slide, driver, 9.84, yy - 0.02, 2.48, 0.38, 7.5, body_color, line_spacing=0.92)
    if idx < 4:
        add_line(slide, 8.2, yy + 0.5, 12.48, yy + 0.5, "D4DEE2", 0.6)
pill(slide, "YEAR 5  /  $24.2M", 10.56, 6.48, 1.72, 0.32, TEAL, WHITE, 7.5)


# Slide 17 — The ask
slide = base_slide("The Ask", "An $18M Series A to Reach First Commercial Flight", 17)
card(slide, 0.55, 1.56, 3.16, 5.41, NAVY, NAVY)
add_text(slide, "ROUND SIZE", 0.86, 1.93, 1.8, 0.2, 7.8, TEAL, True)
add_text(slide, "$18M", 0.86, 2.34, 2.36, 0.72, 30, WHITE, True)
add_text(slide, "Series A Preferred Equity", 0.86, 3.12, 2.3, 0.26, 9.2, "C7D2DB", True)
add_line(slide, 0.86, 3.67, 3.38, 3.67, "3B5064", 0.8)
ask_metrics = [("~30 months", "Runway to next value-inflection"), ("2", "Product verticals reaching clinical pilot"), ("3–5", "Target commercial flight missions secured")]
for idx, (value, label) in enumerate(ask_metrics):
    yy = 3.95 + idx * 0.88
    add_text(slide, value, 0.86, yy, 2.18, 0.3, 12.5, TEAL, True)
    add_text(slide, label, 0.86, yy + 0.35, 2.18, 0.34, 7.6, "C7D2DB", line_spacing=0.93)

card(slide, 3.94, 1.56, 4.1, 5.41, WHITE, LINE)
add_text(slide, "USE OF FUNDS", 4.23, 1.86, 2.0, 0.22, 8.2, NAVY, True)
allocations = [
    ("R&D / Payload Development", 40, TEAL),
    ("Flight + Integration", 25, "5E91B5"),
    ("Regulatory + Clinical", 15, GOLD),
    ("Team Build-Out", 15, "7D6CB0"),
    ("G&A", 5, CORAL),
]
bar_x = 4.23
bar_y = 2.3
bar_total_w = 3.52
cursor = bar_x
for _, pct, color in allocations:
    ww = bar_total_w * pct / 100
    add_shape(slide, MSO_SHAPE.RECTANGLE, cursor, bar_y, ww, 0.34, color, WHITE, 0.5)
    cursor += ww
for idx, (label, pct, color) in enumerate(allocations):
    yy = 2.93 + idx * 0.68
    add_shape(slide, MSO_SHAPE.RECTANGLE, 4.23, yy + 0.05, 0.12, 0.12, color, color)
    add_text(slide, label, 4.52, yy, 2.42, 0.28, 8.7, INK, True)
    add_text(slide, f"{pct}%", 7.22, yy, 0.5, 0.28, 9.2, NAVY, True, align=PP_ALIGN.RIGHT)
add_text(slide, "Allocation is illustrative and subject to diligence.", 4.23, 6.52, 3.4, 0.18, 6.5, MUTED, italic=True)

card(slide, 8.27, 1.56, 4.51, 5.41, BLUE_2, LINE)
add_text(slide, "VALUE-INFLECTION MILESTONES", 8.56, 1.86, 3.92, 0.22, 8.2, NAVY, True)
milestones = [
    "Flight-qualify modular payload hardware",
    "Complete 3–5 commissioned ISS missions",
    "File first NASA InSPA Phase 1 grant",
    "Secure Starlab payload reservation",
    "Build core scientific + engineering team",
]
for idx, item in enumerate(milestones):
    yy = 2.35 + idx * 0.76
    pill(slide, f"{idx + 1:02d}", 8.56, yy, 0.48, 0.31, TEAL if idx < 3 else NAVY, WHITE, 7.2)
    add_text(slide, item, 9.28, yy - 0.01, 3.04, 0.42, 9, INK, True)
    if idx < 4:
        add_line(slide, 9.28, yy + 0.52, 12.35, yy + 0.52, "CFD9DE", 0.6)
card(slide, 8.56, 6.28, 3.92, 0.4, NAVY, NAVY)
add_text(slide, "TARGET  /  FIRST COMMERCIAL FLIGHT", 8.7, 6.28, 3.64, 0.4, 7.4, WHITE, True, valign=MSO_ANCHOR.MIDDLE)


# Slide 18 — Risks
slide = base_slide("Risk Factors", "Key Risks and Mitigating Strategies", 18)
columns = [("RISK CATEGORY", 1.55), ("RISK", 4.85), ("MITIGATING STRATEGY", 5.43)]
cx = 0.55
for label, width in columns:
    add_shape(slide, MSO_SHAPE.RECTANGLE, cx, 1.56, width, 0.46, NAVY, NAVY)
    add_text(slide, label, cx + 0.16, 1.56, width - 0.32, 0.46, 7.7, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
    cx += width + 0.08
risks = [
    ("TECHNICAL", "Scaling from validated tissue patches to full vascularised organs remains scientifically unsolved industry-wide.", "Monetise lower-complexity wound care and orthopedic products before attempting whole-organ printing.", CORAL),
    ("REGULATORY", "FDA pathways for novel tissue-engineered products are long, costly, and precedent-limited.", "Use the Symvess precedent, engage regulatory counsel from Year 1, and target lower-risk 510(k) products first.", GOLD),
    ("PLATFORM / ACCESS", "ISS retirement around 2030 creates dependency risk if successor stations are delayed.", "Design for ISS, Starlab, Axiom, Orbital Reef, and Vast to diversify single-platform exposure.", "5E91B5"),
    ("CAPITAL INTENSITY", "Flight missions, payload hardware, and clinical trials require sustained, significant investment.", "Blend equity with NASA, ESA, and DoD non-dilutive funding, following sector precedent.", "7D6CB0"),
    ("COMPETITIVE", "Redwire and LambdaVision hold first-mover relationships with NASA and station operators.", "Differentiate through multi-vertical breadth and station-agnostic flexibility, not a single-contract strategy.", TEAL),
]
for idx, (category, risk, mitigant, accent) in enumerate(risks):
    y = 2.1 + idx * 0.95
    fill = WHITE if idx % 2 == 0 else BLUE_2
    card(slide, 0.55, y, 1.55, 0.83, fill, LINE)
    add_shape(slide, MSO_SHAPE.RECTANGLE, 0.55, y, 0.08, 0.83, accent, accent)
    add_text(slide, category, 0.78, y, 1.12, 0.83, 7.7, NAVY, True, valign=MSO_ANCHOR.MIDDLE)
    card(slide, 2.18, y, 4.85, 0.83, fill, LINE)
    add_text(slide, risk, 2.4, y + 0.12, 4.41, 0.58, 8.1, SLATE, line_spacing=0.93)
    card(slide, 7.11, y, 5.67, 0.83, fill, LINE)
    add_text(slide, mitigant, 7.34, y + 0.12, 5.21, 0.58, 8.1, INK, True, line_spacing=0.93)


# Slide 19 — Organisation
slide = base_slide("Organisation", "Team & Scientific Advisory Approach", 19)
card(slide, 0.55, 1.53, 12.23, 0.56, BLUE_2, LINE)
add_text(
    slide,
    "The Series A hiring plan is organised around four capabilities required to build, certify, and commercialise the platform.",
    0.8,
    1.54,
    11.73,
    0.54,
    9.3,
    INK,
    valign=MSO_ANCHOR.MIDDLE,
)
roles = [
    ("01", "Founder / Chief\nExecutive Officer", "Commercial space or biotech leadership; regulated hardware or therapeutics experience from seed to scale.", TEAL),
    ("02", "Chief Scientific\nOfficer", "PhD-level tissue engineering, stem-cell biology, or regenerative medicine expertise.", "5E91B5"),
    ("03", "VP, Payload\nEngineering", "Aerospace hardware background with ISS or commercial LEO flight-qualification experience.", GOLD),
    ("04", "Head of Regulatory\nAffairs", "FDA medical device / biologics expertise across 510(k) and PMA submissions.", "7D6CB0"),
]
for idx, (num, title, body, accent) in enumerate(roles):
    x = 0.55 + idx * 3.08
    card(slide, x, 2.29, 2.87, 2.56, WHITE, LINE)
    pill(slide, num, x + 0.22, 2.53, 0.5, 0.31, accent, WHITE, 7.6)
    add_text(slide, "PLANNED HIRE", x + 1.13, 2.59, 1.48, 0.16, 6.8, accent, True, align=PP_ALIGN.RIGHT)
    add_text(slide, title, x + 0.22, 3.08, 2.4, 0.67, 13, NAVY, True, line_spacing=0.9)
    add_text(slide, body, x + 0.22, 3.91, 2.4, 0.68, 8.5, SLATE, line_spacing=0.95)
card(slide, 0.55, 5.1, 12.23, 1.87, NAVY, NAVY)
pill(slide, "IN FORMATION", 0.84, 5.38, 1.18, 0.32, TEAL, WHITE, 7.2)
add_text(slide, "Scientific Advisory Board", 2.25, 5.32, 3.3, 0.4, 15, WHITE, True)
add_text(
    slide,
    "Recruiting across stem-cell and cardiac tissue biology, microgravity payload design, and protein-based biomanufacturing. Target appointments will be announced as the round closes.",
    5.65,
    5.28,
    6.58,
    0.74,
    9.4,
    "CBD6DE",
    line_spacing=0.98,
)
add_line(slide, 0.84, 6.3, 12.22, 6.3, "365068", 0.7)
add_text(slide, "TARGET DISCIPLINES", 0.84, 6.49, 1.55, 0.16, 7, TEAL, True)
add_text(slide, "TISSUE BIOLOGY", 2.62, 6.49, 1.5, 0.16, 7, WHITE, True)
add_text(slide, "PAYLOAD DESIGN", 5.26, 6.49, 1.5, 0.16, 7, WHITE, True)
add_text(slide, "BIOPROCESSING", 8.06, 6.49, 1.5, 0.16, 7, WHITE, True)
add_text(slide, "REGULATORY SCIENCE", 10.68, 6.49, 1.55, 0.16, 7, WHITE, True)


# Slide 20 — Closing
slide = prs.slides.add_slide(blank)
set_background(slide)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0, 0, 0.16, SLIDE_H, TEAL, TEAL)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.16, 0, 0.05, SLIDE_H, NAVY, NAVY)
add_text(slide, "THE OPPORTUNITY", 0.72, 0.62, 3.2, 0.24, 8.5, TEAL_DARK, True)
add_text(
    slide,
    "Every organ printed in orbit\nis a life that didn't have to wait.",
    0.72,
    1.28,
    7.0,
    1.46,
    25,
    NAVY,
    True,
    line_spacing=0.93,
)
add_text(
    slide,
    "Samuel Space stands at the intersection of a validated microgravity bioprinting science base and the arrival of commercial stations built for exactly this work.",
    0.74,
    3.15,
    6.52,
    0.9,
    12.2,
    INK,
    line_spacing=1.0,
)
add_shape(slide, MSO_SHAPE.RECTANGLE, 0.74, 4.44, 0.84, 0.06, TEAL, TEAL)
add_text(slide, "The science already works.", 0.74, 4.74, 5.8, 0.35, 13, NAVY, True)
add_text(slide, "What's needed now is a company built to commercialise it.", 0.74, 5.16, 5.95, 0.35, 12.2, SLATE)
card(slide, 0.74, 6.2, 6.2, 0.63, NAVY, NAVY)
add_text(slide, "SAMUEL SPACE ORGAN AND TISSUE ENGINEERING", 0.98, 6.2, 3.85, 0.63, 7.7, WHITE, True, valign=MSO_ANCHOR.MIDDLE)
add_text(slide, "investors@samuelspace.com  (placeholder)", 4.63, 6.2, 2.02, 0.63, 7.2, "C6D2DB", align=PP_ALIGN.RIGHT, valign=MSO_ANCHOR.MIDDLE)

card(slide, 7.7, 0.42, 5.08, 6.62, MINT_2, "BFDCD5", 0.9)
add_text(slide, "THE INVESTMENT CASE", 8.05, 0.84, 4.38, 0.22, 8, TEAL_DARK, True)
case_points = [
    ("01", "Validated science", "Eight years of independent in-orbit milestones."),
    ("02", "Expanding access", "Commercial LEO infrastructure lowers platform risk."),
    ("03", "Layered economics", "Services fund the path to high-value tissue IP."),
]
for idx, (num, title, body) in enumerate(case_points):
    y = 1.46 + idx * 1.34
    pill(slide, num, 8.05, y, 0.48, 0.31, TEAL if idx == 0 else NAVY, WHITE, 7.4)
    add_text(slide, title, 8.75, y - 0.01, 3.4, 0.28, 12, NAVY, True)
    add_text(slide, body, 8.75, y + 0.42, 3.4, 0.44, 8.7, SLATE)
    add_line(slide, 8.05, y + 1.03, 12.42, y + 1.03, "C6DCD6", 0.7)
add_shape(slide, MSO_SHAPE.OVAL, 9.02, 5.37, 2.62, 1.06, "B9E7DC", TEAL, 1.1)
add_shape(slide, MSO_SHAPE.OVAL, 9.45, 5.55, 1.76, 0.7, "8DD8C8", TEAL_DARK, 1.0)
for idx in range(7):
    xx = 9.58 + idx * 0.24
    add_shape(slide, MSO_SHAPE.OVAL, xx, 5.77 + (0.06 if idx % 2 else 0), 0.15, 0.15, TEAL if idx % 2 else WHITE, TEAL, 0.5)
add_text(slide, "20 / 21", 11.84, 6.73, 0.62, 0.16, 7, MUTED, True, align=PP_ALIGN.RIGHT)


# Slide 21 — Sources
slide = base_slide("Appendix", "Sources & References", 21)
source_groups = [
    (
        "BIOPRINTING SCIENCE + TECHNOLOGY",
        [
            'NASA — "3D Bioprinting" research overview, ISS research program page',
            "NASA — BioFabrication Facility mission page, Redwire Corporation",
            "Redwire — BFF-Meniscus first human knee meniscus, Sept 2023",
            "Redwire — BFF-Cardiac investigation and ISS return coverage, 2024",
            "Redwire — NASA IDIQ contract award coverage, Sept 2025",
            "ESA / DLR — Bioprint FirstAid, Cosmic Kiss mission materials",
            "LambdaVision — Space Tango partnership, InSPA Phase 2 award, seed financing, 2025–2026",
        ],
        0.55,
    ),
    (
        "MARKET, REGULATORY + INDUSTRY DATA",
        [
            "ROSCOSMOS — 3D MBP / Organ.Aut mission summaries, 2018–2020",
            "American Heart Association — 2025 Heart Disease & Stroke Statistics",
            "OPTN / organdonor.gov (HRSA) — U.S. transplant waitlist, 2025",
            "WHO — Global Observatory on Donation and Transplantation",
            "Mordor Intelligence — 3D Bioprinting Market Report, Jan 2026",
            "U.S. FDA — Symvess clearance announcement, Dec 2024",
            "NASA CLD program — Starlab, Axiom, Orbital Reef, Vast program pages",
        ],
        6.72,
    ),
]
for title, sources, x in source_groups:
    card(slide, x, 1.56, 6.0, 3.95, WHITE, LINE)
    add_text(slide, title, x + 0.26, 1.84, 5.46, 0.22, 8, TEAL_DARK, True)
    add_line(slide, x + 0.26, 2.18, x + 5.74, 2.18, LINE, 0.7)
    for idx, source in enumerate(sources):
        yy = 2.4 + idx * 0.41
        add_text(slide, f"{idx + 1:02d}", x + 0.26, yy, 0.32, 0.2, 6.8, TEAL_DARK, True)
        add_text(slide, source, x + 0.72, yy - 0.01, 4.98, 0.3, 7.5, SLATE, line_spacing=0.9)
card(slide, 0.55, 5.76, 12.23, 1.2, NAVY, NAVY)
pill(slide, "METHODOLOGY", 0.83, 6.02, 1.2, 0.3, TEAL, WHITE, 7.2)
add_text(
    slide,
    "All third-party milestones, dates, and figures are drawn from public sources and were current as of the cited date. References to NASA, Redwire, LambdaVision, ROSCOSMOS, and ESA/DLR describe industry validation — not partnerships, endorsements, or affiliations with Samuel Space.",
    2.28,
    5.96,
    10.1,
    0.66,
    8.2,
    "CBD6DE",
    line_spacing=0.96,
)


prs.save(OUTPUT)
print(f"Saved {OUTPUT} with {len(prs.slides)} slides")
